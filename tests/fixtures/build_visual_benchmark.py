#!/usr/bin/env python3
"""Validate and freeze the offline visual benchmark contract without rendering."""

from __future__ import annotations

import argparse
import hashlib
import hmac
import json
import re
import secrets
import shutil
import sys
import tempfile
from collections import Counter, defaultdict
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[2]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))
CASES_PATH = Path(__file__).with_name("visual_quality_cases.json")
PREFERENCES_PATH = Path(__file__).with_name("visual_preferences.json")
SCENARIOS_PATH = Path(__file__).with_name("scenarios.json")
PRECHANGE_SVG_DIR = Path(__file__).with_name("visual_prechange_svg")
FINAL_CONTRACT_PATH = Path(__file__).with_name("final_deterministic_contract.json")
PURPOSES = {"architecture", "process", "comparison", "timeline", "metrics", "summary"}
PARTITIONS = {"calibration", "pilot-holdout", "final-holdout"}
CONSUMER_PARTITION = {
    "tuning": "calibration",
    "u11-pilot": "pilot-holdout",
    "final-evaluation": "final-holdout",
}

from scripts.ghb_visual_quality import (  # noqa: E402
    analyze_deck_quality,
    evaluate_page_quality,
    measure_svg,
)
from scripts.ppt_master.svg_layouts import LAYOUT_CONTRACTS, LayoutSpec, render_layout  # noqa: E402
from scripts.render_ghb_pptx import font_evidence  # noqa: E402
from scripts.validate_ghb_pptx import validate_pptx  # noqa: E402
from scripts.validate_project_contract import default_visual_profile  # noqa: E402


def canonical_sha256(value: Any) -> str:
    payload = json.dumps(value, ensure_ascii=False, sort_keys=True, separators=(",", ":"))
    return hashlib.sha256(payload.encode("utf-8")).hexdigest()


def partition_digest(cases: list[dict[str, Any]]) -> str:
    assignments = sorted((case["case_id"], case["partition"]) for case in cases)
    return canonical_sha256(assignments)


def advisory_denominator_digest(pairs: list[str]) -> str:
    return canonical_sha256(sorted(pairs))


def cases_for_consumer(corpus: dict[str, Any], consumer: str) -> list[dict[str, Any]]:
    try:
        partition = CONSUMER_PARTITION[consumer]
    except KeyError as exc:
        raise ValueError(f"unknown benchmark consumer: {consumer}") from exc
    return [dict(case) for case in corpus["cases"] if case["partition"] == partition]


def _require(condition: bool, message: str) -> None:
    if not condition:
        raise ValueError(message)


def validate_contracts(corpus: dict[str, Any], preferences: dict[str, Any]) -> None:
    _require(corpus.get("schema") == "ghb.visual-quality-corpus.v1", "unknown corpus schema")
    _require(preferences.get("schema") == "ghb.visual-preferences.v1", "unknown preferences schema")
    cases = corpus.get("cases")
    _require(isinstance(cases, list) and len(cases) >= 30, "corpus requires at least 30 cases")
    ids = [case.get("case_id") for case in cases]
    _require(all(isinstance(case_id, str) and case_id for case_id in ids), "case_id is required")
    _require(len(ids) == len(set(ids)), "duplicate case_id")
    source_ids = [
        (case.get("source", {}).get("scenario_id"), case.get("source", {}).get("body_slide_index"))
        for case in cases
    ]
    _require(len(source_ids) == len(set(source_ids)), "source identity cannot cross cases or partitions")
    _require({case.get("page_purpose") for case in cases} == PURPOSES, "all six page purposes are required")

    coverage: dict[str, Counter[str]] = defaultdict(Counter)
    for case in cases:
        purpose = case.get("page_purpose")
        partition = case.get("partition")
        _require(partition in PARTITIONS, f"{case['case_id']}: invalid partition")
        coverage[purpose][partition] += 1
        _require(case.get("density") in {"breathing", "anchor", "dense"}, f"{case['case_id']}: invalid density")
        _require(case.get("content_length_band") in {"short", "medium", "long"}, f"{case['case_id']}: invalid content band")
        _require(case.get("item_count_band") in {"few", "standard", "many"}, f"{case['case_id']}: invalid item band")
        source = case.get("source", {})
        _require(source.get("scenario_file") == "tests/fixtures/scenarios.json", f"{case['case_id']}: invalid scenario source")
        _require(isinstance(source.get("scenario_id"), str), f"{case['case_id']}: scenario_id is required")
        _require(isinstance(source.get("body_slide_index"), int), f"{case['case_id']}: body_slide_index is required")
        evidence = case.get("pre_change_evidence", {})
        _require(evidence.get("categorical_status") in {"approved-current", "known-limitation"}, f"{case['case_id']}: categorical evidence is required")
        provenance = evidence.get("provenance", {})
        _require(isinstance(provenance.get("renderer"), str), f"{case['case_id']}: renderer is required")
        _require(isinstance(provenance.get("dpi"), int), f"{case['case_id']}: dpi is required")
        _require(isinstance(provenance.get("fonts"), list), f"{case['case_id']}: fonts are required")
        digest = evidence.get("authored_svg_sha256")
        _require(isinstance(digest, str) and len(digest) == 64, f"{case['case_id']}: authored SVG digest is required")
        _require(evidence.get("rendered_png") is None, f"{case['case_id']}: rendered PNG must not be fabricated")
        _require(evidence.get("render_status") == "unavailable-without-render", f"{case['case_id']}: render limitation is required")
    for purpose in PURPOSES:
        _require(coverage[purpose]["calibration"] >= 2, f"{purpose}: needs two calibration cases")
        _require(coverage[purpose]["pilot-holdout"] >= 1, f"{purpose}: needs one pilot holdout")
        _require(coverage[purpose]["final-holdout"] >= 2, f"{purpose}: needs two final holdouts")
    _require(partition_digest(cases) == corpus.get("partition_assignment_sha256"), "partition assignment digest changed")

    scenarios = json.loads(SCENARIOS_PATH.read_text(encoding="utf-8"))
    _require(sum(len(scenario["slides"]) for scenario in scenarios.values()) == 30, "scenario corpus must contain exactly 30 independent pages")
    for case in cases:
        source = case["source"]
        try:
            slide = scenarios[source["scenario_id"]]["slides"][source["body_slide_index"] - 1]
        except (KeyError, IndexError) as exc:
            raise ValueError(f"{case['case_id']}: source slide does not exist") from exc
        _require(slide.get("page_purpose") == case["page_purpose"], f"{case['case_id']}: page-purpose drift")
        _require(slide.get("density") == case["density"], f"{case['case_id']}: density drift")

    protocol = preferences.get("blind_review_protocol", {})
    _require(protocol.get("pair_order_randomized") is True, "pair order must be randomized")
    _require(protocol.get("identity_anonymized") is True, "pair identities must be anonymized")
    _require(protocol.get("minimum_independent_eligible_reviewers_per_page", 0) >= 3, "three reviewers are required")
    _require(protocol.get("minimum_eligible_non_tie_judgments_per_page", 0) >= 2, "two non-tie judgments are required")
    _require(protocol.get("structural_regression") == "page-level-veto", "structural regression must veto the page")
    gate = preferences.get("deterministic_pilot_false_positive_gate", {})
    pairs = gate.get("advisory_rule_case_pairs", [])
    _require(gate.get("blocking_false_positive_maximum") == 0, "blocking false positives must be zero")
    _require(gate.get("advisory_false_positive_rate_maximum") == 0.10, "advisory ceiling must be 10 percent")
    _require(gate.get("advisory_denominator_count") == len(pairs), "advisory denominator count changed")
    _require(len(pairs) == len(set(pairs)), "duplicate advisory rule/case pair")
    _require(advisory_denominator_digest(pairs) == gate.get("advisory_denominator_sha256"), "advisory denominator digest changed")
    pilot_ids = {case["case_id"] for case in cases if case["partition"] == "pilot-holdout"}
    _require({pair.split("::", 1)[0] for pair in pairs} == pilot_ids, "advisory denominator must cover every pilot case")


def build(output: Path, corpus: dict[str, Any], preferences: dict[str, Any]) -> dict[str, Any]:
    validate_contracts(corpus, preferences)
    if output.exists():
        raise FileExistsError(f"refusing to overwrite frozen benchmark output: {output}")
    output.mkdir(parents=True)
    svg_dir = output / "authored-svg"
    svg_dir.mkdir()
    generated: list[dict[str, str]] = []
    for case in corpus["cases"]:
        frozen_svg = PRECHANGE_SVG_DIR / f"{case['case_id']}.svg"
        if not frozen_svg.is_file() or frozen_svg.is_symlink():
            raise ValueError(f"{case['case_id']}: frozen pre-change SVG is missing or unsafe")
        digest = hashlib.sha256(frozen_svg.read_bytes()).hexdigest()
        if digest != case["pre_change_evidence"]["authored_svg_sha256"]:
            raise ValueError(f"{case['case_id']}: pre-change SVG digest changed")
        svg_path = svg_dir / frozen_svg.name
        shutil.copy2(frozen_svg, svg_path)
        generated.append({"case_id": case["case_id"], "authored_svg_sha256": digest})
    shutil.copy2(CASES_PATH, output / CASES_PATH.name)
    shutil.copy2(PREFERENCES_PATH, output / PREFERENCES_PATH.name)
    result = {
        "schema": "ghb.visual-benchmark-manifest.v1",
        "case_count": len(corpus["cases"]),
        "corpus_sha256": canonical_sha256(corpus),
        "preferences_sha256": canonical_sha256(preferences),
        "partition_assignment_sha256": corpus["partition_assignment_sha256"],
        "network": "disabled-not-required",
        "model_adapter": "not-used",
        "credentials": "not-used",
        "binary_evidence": "copied-from-versioned-pre-change-svg-fixtures",
        "generated_authored_svg": generated,
    }
    (output / "benchmark-manifest.json").write_text(
        json.dumps(result, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
    )
    return result


def _final_page_schema(case: dict[str, Any], slide: dict[str, Any]) -> dict[str, Any]:
    family = str(slide["layout_type"])
    density = "balanced" if case["density"] == "anchor" else case["density"]
    emphasis = "single-focal" if case["density"] == "anchor" else "distributed"
    variant = f"{family}/default"
    if family == "matrix":
        variant = "matrix/metric-callout" if case["page_purpose"] == "metrics" else "matrix/comparison"
    contract = LAYOUT_CONTRACTS[family]
    schema: dict[str, Any] = {
        "schema_version": 1,
        "page_purpose": case["page_purpose"],
        "layout_variant": variant,
        "density": density,
        "emphasis": emphasis,
        "budgets": {
            "max_items": contract.max_items,
            "max_text_chars": contract.max_text_chars,
        },
    }
    if emphasis == "single-focal":
        schema["focal_target"] = slide["items"][1 if len(slide["items"]) > 1 else 0]
    return schema


def _after_svg(case: dict[str, Any], slide: dict[str, Any], schema: dict[str, Any]) -> str:
    title = str(slide["key_message"])
    fragment = render_layout(
        LayoutSpec(
            str(slide["layout_type"]),
            [str(item) for item in slide["items"]],
            x=100,
            y=240,
            width=1080,
            height=390,
            density=str(schema["density"]),
            variant=str(schema["layout_variant"]),
            emphasis=str(schema["emphasis"]),
            focal_target=schema.get("focal_target"),
            max_items=int(schema["budgets"]["max_items"]),
            max_text_chars=int(schema["budgets"]["max_text_chars"]),
        )
    )
    def escaped(value: str) -> str:
        return value.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    if len(title) > 42:
        target = len(title) // 2
        breakpoints = [match.start() for match in re.finditer(r"[ ，、；|/]", title)]
        split_at = min(breakpoints, key=lambda position: abs(position - target)) if breakpoints else target
        title_svg = (
            f'<text x="108" y="148" font-size="22" font-weight="bold" fill="#2B2B2B">{escaped(title[: split_at + 1].strip())}</text>'
            f'<text x="108" y="180" font-size="22" font-weight="bold" fill="#2B2B2B">{escaped(title[split_at + 1 :].strip())}</text>'
        )
    else:
        title_svg = f'<text x="108" y="162" font-size="30" font-weight="bold" fill="#2B2B2B">{escaped(title)}</text>'
    return (
        '<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" '
        'viewBox="0 0 1280 720">\n'
        '  <g id="bg"><rect width="1280" height="720" fill="#FFFFFF"/></g>\n'
        '  <g id="bg-surface"><rect x="56" y="96" width="1168" height="608" '
        'rx="12" fill="#FFFFFF" stroke="#E0E0E0"/></g>\n'
        '  <g id="header"><rect x="88" y="132" width="6" height="40" fill="#AB1F29"/>'
        f'{title_svg}</g>\n'
        f'  {fragment}\n'
        '</svg>\n'
    )


def _blind_assignment(case_id: str, seed: str) -> tuple[dict[str, str], str, str]:
    left = f"{case_id}::A"
    right = f"{case_id}::B"
    digest = hmac.new(bytes.fromhex(seed), case_id.encode("utf-8"), hashlib.sha256).digest()
    optimized_left = digest[0] & 1 == 0
    roles = {
        left: "optimized" if optimized_left else "baseline",
        right: "baseline" if optimized_left else "optimized",
    }
    return roles, left, right


def _write_blind_svg(source: Path, destination: Path) -> None:
    """Remove inert authoring metadata that would reveal the generated variant."""

    svg = source.read_text(encoding="utf-8")
    reviewer_safe = re.sub(r'\sdata-[A-Za-z0-9_-]+="[^"]*"', "", svg)
    destination.write_text(reviewer_safe, encoding="utf-8")


def build_final_evaluation(
    output: Path,
    corpus: dict[str, Any],
    preferences: dict[str, Any],
) -> dict[str, Any]:
    """Build the sealed final-holdout handoff without making a human decision."""
    validate_contracts(corpus, preferences)
    if output.exists():
        raise FileExistsError(f"refusing to overwrite final benchmark output: {output}")
    scenarios = json.loads(SCENARIOS_PATH.read_text(encoding="utf-8"))
    final_cases = cases_for_consumer(corpus, "final-evaluation")
    output.mkdir(parents=True)
    before_dir = output / "before"
    after_dir = output / "after"
    blind_dir = output / "blind"
    for directory in (before_dir, after_dir, blind_dir):
        directory.mkdir()

    metrics: list[dict[str, Any]] = []
    private_pairs: list[dict[str, Any]] = []
    public_pairs: list[dict[str, Any]] = []
    randomization_seed = secrets.token_hex(32)
    for case in final_cases:
        case_id = case["case_id"]
        frozen = PRECHANGE_SVG_DIR / f"{case_id}.svg"
        if not frozen.is_file() or frozen.is_symlink():
            raise ValueError(f"{case_id}: frozen pre-change SVG is missing or unsafe")
        before = frozen.read_bytes()
        before_sha = hashlib.sha256(before).hexdigest()
        if before_sha != case["pre_change_evidence"]["authored_svg_sha256"]:
            raise ValueError(f"{case_id}: frozen-before-digest-mismatch")
        source = case["source"]
        slide = scenarios[source["scenario_id"]]["slides"][source["body_slide_index"] - 1]
        schema = _final_page_schema(case, slide)
        after = _after_svg(case, slide, schema)
        after_sha = hashlib.sha256(after.encode("utf-8")).hexdigest()
        (before_dir / frozen.name).write_bytes(before)
        (after_dir / frozen.name).write_text(after, encoding="utf-8")

        roles, left, right = _blind_assignment(case_id, randomization_seed)
        case_blind = blind_dir / case_id
        case_blind.mkdir()
        sources = {"baseline": before_dir / frozen.name, "optimized": after_dir / frozen.name}
        _write_blind_svg(sources[roles[left]], case_blind / "A.svg")
        _write_blind_svg(sources[roles[right]], case_blind / "B.svg")
        private_pairs.append({
            "page_case_id": case_id,
            "page_purpose": case["page_purpose"],
            "masked_left_id": left,
            "masked_right_id": right,
            "roles": roles,
            "blind_artifacts": {
                left: {
                    "artifact_path": f"blind/{case_id}/A.svg",
                    "sha256": hashlib.sha256(
                        (case_blind / "A.svg").read_bytes()
                    ).hexdigest(),
                },
                right: {
                    "artifact_path": f"blind/{case_id}/B.svg",
                    "sha256": hashlib.sha256(
                        (case_blind / "B.svg").read_bytes()
                    ).hexdigest(),
                },
            },
        })
        public_pairs.append({
            "page_case_id": case_id,
            "page_purpose": case["page_purpose"],
            "masked_left_id": left,
            "masked_right_id": right,
            "left_artifact": f"blind/{case_id}/A.svg",
            "right_artifact": f"blind/{case_id}/B.svg",
        })
        metrics.append({
            "case_id": case_id,
            "page_purpose": case["page_purpose"],
            "family": slide["layout_type"],
            "page_schema": schema,
            "before_sha256": before_sha,
            "after_sha256": after_sha,
            "before_metrics": measure_svg(before.decode("utf-8")),
            "after_metrics": measure_svg(after),
            "provenance": {
                "before": "versioned-frozen-pre-change-svg",
                "after": "current-schema-aware-render-layout",
                "semantic_source": source,
            },
            "limitations": [
                "svg-proxy-metrics-do-not-prove-human-aesthetic-preference",
                "png-pdf-office-render-not-produced-by-offline-builder",
            ],
        })

    protocol = preferences["blind_review_protocol"]
    deterministic = evaluate_deterministic_fixtures()
    private = {
        "schema": "ghb.visual-final-review.v1",
        "eligible_case_ids": [case["case_id"] for case in final_cases],
        "pair_assignments": private_pairs,
        "randomization": {
            "algorithm": "hmac-sha256",
            "seed": randomization_seed,
            "visibility": "private-evaluator-only",
        },
        "reviewer_eligibility_roster": [],
        "structural_evidence": {
            "pptx": {"status": "unavailable", "limitations": ["svg-only-bundle"]},
            "render": {"status": "unavailable", "limitations": ["no-office-render"]},
            "target_fonts": {"status": "unavailable", "limitations": ["font-fidelity-not-evaluated"]},
            "contact_sheet": {"status": "unavailable", "limitations": ["no-gui-render-evidence"]},
        },
        "judgments": [],
        "decision": "pending",
    }
    public = {
        "schema": "ghb.visual-final-blind-template.v1",
        "protocol": protocol,
        "eligible_case_ids": [case["case_id"] for case in final_cases],
        "pairs": public_pairs,
        "eligibility_attestation_template": {
            "reviewer_id_hash": "<stable-opaque-reviewer-id>",
            "independent": True,
            "did_not_author": True,
            "did_not_tune": True,
        },
        "evidence_status": private["structural_evidence"],
        "unreviewable_dimensions": [
            "readability", "editability", "office-clipping", "font-fidelity"
        ],
        "judgment_template": {
            "reviewer_id_hash": "<stable-opaque-reviewer-id>",
            "page_case_id": "<page-case-id>",
            "masked_left_id": "<masked-left-id>",
            "masked_right_id": "<masked-right-id>",
            "presented_order": ["<masked-left-id>", "<masked-right-id>"],
            "judgment": "<masked-left-id|masked-right-id|tie|abstain>",
            "rubric": {dimension: "<1-5|not-scored>" for dimension in protocol["rubric_dimensions"]},
            "structural_veto": False,
            "recorded_at": "<ISO-8601 timestamp>",
        },
        "judgments": [],
        "decision": "pending",
    }
    (output / "evaluator-record.json").write_text(json.dumps(private, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    (output / "blind-review-template.json").write_text(json.dumps(public, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    (output / "deterministic-fixtures.json").write_text(json.dumps(deterministic, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    manifest = {
        "schema": "ghb.visual-final-benchmark.v1",
        "eligible_case_count": len(final_cases),
        "eligible_case_ids": [case["case_id"] for case in final_cases],
        "partition": "final-holdout",
        "partition_assignment_sha256": corpus["partition_assignment_sha256"],
        "case_metrics": metrics,
        "deterministic_fixture_report": "deterministic-fixtures.json",
        "decision": "pending",
        "reason": "independent eligible human judgments have not been recorded",
        "network": "not-used",
        "model_adapter": "not-used",
        "credentials": "not-used",
        "provenance": {
            "corpus_sha256": canonical_sha256(corpus),
            "preferences_sha256": canonical_sha256(preferences),
            "renderer": "scripts.ppt_master.svg_layouts.render_layout",
        },
        "limitations": [
            "no-human-preference-claim-until-evaluator-record-is-complete",
            "svg-proxy-metrics-cannot-replace-office-render-and-human-review",
            "contact-sheet-unavailable-no-office-render",
            "structural-gate-blocked-until-pptx-render-font-and-contact-sheet-evidence-exist",
        ],
    }
    (output / "final-benchmark-manifest.json").write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return manifest


def _fixture_svg(rects: str) -> str:
    return f'<svg viewBox="0 0 200 100"><g id="layout-fixture">{rects}</g></svg>'


def _codes(result: dict[str, Any]) -> list[str]:
    return sorted({finding["code"] for finding in result["findings"]})


def _empty_text_fixture_evidence() -> tuple[list[str], dict[str, float]]:
    with tempfile.TemporaryDirectory() as tmp:
        path = Path(tmp) / "empty-text.pptx"
        presentation = Presentation()
        blank = presentation.slide_layouts[6]
        for index, text in enumerate(("Cover", "", "Ending")):
            slide = presentation.slides.add_slide(blank)
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
            box.text = text
            if index == 1:
                box.name = "intentional-empty-body-text-box"
        presentation.save(path)
        report = validate_pptx(path, expected_body_count=1, expect_ending=True)
    codes = sorted({issue.code for issue in report.issues if issue.code == "empty-text-box"})
    count = sum(issue.code == "empty-text-box" for issue in report.issues)
    return codes, {"empty_text_box_issue_count": float(count)}


def evaluate_deterministic_fixtures() -> dict[str, Any]:
    """Characterize stable SVG proxy rules and state unsupported cases plainly."""
    contract_document = json.loads(FINAL_CONTRACT_PATH.read_text(encoding="utf-8"))
    contract = contract_document.get("fixtures")
    if (
        contract_document.get("schema")
        != "ghb.visual-final-deterministic-contract.v1"
        or not isinstance(contract, dict)
        or not contract
    ):
        raise ValueError("missing-final-deterministic-fixture-contract")
    profile = default_visual_profile()
    tolerance = {"occupancy": 1e-6, "emphasis-color-area": 1e-6}

    approved_svg = _fixture_svg(
        '<rect x="10" y="20" width="80" height="60" fill="#FFFFFF"/>'
        '<rect x="110" y="20" width="80" height="60" fill="#FFFFFF"/>'
    )
    approved = evaluate_page_quality(
        approved_svg, slide_id="approved-balanced", profile=profile,
        page_schema={"density": "balanced", "emphasis": "distributed"},
    )

    tiny_svg = _fixture_svg(
        "".join(f'<rect x="{8 + index * 40}" y="40" width="16" height="16" fill="#FFFFFF"/>' for index in range(4))
    )
    underscaled = evaluate_page_quality(
        tiny_svg, slide_id="underscaled-content", profile=profile,
        page_schema={"density": "balanced", "emphasis": "distributed"},
    )
    overuse_svg = _fixture_svg('<rect x="16" y="16" width="168" height="64" fill="#AB1F29"/>')
    overuse = evaluate_page_quality(
        overuse_svg, slide_id="primary-color-overuse", profile=profile,
        page_schema={"density": "balanced", "emphasis": "distributed"},
    )
    overflow_svg = _fixture_svg(
        '<g data-qa-role="title" data-qa-box="2 2 196 20"><text font-size="24">A title that exceeds its explicit safe bounds</text></g>'
        '<rect x="20" y="40" width="160" height="45" fill="#FFFFFF"/>'
    )
    overflow = evaluate_page_quality(
        overflow_svg, slide_id="long-title-overflow", profile=profile,
        page_schema={
            "density": "balanced", "emphasis": "distributed",
            "bounds_override": {"x": 10, "y": 10, "width": 180, "height": 80},
        },
    )

    repeated_pages = [
        evaluate_page_quality(
            approved_svg.replace('id="layout-fixture"', f'id="layout-marker-{index}"'),
            slide_id=f"fake-diversity-{index}", profile=profile,
            page_schema={"density": "balanced", "rhythm_role": "continuity", "layout_variant": f"variant-{index}"},
        )
        for index in range(2)
    ]
    fake_diversity = analyze_deck_quality(repeated_pages, profile=profile)
    focal_svg = _fixture_svg(
        '<rect x="10" y="20" width="80" height="60" fill="#AB1F29" data-focal="true"/>'
        '<rect x="110" y="20" width="80" height="60" fill="#FFFFFF"/>'
    )
    focal_pages = [
        evaluate_page_quality(
            focal_svg, slide_id=f"repeated-focal-{index}", profile=profile,
            page_schema={"density": "balanced", "rhythm_role": "continuity", "layout_variant": f"focal-{index}"},
        )
        for index in range(3)
    ]
    repeated_focal = analyze_deck_quality(focal_pages, profile=profile)
    empty_codes, empty_metrics = _empty_text_fixture_evidence()
    font = font_evidence("Microsoft YaHei target font is unavailable")

    def metrics(result: dict[str, Any]) -> dict[str, float]:
        measured = result["measurements"]
        return {
            "occupancy": measured["occupancy"]["value"],
            "emphasis-color-area": measured["emphasis-color-area"]["value"],
        }

    def fixture(
        fixture_id: str,
        status: str,
        expected_codes: list[str],
        observed_codes: list[str],
        observed_metrics: dict[str, float],
        expected_metrics: dict[str, float],
        source: str,
        *,
        tolerances: dict[str, float] = tolerance,
        limitations: list[str] | None = None,
    ) -> dict[str, Any]:
        return {
            "fixture_id": fixture_id,
            "status": status,
            "expected_issue_codes": expected_codes,
            "observed_issue_codes": observed_codes,
            "metrics": observed_metrics,
            "expected_metrics": expected_metrics,
            "metric_tolerances": tolerances,
            "limitations": limitations or [],
            "source": source,
        }

    rows = [
        fixture("approved-balanced", "approved", [], _codes(approved), metrics(approved),
                {"occupancy": 0.48, "emphasis-color-area": 0.0}, "ghb_visual_quality.evaluate_page_quality"),
        fixture("fake-diversity", "intentionally-failing", ["visual-composition-repeated"],
                _codes(fake_diversity), metrics(repeated_pages[0]),
                {"occupancy": 0.48, "emphasis-color-area": 0.0}, "ghb_visual_quality.analyze_deck_quality",
                limitations=["deck-rule-requires-at-least-two-pages"]),
        fixture("empty-text", "intentionally-failing", ["empty-text-box"], empty_codes,
                empty_metrics, {"empty_text_box_issue_count": 1.0}, "validate_ghb_pptx.validate_pptx",
                tolerances={"empty_text_box_issue_count": 0.0}, limitations=["requires-pptx-structure-validation"]),
        fixture("font-limited-rendering", "known-limitation", ["target-font-missing"],
                font["limitation_codes"], {"limited": 1.0}, {"limited": 1.0},
                "render_ghb_pptx.font_evidence", tolerances={"limited": 0.0},
                limitations=font["limitation_codes"]),
        fixture("long-title-overflow", "intentionally-failing", ["visual-explicit-bounds-violation"],
                _codes(overflow), metrics(overflow), {"occupancy": 0.556, "emphasis-color-area": 0.0},
                "ghb_visual_quality.evaluate_page_quality", limitations=["requires-explicit-measurable-bounds"]),
        fixture("underscaled-content", "intentionally-failing", ["visual-occupancy-below-min"],
                _codes(underscaled), metrics(underscaled), {"occupancy": 0.0512, "emphasis-color-area": 0.0},
                "ghb_visual_quality.evaluate_page_quality"),
        fixture("primary-color-overuse", "intentionally-failing", ["visual-primary-color-overuse"],
                _codes(overuse), metrics(overuse), {"occupancy": 0.5376, "emphasis-color-area": 0.5376},
                "ghb_visual_quality.evaluate_page_quality"),
        fixture("repeated-focal-zones", "intentionally-failing",
                ["visual-composition-repeated", "visual-focal-zone-streak"], _codes(repeated_focal),
                metrics(focal_pages[0]), {"occupancy": 0.48, "emphasis-color-area": 0.24},
                "ghb_visual_quality.analyze_deck_quality", limitations=["deck-rule-requires-at-least-three-pages"]),
    ]
    for row in rows:
        expected_contract = contract.get(row["fixture_id"])
        if not isinstance(expected_contract, dict):
            raise ValueError("final-deterministic-fixture-contract-mismatch")
        if row["expected_issue_codes"] != expected_contract.get("expected_issue_codes"):
            raise ValueError("final-deterministic-fixture-contract-mismatch")
        if row["expected_metrics"] != expected_contract.get("expected_metrics"):
            raise ValueError("final-deterministic-fixture-contract-mismatch")
        if row["metric_tolerances"] != expected_contract.get("metric_tolerances"):
            raise ValueError("final-deterministic-fixture-contract-mismatch")
        expected = set(row["expected_issue_codes"])
        observed = set(row["observed_issue_codes"])
        row["unexpected_issue_codes"] = sorted(observed - expected)
        row["missing_expected_issue_codes"] = sorted(expected - observed)
        row["metric_regressions"] = sorted(
            metric
            for metric, expected_value in row["expected_metrics"].items()
            if abs(row["metrics"][metric] - expected_value) > row["metric_tolerances"][metric]
        )
    if {row["fixture_id"] for row in rows} != set(contract):
        raise ValueError("final-deterministic-fixture-contract-mismatch")
    return {
        "schema": "ghb.visual-benchmark-deterministic.v1",
        "contract_sha256": canonical_sha256(contract_document),
        "fixtures": rows,
        "provenance": {
            "metrics": "scripts.ghb_visual_quality",
            "profile": "scripts.validate_project_contract.default_visual_profile",
        },
        "limitations": ["categorical-and-numeric-proxy-evidence-only", "no-pixel-equality-gate"],
    }


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output", type=Path)
    parser.add_argument("--validate-only", action="store_true")
    parser.add_argument("--consumer", choices=sorted(CONSUMER_PARTITION))
    parser.add_argument("--final-evaluation", action="store_true")
    args = parser.parse_args()
    corpus = json.loads(CASES_PATH.read_text(encoding="utf-8"))
    preferences = json.loads(PREFERENCES_PATH.read_text(encoding="utf-8"))
    validate_contracts(corpus, preferences)
    if args.consumer:
        print(json.dumps(cases_for_consumer(corpus, args.consumer), ensure_ascii=False, indent=2))
        return 0
    if args.validate_only:
        print(f"validated {len(corpus['cases'])} frozen visual benchmark cases")
        return 0
    if args.final_evaluation:
        if args.output is None:
            parser.error("--output is required with --final-evaluation")
        result = build_final_evaluation(args.output, corpus, preferences)
        print(f"built pending final evaluation for {result['eligible_case_count']} cases at {args.output}")
        return 0
    if args.output is None:
        parser.error("--output is required unless --validate-only or --consumer is used")
    result = build(args.output, corpus, preferences)
    print(f"frozen {result['case_count']} cases at {args.output}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
