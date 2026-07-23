from __future__ import annotations

import hashlib
import tempfile
import unittest
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from scripts.merge_template_master import (
    CONTENT_TYPES,
    MergeError,
    NS,
    merge_pptx,
    parse_rels,
    presentation_slide_parts,
    qn,
    relation_type,
    resolve_target,
    slide_layout_part,
)


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "templates" / "GHB_PPT_模板.pptx"


def make_content(
    path: Path,
    count: int,
    *,
    with_image: bool = False,
    section_labels: list[str] | None = None,
) -> bytes | None:
    presentation = Presentation()
    presentation.slide_width = 12192000
    presentation.slide_height = 6858000
    blank = presentation.slide_layouts[6]
    picture_bytes = None
    for index in range(count):
        slide = presentation.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(7), Inches(1))
        box.text = f"Body {index + 1} 可编辑正文"
        if section_labels is not None:
            label = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(5), Inches(0.3))
            label.name = "template-section-label"
            label.text = section_labels[index]
        if with_image and index == 0:
            image_path = path.with_suffix(".png")
            Image.new("RGB", (120, 80), (0, 80, 180)).save(image_path)
            picture_bytes = image_path.read_bytes()
            slide.shapes.add_picture(str(image_path), Inches(1), Inches(2), width=Inches(2))
    presentation.save(path)
    return picture_bytes


def package(path: Path) -> dict[str, bytes]:
    with zipfile.ZipFile(path) as archive:
        return {info.filename: archive.read(info.filename) for info in archive.infolist() if not info.is_dir()}


class MergeTemplateMasterTest(unittest.TestCase):
    def test_semantic_section_labels_move_into_native_template_title_frame(self):
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            content = tmp_path / "content.pptx"
            labels = ["SpaceX 投研叙事 · Part 1 · Thesis", "SpaceX 投研叙事 · Part 2 · Supply"]
            make_content(content, len(labels), section_labels=labels)
            result = merge_pptx(
                content_path=content,
                template_path=TEMPLATE,
                cover_path=TEMPLATE,
                output_path=tmp_path / "final.pptx",
            )

            parts = package(result.output)
            ordered = presentation_slide_parts(parts)
            self.assertEqual(len(ordered[1:-1]), len(labels))
            for body_part, expected_label in zip(ordered[1:-1], labels):
                slide = ET.fromstring(parts[body_part])
                names = [node.get("name") for node in slide.findall(".//p:cNvPr", NS)]
                self.assertEqual(names.count("GHB Template Section Frame"), 1)
                self.assertNotIn("template-section-label", names)
                text = "".join(node.text or "" for node in slide.findall(".//a:t", NS))
                self.assertIn(expected_label, text)
                self.assertNotIn("XXX", text)
                shape_ids = [node.get("id") for node in slide.findall(".//p:cNvPr", NS)]
                self.assertEqual(len(shape_ids), len(set(shape_ids)))
                frame = next(
                    node
                    for node in slide.findall("p:cSld/p:spTree/p:grpSp", NS)
                    if node.find("p:nvGrpSpPr/p:cNvPr", NS).get("name") == "GHB Template Section Frame"
                )
                off = frame.find("p:grpSpPr/a:xfrm/a:off", NS)
                ext = frame.find("p:grpSpPr/a:xfrm/a:ext", NS)
                presentation = ET.fromstring(parts["ppt/presentation.xml"])
                slide_size = presentation.find("p:sldSz", NS)
                self.assertLessEqual(int(off.get("x")) + int(ext.get("cx")), int(slide_size.get("cx")))

    def test_profiled_section_frame_uses_kaiti_and_is_right_flush(self):
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            content = tmp_path / "content.pptx"
            make_content(content, 1, section_labels=["投资判断"])
            default = merge_pptx(
                content_path=content,
                template_path=TEMPLATE,
                cover_path=TEMPLATE,
                output_path=tmp_path / "default.pptx",
            )
            profiled = merge_pptx(
                content_path=content,
                template_path=TEMPLATE,
                cover_path=TEMPLATE,
                output_path=tmp_path / "profiled.pptx",
                section_frame_font="KaiTi",
                section_frame_left_inset_px=24,
            )

            def frame_geometry(path: Path):
                parts = package(path)
                presentation = ET.fromstring(parts["ppt/presentation.xml"])
                width = int(presentation.find("p:sldSz", NS).get("cx"))
                slide = ET.fromstring(parts[presentation_slide_parts(parts)[1]])
                frame = next(
                    node
                    for node in slide.findall("p:cSld/p:spTree/p:grpSp", NS)
                    if node.find("p:nvGrpSpPr/p:cNvPr", NS).get("name")
                    == "GHB Template Section Frame"
                )
                xfrm = frame.find("p:grpSpPr/a:xfrm", NS)
                off = xfrm.find("a:off", NS)
                ext = xfrm.find("a:ext", NS)
                fonts = [
                    node.get("typeface")
                    for node in (
                        frame.findall(".//a:latin", NS)
                        + frame.findall(".//a:ea", NS)
                        + frame.findall(".//a:cs", NS)
                    )
                ]
                return width, int(off.get("x")), int(ext.get("cx")), fonts

            width, default_x, default_width, _ = frame_geometry(default.output)
            profiled_width, profiled_x, profiled_frame_width, fonts = frame_geometry(profiled.output)
            inset = round(24 * width / 1280)
            self.assertEqual(profiled_width, width)
            self.assertEqual(profiled_x + profiled_frame_width, width)
            self.assertEqual(profiled_x - default_x, inset)
            self.assertEqual(default_width - profiled_frame_width, inset)
            self.assertTrue(fonts)
            self.assertTrue(all(font == "KaiTi" for font in fonts))

    def test_body_without_semantic_section_label_does_not_gain_title_frame(self):
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            content = tmp_path / "content.pptx"
            make_content(content, 1)
            result = merge_pptx(
                content_path=content,
                template_path=TEMPLATE,
                cover_path=TEMPLATE,
                output_path=tmp_path / "final.pptx",
            )
            parts = package(result.output)
            body = ET.fromstring(parts[presentation_slide_parts(parts)[1]])
            names = [node.get("name") for node in body.findall(".//p:cNvPr", NS)]
            self.assertNotIn("GHB Template Section Frame", names)

    def test_body_count_matrix_and_default_ending_mounts_all_roles(self):
        for count in (1, 3, 10):
            with self.subTest(body_count=count), tempfile.TemporaryDirectory() as tmp:
                tmp_path = Path(tmp)
                content = tmp_path / "content.pptx"
                output = tmp_path / "final.pptx"
                make_content(content, count)
                result = merge_pptx(
                    content_path=content,
                    template_path=TEMPLATE,
                    cover_path=TEMPLATE,
                    output_path=output,
                    content_layout_index=2,
                )
                self.assertTrue(result.has_ending)
                self.assertEqual(len(Presentation(output).slides), count + 2)
                parts = package(output)
                ordered = presentation_slide_parts(parts)
                self.assertEqual(slide_layout_part(parts, ordered[0]), result.cover_layout_part)
                for body in ordered[1:-1]:
                    self.assertEqual(slide_layout_part(parts, body), result.content_layout_part)
                self.assertEqual(slide_layout_part(parts, ordered[-1]), result.ending_layout_part)

                master_xml = ET.fromstring(parts[result.master_part])
                registered = {
                    node.get(qn("r", "id"))
                    for node in master_xml.findall(".//p:sldLayoutId", NS)
                }
                master_rels = parse_rels(parts, result.master_part)
                layout_rel_ids = {
                    rel.get("Id") for rel in master_rels if relation_type(rel) == "slideLayout"
                }
                self.assertEqual(registered, layout_rel_ids)
                self.assertEqual(len(layout_rel_ids), 3)

    def test_no_ending_and_explicit_ending(self):
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            content = tmp_path / "content.pptx"
            make_content(content, 3)
            no_ending = merge_pptx(
                content_path=content,
                template_path=TEMPLATE,
                cover_path=TEMPLATE,
                output_path=tmp_path / "no-ending.pptx",
                no_ending=True,
            )
            self.assertFalse(no_ending.has_ending)
            self.assertIsNone(no_ending.ending_layout_part)
            self.assertEqual(len(Presentation(no_ending.output).slides), 4)

            explicit = merge_pptx(
                content_path=content,
                template_path=TEMPLATE,
                cover_path=TEMPLATE,
                output_path=tmp_path / "explicit.pptx",
                ending_slide_index=4,
            )
            self.assertTrue(explicit.has_ending)
            self.assertEqual(len(Presentation(explicit.output).slides), 5)

    def test_content_layout_can_share_ending_layout(self):
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            content = tmp_path / "content.pptx"
            make_content(content, 1)
            result = merge_pptx(
                content_path=content,
                template_path=TEMPLATE,
                cover_path=TEMPLATE,
                output_path=tmp_path / "final.pptx",
                content_layout_index=3,
            )
            self.assertEqual(result.content_layout_part, result.ending_layout_part)
            parts = package(result.output)
            master_rels = parse_rels(parts, result.master_part)
            self.assertEqual(sum(relation_type(rel) == "slideLayout" for rel in master_rels), 2)

    def test_media_filename_collision_preserves_content_and_rewires_template(self):
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            content = tmp_path / "content.pptx"
            picture_bytes = make_content(content, 1, with_image=True)
            output = tmp_path / "final.pptx"
            result = merge_pptx(
                content_path=content,
                template_path=TEMPLATE,
                cover_path=TEMPLATE,
                output_path=output,
            )
            parts = package(output)
            self.assertEqual(parts["ppt/media/image1.png"], picture_bytes)
            template_hash = hashlib.sha256(package(TEMPLATE)["ppt/media/image1.png"]).hexdigest()
            copied_template_images = [
                name for name, payload in parts.items()
                if name.startswith("ppt/media/") and hashlib.sha256(payload).hexdigest() == template_hash
            ]
            self.assertTrue(copied_template_images)
            self.assertNotIn("ppt/media/image1.png", copied_template_images)
            master_rels = parse_rels(parts, result.master_part)
            for rel in master_rels:
                if relation_type(rel) == "image":
                    self.assertIn(resolve_target(result.master_part, rel.get("Target", "")), parts)

    def test_ids_content_types_and_package_members_are_unique(self):
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            content = tmp_path / "content.pptx"
            output = tmp_path / "final.pptx"
            make_content(content, 3)
            result = merge_pptx(
                content_path=content,
                template_path=TEMPLATE,
                cover_path=TEMPLATE,
                output_path=output,
            )
            with zipfile.ZipFile(output) as archive:
                names = archive.namelist()
                self.assertEqual(len(names), len(set(names)))
            parts = package(output)
            self.assertIn(
                b'<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">',
                parts["[Content_Types].xml"],
            )
            self.assertNotIn(b"ns0:Types", parts["[Content_Types].xml"])
            presentation = ET.fromstring(parts["ppt/presentation.xml"])
            slide_ids = [node.get("id") for node in presentation.findall(".//p:sldId", NS)]
            self.assertEqual(len(slide_ids), len(set(slide_ids)))
            for name, payload in parts.items():
                if not name.endswith(".rels"):
                    continue
                rels = ET.fromstring(payload)
                ids = [rel.get("Id") for rel in rels]
                self.assertEqual(len(ids), len(set(ids)), name)
            ct = ET.fromstring(parts["[Content_Types].xml"])
            overrides = [node.get("PartName") for node in ct if node.tag.endswith("Override")]
            self.assertEqual(len(overrides), len(set(overrides)))
            expected = {
                "/" + result.master_part: CONTENT_TYPES["master"],
                "/" + result.cover_layout_part: CONTENT_TYPES["layout"],
                "/" + result.content_layout_part: CONTENT_TYPES["layout"],
            }
            actual = {node.get("PartName"): node.get("ContentType") for node in ct if node.tag.endswith("Override")}
            for part_name, content_type in expected.items():
                self.assertEqual(actual.get(part_name), content_type)

    def test_invalid_inputs_fail_without_replacing_existing_output(self):
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            content = tmp_path / "content.pptx"
            output = tmp_path / "final.pptx"
            make_content(content, 1)
            output.write_bytes(b"keep-me")
            with self.assertRaisesRegex(MergeError, "missing part|not found"):
                merge_pptx(
                    content_path=content,
                    template_path=TEMPLATE,
                    cover_path=TEMPLATE,
                    output_path=output,
                    content_layout_index=999,
                )
            self.assertEqual(output.read_bytes(), b"keep-me")
            corrupt = tmp_path / "corrupt.pptx"
            corrupt.write_bytes(b"not a zip")
            with self.assertRaisesRegex(MergeError, "invalid PPTX ZIP"):
                merge_pptx(
                    content_path=corrupt,
                    template_path=TEMPLATE,
                    cover_path=TEMPLATE,
                    output_path=output,
                )


if __name__ == "__main__":
    unittest.main()
