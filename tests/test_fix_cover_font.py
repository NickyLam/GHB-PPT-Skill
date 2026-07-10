from __future__ import annotations

import tempfile
import unittest
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from scripts.fix_cover_font import FontFixError, fix_cover_font


def make_cover(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    paragraph = box.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = "封面标题"
    run.font.name = "楷体"
    presentation.save(path)


def slide_xml(path: Path) -> str:
    with zipfile.ZipFile(path) as archive:
        return archive.read("ppt/slides/slide1.xml").decode("utf-8")


class FixCoverFontTest(unittest.TestCase):
    def test_replaces_cover_font_and_is_idempotent(self):
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "cover.pptx"
            make_cover(path)
            first = fix_cover_font(path)
            self.assertGreaterEqual(first.replacements, 1)
            self.assertIn('typeface="Microsoft YaHei"', slide_xml(path))
            second = fix_cover_font(path)
            self.assertEqual(second.replacements, 0)
            self.assertEqual(second.changed_parts, 0)

    def test_invalid_input_fails_without_overwriting(self):
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "cover.pptx"
            path.write_bytes(b"not-a-pptx")
            with self.assertRaisesRegex(FontFixError, "invalid PPTX ZIP"):
                fix_cover_font(path)
            self.assertEqual(path.read_bytes(), b"not-a-pptx")


if __name__ == "__main__":
    unittest.main()
