from __future__ import annotations

import hashlib
import json
import tempfile
import unittest
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from scripts.merge_template_master import merge_pptx, relation_type, resolve_target
from scripts.validate_ghb_pptx import main, markdown_report, report_dict, validate_pptx


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "templates" / "GHB_PPT_模板.pptx"


def make_cover(path: Path) -> None:
    presentation = Presentation(TEMPLATE)
    slide = presentation.slides[0]
    replacements = ["企业级基线", "离线验证", "2026年7月"]
    index = 0
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or not shape.text.strip():
            continue
        shape.text = replacements[min(index, len(replacements) - 1)]
        index += 1
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "Microsoft YaHei"
    presentation.save(path)


def make_content(path: Path, count: int, *, full_slide_image: bool = False) -> None:
    presentation = Presentation()
    presentation.slide_width = 12192000
    presentation.slide_height = 6858000
    blank = presentation.slide_layouts[6]
    for index in range(count):
        slide = presentation.slides.add_slide(blank)
        if full_slide_image and index == 0:
            image_path = path.with_suffix(".png")
            Image.new("RGB", (1280, 720), "white").save(image_path)
            slide.shapes.add_picture(
                str(image_path), 0, 0,
                width=presentation.slide_width,
                height=presentation.slide_height,
            )
        else:
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
            box.text = f"Body {index + 1} 可编辑正文"
    presentation.save(path)


def rewrite_package(path: Path, transform) -> None:
    with zipfile.ZipFile(path) as archive:
        entries = {info.filename: archive.read(info.filename) for info in archive.infolist() if not info.is_dir()}
    transform(entries)
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as archive:
        for name, payload in entries.items():
            archive.writestr(name, payload)


class ValidateGhbPptxTest(unittest.TestCase):
    def build(self, directory: Path, *, count: int = 3, no_ending: bool = False, full_slide_image: bool = False):
        content = directory / "content.pptx"
        cover = directory / "cover.pptx"
        output = directory / "final.pptx"
        make_content(content, count, full_slide_image=full_slide_image)
        make_cover(cover)
        result = merge_pptx(
            content_path=content,
            template_path=TEMPLATE,
            cover_path=cover,
            output_path=output,
            no_ending=no_ending,
        )
        return output, result

    def test_valid_final_deck_passes_with_per_slide_editability_summary(self):
        with tempfile.TemporaryDirectory() as tmp:
            output, _result = self.build(Path(tmp))
            report = validate_pptx(output, expected_body_count=3, expect_ending=True)
            self.assertTrue(report.passed, report.errors)
            self.assertEqual(report.page_count, 5)
            self.assertEqual([slide.role for slide in report.slides], ["cover", "body", "body", "body", "ending"])
            self.assertTrue(all(slide.text_objects >= 1 for slide in report.slides[1:-1]))
            self.assertEqual(len(report.package["masters_used"]), 1)

    def test_wrong_page_count_is_an_error(self):
        with tempfile.TemporaryDirectory() as tmp:
            output, _result = self.build(Path(tmp), count=1)
            report = validate_pptx(output, expected_body_count=3, expect_ending=True)
            self.assertFalse(report.passed)
            self.assertTrue(any(issue.code == "page-count" for issue in report.errors))

    def test_unregistered_used_ending_layout_is_detected(self):
        with tempfile.TemporaryDirectory() as tmp:
            output, result = self.build(Path(tmp))

            def corrupt(entries: dict[str, bytes]) -> None:
                rels_name = result.master_part.replace("ppt/slideMasters/", "ppt/slideMasters/_rels/") + ".rels"
                root = ET.fromstring(entries[rels_name])
                for rel in list(root):
                    if relation_type(rel) != "slideLayout":
                        continue
                    target = resolve_target(result.master_part, rel.get("Target", ""))
                    if target == result.ending_layout_part:
                        root.remove(rel)
                entries[rels_name] = ET.tostring(root, encoding="utf-8", xml_declaration=True)

            rewrite_package(output, corrupt)
            report = validate_pptx(output, expected_body_count=3, expect_ending=True)
            self.assertFalse(report.passed)
            self.assertTrue(any(issue.code == "unregistered-used-layout" for issue in report.errors))

    def test_full_slide_body_image_without_editable_text_is_rejected(self):
        with tempfile.TemporaryDirectory() as tmp:
            output, _result = self.build(Path(tmp), count=1, full_slide_image=True)
            report = validate_pptx(output, expected_body_count=1, expect_ending=True)
            self.assertFalse(report.passed)
            self.assertTrue(any(issue.code == "full-slide-image-body" for issue in report.errors))

    def test_missing_planned_body_item_is_rejected_even_when_title_exists(self):
        with tempfile.TemporaryDirectory() as tmp:
            directory = Path(tmp)
            output, _result = self.build(directory, count=1)
            plan = directory / "layout_plan.json"
            plan.write_text(
                json.dumps(
                    [
                        {
                            "slide": 1,
                            "key_message": "Body 1 可编辑正文",
                            "items": ["必须出现的卡片标签"],
                        }
                    ],
                    ensure_ascii=False,
                ),
                encoding="utf-8",
            )
            report = validate_pptx(
                output,
                expected_body_count=1,
                expect_ending=True,
                layout_plan_path=plan,
            )
            self.assertFalse(report.passed)
            self.assertTrue(any(issue.code == "missing-planned-item" for issue in report.errors))

    def test_cli_writes_json_and_markdown_reports(self):
        with tempfile.TemporaryDirectory() as tmp:
            directory = Path(tmp)
            output, _result = self.build(directory, count=1, no_ending=True)
            json_path = directory / "report.json"
            markdown_path = directory / "report.md"
            code = main(
                [
                    str(output),
                    "--body-count", "1",
                    "--no-ending",
                    "--json-output", str(json_path),
                    "--markdown-output", str(markdown_path),
                ]
            )
            self.assertEqual(code, 0)
            self.assertTrue(json.loads(json_path.read_text(encoding="utf-8"))["passed"])
            self.assertIn("Per-slide object summary", markdown_path.read_text(encoding="utf-8"))

    def test_json_and_markdown_share_mandated_evidence_hierarchy(self):
        with tempfile.TemporaryDirectory() as tmp:
            output, _result = self.build(Path(tmp), count=1, no_ending=True)
            report = validate_pptx(output, expected_body_count=1, expect_ending=False)
            payload = report_dict(report)
            self.assertEqual(
                list(payload["quality"].keys()),
                [
                    "deterministic_outcome",
                    "freshness",
                    "reviewability",
                    "blocking_findings",
                    "advisory_findings",
                    "per_slide_evidence",
                ],
            )
            self.assertEqual(payload["quality"]["reviewability"]["review_outcome"], "unavailable")
            markdown = markdown_report(report)
            headings = [
                "## Deterministic outcome",
                "## Freshness",
                "## Reviewability and limitations",
                "## Blocking findings",
                "## Advisory findings",
                "## Per-slide evidence and actions",
            ]
            self.assertEqual([markdown.index(item) for item in headings], sorted(markdown.index(item) for item in headings))
            self.assertIn("Review outcome: `unavailable`", markdown)

    def test_markdown_preserves_freshness_issues_from_json_hierarchy(self):
        with tempfile.TemporaryDirectory() as tmp:
            output, _result = self.build(Path(tmp), count=1, no_ending=True)
            freshness = {
                "status": "stale",
                "states": {"pptx": "stale"},
                "issues": [{"code": "evidence-byte-digest-mismatch", "identity": "pptx"}],
            }
            report = validate_pptx(
                output,
                expected_body_count=1,
                expect_ending=False,
                freshness=freshness,
            )
            self.assertEqual(report_dict(report)["quality"]["freshness"], freshness)
            markdown = markdown_report(report)
            self.assertIn("evidence-byte-digest-mismatch", markdown)
            self.assertIn('"identity": "pptx"', markdown)

    def test_successful_render_is_review_skipped_not_visual_pass(self):
        with tempfile.TemporaryDirectory() as tmp:
            directory = Path(tmp)
            output, _result = self.build(directory, count=1, no_ending=True)
            render_dir = directory / "render"
            render_dir.mkdir()
            (render_dir / "slide-01.png").write_bytes(b"png")
            (render_dir / "slide-02.png").write_bytes(b"png")
            (render_dir / "render-report.json").write_text(
                json.dumps({
                    "schema": "ghb.render-report.v1",
                    "status": "passed",
                    "passed": True,
                    "pptx": str(output.resolve()),
                    "pptx_sha256": hashlib.sha256(output.read_bytes()).hexdigest(),
                    "renderer": "soffice",
                    "dpi": 144,
                    "font": {"status": "available", "warnings": []},
                    "outputs": ["render.pdf", "slide-01.png", "slide-02.png"],
                    "warnings": [],
                    "errors": [],
                }),
                encoding="utf-8",
            )
            report = validate_pptx(
                output,
                expected_body_count=1,
                expect_ending=False,
                render_dir=render_dir,
            )
            quality = report_dict(report)["quality"]
            self.assertEqual(quality["reviewability"]["review_outcome"], "skipped")
            self.assertNotEqual(quality["deterministic_outcome"]["status"], "visual-pass")

    def test_final_reports_preserve_detailed_svg_visual_findings(self):
        with tempfile.TemporaryDirectory() as tmp:
            directory = Path(tmp)
            output, _result = self.build(directory, count=1, no_ending=True)
            svg_report = directory / "svg-authored.json"
            svg_report.write_text(
                json.dumps({
                    "stage": "authored",
                    "passed": True,
                    "error_count": 0,
                    "warning_count": 1,
                    "files": [{
                        "file": "01_matrix.svg",
                        "slide_id": "body-01",
                        "visual_findings": [{
                            "code": "visual-component-gap-small",
                            "severity": "warning",
                            "slide_id": "body-01",
                            "evidence": {"minimum_gap": 0},
                            "expected": {"min": 16},
                            "suggested_action": "Separate peer cards.",
                        }],
                    }],
                }),
                encoding="utf-8",
            )
            report = validate_pptx(
                output,
                expected_body_count=1,
                expect_ending=False,
                svg_report_paths=[svg_report],
            )
            payload = report_dict(report)
            advisory = payload["quality"]["advisory_findings"]
            self.assertIn("visual-component-gap-small", {item["code"] for item in advisory})
            body = next(
                item for item in payload["quality"]["per_slide_evidence"]
                if item.get("slide_id") == "body-01"
            )
            self.assertEqual(body["actions"], ["Separate peer cards."])
            self.assertIn("visual-component-gap-small", markdown_report(report))

    def test_render_failure_reason_remains_visible_in_quality_report(self):
        with tempfile.TemporaryDirectory() as tmp:
            directory = Path(tmp)
            output, _result = self.build(directory, count=1, no_ending=True)
            render_dir = directory / "render"
            render_dir.mkdir()
            (render_dir / "render-report.json").write_text(
                json.dumps({
                    "schema": "ghb.render-report.v1",
                    "status": "unavailable",
                    "passed": False,
                    "renderer": "auto",
                    "dpi": 144,
                    "font": {"status": "unknown", "warnings": []},
                    "outputs": [],
                    "warnings": [],
                    "errors": ["no renderer detected"],
                }),
                encoding="utf-8",
            )
            report = validate_pptx(
                output,
                expected_body_count=1,
                expect_ending=False,
                render_dir=render_dir,
            )
            self.assertIn("no renderer detected", "\n".join(report.known_limitations))
            self.assertIn("render-unavailable", {item.code for item in report.warnings})

    def test_failed_render_report_does_not_reuse_old_page_images(self):
        with tempfile.TemporaryDirectory() as tmp:
            directory = Path(tmp)
            output, _result = self.build(directory, count=1, no_ending=True)
            render_dir = directory / "render"
            render_dir.mkdir()
            (render_dir / "slide-01.png").write_bytes(b"old")
            (render_dir / "slide-02.png").write_bytes(b"old")
            (render_dir / "render-report.json").write_text(
                json.dumps({
                    "schema": "ghb.render-report.v1",
                    "status": "error",
                    "passed": False,
                    "renderer": "soffice",
                    "dpi": 144,
                    "font": {"status": "available", "warnings": []},
                    "outputs": [],
                    "warnings": [],
                    "errors": ["conversion failed"],
                }),
                encoding="utf-8",
            )
            report = validate_pptx(
                output,
                expected_body_count=1,
                expect_ending=False,
                render_dir=render_dir,
            )
            self.assertEqual(report.package["rendered_pages"], 0)
            reviewability = report_dict(report)["quality"]["reviewability"]
            self.assertEqual(reviewability["render_status"], "error")
            self.assertEqual(reviewability["provenance"]["outputs"], [])

    def test_successful_render_for_overwritten_pptx_is_stale(self):
        with tempfile.TemporaryDirectory() as tmp:
            directory = Path(tmp)
            output, _result = self.build(directory, count=1, no_ending=True)
            render_dir = directory / "render"
            render_dir.mkdir()
            pages = []
            for index in (1, 2):
                page = render_dir / f"slide-{index:02d}.png"
                page.write_bytes(b"old")
                pages.append(str(page.resolve()))
            (render_dir / "render-report.json").write_text(
                json.dumps({
                    "schema": "ghb.render-report.v1",
                    "status": "passed",
                    "passed": True,
                    "pptx": str(output.resolve()),
                    "pptx_sha256": "0" * 64,
                    "renderer": "soffice",
                    "dpi": 144,
                    "font": {"status": "available", "warnings": []},
                    "outputs": pages,
                    "warnings": [],
                    "errors": [],
                }),
                encoding="utf-8",
            )
            report = validate_pptx(
                output,
                expected_body_count=1,
                expect_ending=False,
                render_dir=render_dir,
            )
            self.assertEqual(report.package["rendered_pages"], 0)
            reviewability = report_dict(report)["quality"]["reviewability"]
            self.assertEqual(reviewability["review_outcome"], "stale")
            self.assertIn("render-pptx-mismatch", {item.code for item in report.warnings})

    def test_stale_freshness_is_blocking(self):
        with tempfile.TemporaryDirectory() as tmp:
            output, _result = self.build(Path(tmp), count=1, no_ending=True)
            report = validate_pptx(
                output,
                expected_body_count=1,
                expect_ending=False,
                freshness={
                    "status": "stale",
                    "states": {"pptx": "stale"},
                    "issues": [{"code": "invalid-freshness-evidence"}],
                },
            )
            self.assertFalse(report.passed)
            self.assertIn("stale-evidence", {item.code for item in report.errors})

    def test_optional_review_projection_preserves_deterministic_priority_and_markdown_parity(self):
        with tempfile.TemporaryDirectory() as tmp:
            directory = Path(tmp)
            output, _result = self.build(directory, count=1, no_ending=True)
            render_dir = directory / "render"
            render_dir.mkdir()
            rendered_pages = []
            for index in (1, 2):
                page = render_dir / f"slide-{index:02d}.png"
                page.write_bytes(b"png")
                rendered_pages.append(str(page.resolve()))
            (render_dir / "render-report.json").write_text(json.dumps({
                "schema": "ghb.render-report.v1", "status": "passed", "passed": True,
                "pptx": str(output.resolve()),
                "pptx_sha256": hashlib.sha256(output.read_bytes()).hexdigest(),
                "renderer": "soffice", "dpi": 144,
                "font": {"status": "available", "warnings": []},
                "page_count": 2, "outputs": rendered_pages,
                "warnings": [], "errors": [],
            }), encoding="utf-8")
            review = directory / "visual-review.json"
            review.write_text(json.dumps({
                "schema": "ghb.visual-review-report.v1",
                "outcome": "needs-revision",
                "deterministic_status": "passed",
                "completion_status": "completed",
                "freshness": "fresh",
                "findings": [{
                    "code": "review-weak-hierarchy", "slide_id": "slide-01",
                    "dimension": "hierarchy", "reviewability": "reviewed",
                    "severity": "advisory",
                    "location": {"x": 0.1, "y": 0.1, "width": 0.2, "height": 0.2},
                    "evidence": "Weak focal contrast.", "action": "Increase contrast.",
                }],
                "dimension_reviewability": [
                    {"dimension": dimension, "status": "reviewed", "limitations": []}
                    for dimension in (
                        "hierarchy", "spacing", "typography", "cjk", "geometry", "composition"
                    )
                ],
                "limitations": [],
                "provenance": {
                    "adapter_sha256": "a" * 64, "launcher_sha256": None,
                    "capability": "local", "model_id": "fixture-model",
                    "tool_contract": "fixture-tool-v1", "credentials": [],
                    "disclosure": None,
                    "direct_subprocess_isolation": "trusted-same-user",
                },
                "request_digest": "b" * 64,
                "reviewer_metadata": {"adapter_version": "fixture-1"},
            }), encoding="utf-8")
            report = validate_pptx(
                output, expected_body_count=1, expect_ending=False,
                render_dir=render_dir,
                review_report_path=review, review_required=True,
                freshness={
                    "status": "fresh",
                    "states": {"adapter-review": "fresh"},
                    "issues": [],
                },
            )
            quality = report_dict(report)["quality"]
            self.assertEqual(quality["deterministic_outcome"]["status"], "passed")
            self.assertEqual(quality["reviewability"]["review_outcome"], "needs-revision")
            self.assertTrue(quality["reviewability"]["requirement_satisfied"])
            self.assertEqual(quality["advisory_findings"][-1]["severity"], "warning")
            markdown = markdown_report(report)
            self.assertLess(markdown.index("## Blocking findings"), markdown.index("review-weak-hierarchy"))
            self.assertIn("Review outcome: `needs-revision`", markdown)

            unsafe = json.loads(review.read_text(encoding="utf-8"))
            unsafe["findings"][0]["evidence"] = "../../secret"
            unsafe["findings"][0]["action"] = "open /tmp/private"
            review.write_text(json.dumps(unsafe), encoding="utf-8")
            rejected = validate_pptx(
                output, expected_body_count=1, expect_ending=False,
                render_dir=render_dir, review_report_path=review, review_required=True,
                freshness={
                    "status": "fresh",
                    "states": {"adapter-review": "fresh"},
                    "issues": [],
                },
            )
            rejected_text = json.dumps(report_dict(rejected)) + markdown_report(rejected)
            self.assertFalse(rejected.quality["reviewability"]["requirement_satisfied"])
            self.assertNotIn("../../secret", rejected_text)
            self.assertNotIn("/tmp/private", rejected_text)

    def test_required_review_error_does_not_reclassify_deterministic_outcome(self):
        with tempfile.TemporaryDirectory() as tmp:
            directory = Path(tmp)
            output, _result = self.build(directory, count=1, no_ending=True)
            review = directory / "visual-review.json"
            review.write_text(json.dumps({
                "schema": "ghb.visual-review-report.v1", "outcome": "error",
                "freshness": "fresh", "findings": [], "dimension_reviewability": [],
                "limitations": ["adapter-timeout"], "provenance": {},
                "error": {"code": "adapter-timeout", "message": "did not complete"},
            }), encoding="utf-8")
            report = validate_pptx(
                output, expected_body_count=1, expect_ending=False,
                review_report_path=review, review_required=True,
            )
            quality = report_dict(report)["quality"]
            self.assertEqual(quality["deterministic_outcome"]["status"], "passed")
            self.assertFalse(quality["reviewability"]["requirement_satisfied"])
            self.assertEqual(
                main([
                    str(output),
                    "--body-count", "1",
                    "--no-ending",
                    "--review-report", str(review),
                    "--review-required",
                ]),
                1,
            )

    def test_malicious_optional_advisory_is_not_injected_or_promoted(self):
        with tempfile.TemporaryDirectory() as tmp:
            directory = Path(tmp)
            output, _result = self.build(directory, count=1, no_ending=True)
            review = directory / "visual-review.json"
            review.write_text(json.dumps({
                "schema": "ghb.visual-review-report.v1", "outcome": "passed",
                "freshness": "fresh", "findings": [{
                    "code": "review-malicious", "slide_id": "slide-01",
                    "dimension": "hierarchy", "reviewability": "reviewed",
                    "severity": "error", "location": {"x":0,"y":0,"width":1,"height":1},
                    "evidence": "<script>alert(1)</script>", "action": "[run](file:///tmp/x)",
                }], "dimension_reviewability": [], "limitations": [], "provenance": {},
            }), encoding="utf-8")
            report = validate_pptx(
                output, expected_body_count=1, expect_ending=False,
                review_report_path=review,
            )
            rendered = json.dumps(report_dict(report)) + markdown_report(report)
            self.assertNotIn("<script>", rendered)
            self.assertNotIn("file:///tmp/x", rendered)
            self.assertEqual(report.quality["reviewability"]["review_outcome"], "error")

    def test_optional_review_rejects_unknown_provenance_and_fabricated_slide(self):
        with tempfile.TemporaryDirectory() as tmp:
            directory = Path(tmp)
            output, _result = self.build(directory, count=1, no_ending=True)
            review = directory / "visual-review.json"
            secret = "credential-value-123"
            review.write_text(json.dumps({
                "schema": "ghb.visual-review-report.v1", "outcome": "passed",
                "freshness": "fresh", "findings": [{
                    "code": "review-fabricated-slide", "slide_id": "slide-99",
                    "dimension": "hierarchy", "reviewability": "reviewed",
                    "severity": "advisory",
                    "location": {"x": 0, "y": 0, "width": 1, "height": 1},
                    "evidence": "Fabricated evidence.", "action": "Ignore it.",
                }], "dimension_reviewability": [], "limitations": [],
                "provenance": {"api_token": secret},
            }), encoding="utf-8")
            report = validate_pptx(
                output, expected_body_count=1, expect_ending=False,
                review_report_path=review,
            )
            rendered = json.dumps(report_dict(report)) + markdown_report(report)
            self.assertEqual(report.quality["reviewability"]["review_outcome"], "error")
            self.assertNotIn(secret, rendered)
            self.assertNotIn("slide-99", rendered)

    def test_readback_markdown_is_required_when_requested(self):
        with tempfile.TemporaryDirectory() as tmp:
            directory = Path(tmp)
            output, _result = self.build(directory, count=1)
            missing = directory / "missing-readback.md"
            report = validate_pptx(
                output,
                expected_body_count=1,
                expect_ending=True,
                readback_markdown_path=missing,
            )
            self.assertFalse(report.passed)
            self.assertTrue(any(issue.code == "missing-ppt-readback" for issue in report.errors))

    def test_readback_markdown_page_count_is_verified(self):
        with tempfile.TemporaryDirectory() as tmp:
            directory = Path(tmp)
            output, _result = self.build(directory, count=1)
            readback = directory / "readback.md"
            readback.write_text("# final\n\n## Slide 1\n\ncover only\n", encoding="utf-8")
            report = validate_pptx(
                output,
                expected_body_count=1,
                expect_ending=True,
                readback_markdown_path=readback,
            )
            self.assertFalse(report.passed)
            self.assertTrue(any(issue.code == "ppt-readback-page-count" for issue in report.errors))


if __name__ == "__main__":
    unittest.main()
