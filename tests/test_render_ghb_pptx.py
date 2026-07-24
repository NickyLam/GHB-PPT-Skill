from __future__ import annotations

import tempfile
import unittest
import shutil
import json
from pathlib import Path
from unittest import mock

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from scripts.render_ghb_pptx import (
    RenderCommand,
    RenderError,
    _font_warning,
    _render_environment,
    build_parser,
    make_contact_sheet,
    render_pptx,
)


class RenderGhbPptxTest(unittest.TestCase):
    def test_contact_sheet_preserves_all_pages_and_labels(self):
        with tempfile.TemporaryDirectory() as tmp:
            directory = Path(tmp)
            pages = []
            for index, color in enumerate(("red", "green", "blue", "white"), 1):
                path = directory / f"page-{index}.png"
                Image.new("RGB", (640, 360), color).save(path)
                pages.append(path)
            output = make_contact_sheet(pages, directory / "contact.png", columns=3, thumb_width=320)
            self.assertTrue(output.is_file())
            with Image.open(output) as image:
                self.assertGreater(image.width, 3 * 320)
                self.assertGreater(image.height, 2 * 180)

    def test_contact_sheet_rejects_empty_page_list(self):
        with tempfile.TemporaryDirectory() as tmp:
            with self.assertRaisesRegex(RenderError, "without page images"):
                make_contact_sheet([], Path(tmp) / "contact.png")

    def test_source_han_sans_sc_satisfies_renderer_font_probe(self):
        probe = mock.Mock(stdout="Source Han Sans SC 思源黑体\n", returncode=0)
        with (
            mock.patch("scripts.render_ghb_pptx.shutil.which", return_value="/usr/bin/fc-list"),
            mock.patch("scripts.render_ghb_pptx.subprocess.run", return_value=probe),
        ):
            self.assertIsNone(_font_warning())

    def test_renderer_environment_uses_writable_cache_and_host_fontconfig(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            config = root / "fonts.conf"
            config.write_text("<fontconfig/>", encoding="utf-8")
            cache = root / "cache"
            with (
                mock.patch("scripts.render_ghb_pptx.FONTCONFIG_FILES", (config,)),
                mock.patch.dict("scripts.render_ghb_pptx.os.environ", {}, clear=True),
            ):
                env = _render_environment(cache)
            self.assertEqual(env["XDG_CACHE_HOME"], str(cache))
            self.assertEqual(env["FONTCONFIG_FILE"], str(config))
            self.assertEqual(env["FONTCONFIG_PATH"], str(root))

    def test_explicit_font_file_gets_a_private_fontconfig_environment(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            base_config = root / "fonts.conf"
            base_config.write_text("<fontconfig/>", encoding="utf-8")
            font_file = root / "Kaiti.ttf"
            font_file.write_bytes(b"font payload")
            cache = root / "cache"
            with (
                mock.patch("scripts.render_ghb_pptx.FONTCONFIG_FILES", (base_config,)),
                mock.patch.dict("scripts.render_ghb_pptx.os.environ", {}, clear=True),
            ):
                env = _render_environment(cache, font_paths=[font_file])
            generated_config = Path(env["FONTCONFIG_FILE"])
            self.assertNotEqual(generated_config, base_config)
            config = generated_config.read_text(encoding="utf-8")
            self.assertIn(str(font_file.parent), config)
            self.assertIn(str(cache), config)
            self.assertEqual(env["FONTCONFIG_PATH"], str(root))

    def test_render_cli_accepts_repeatable_font_file(self):
        parsed = build_parser().parse_args(
            ["deck.pptx", "--output-dir", "render", "--font-file", "Kaiti.ttf"]
        )
        self.assertEqual(parsed.font_files, [Path("Kaiti.ttf")])

    def test_missing_renderer_still_writes_atomic_failure_report(self):
        with tempfile.TemporaryDirectory() as tmp:
            directory = Path(tmp)
            pptx = directory / "deck.pptx"
            pptx.write_bytes(b"not-opened-before-renderer-probe")
            output = directory / "render"
            with mock.patch("scripts.render_ghb_pptx.detect_renderer", side_effect=RenderError("missing renderer")):
                report = render_pptx(pptx, output, dpi=144)
            self.assertFalse(report.passed)
            payload = json.loads((output / "render-report.json").read_text(encoding="utf-8"))
            self.assertEqual(payload["schema"], "ghb.render-report.v1")
            self.assertEqual(payload["status"], "unavailable")
            self.assertEqual(payload["dpi"], 144)
            self.assertEqual(payload["outputs"], [])
            self.assertIn("missing renderer", payload["errors"][0])
            self.assertFalse(any(output.glob(".render-report.json.*.tmp")))

    def test_conversion_failure_persists_renderer_font_and_empty_outputs(self):
        with tempfile.TemporaryDirectory() as tmp:
            directory = Path(tmp)
            pptx = directory / "deck.pptx"
            pptx.write_bytes(b"placeholder")
            output = directory / "render"
            failed = RenderCommand([], 1, "", "conversion failed", 0.1)
            with (
                mock.patch("scripts.render_ghb_pptx.detect_renderer", return_value=("soffice", "/fake/soffice")),
                mock.patch("scripts.render_ghb_pptx.shutil.which", return_value="/fake/pdftoppm"),
                mock.patch("scripts.render_ghb_pptx._font_warning", return_value="font substitution possible"),
                mock.patch("scripts.render_ghb_pptx._run", return_value=failed),
            ):
                report = render_pptx(pptx, output, dpi=96)
            self.assertEqual(report.status, "error")
            payload = json.loads((output / "render-report.json").read_text(encoding="utf-8"))
            self.assertEqual(payload["renderer"], "soffice")
            self.assertEqual(payload["dpi"], 96)
            self.assertEqual(payload["font"]["status"], "limited")
            self.assertEqual(payload["font"]["limitation_codes"], ["target-font-missing"])
            self.assertEqual(payload["outputs"], [])

    def test_invalid_input_is_error_not_renderer_unavailable(self):
        with tempfile.TemporaryDirectory() as tmp:
            output = Path(tmp) / "render"
            report = render_pptx(Path(tmp) / "missing.pptx", output, dpi=144)
            self.assertEqual(report.status, "error")
            self.assertIn("PPTX not found", report.errors[0])

    @unittest.skipUnless(shutil.which("soffice") and shutil.which("pdftoppm"), "render tools unavailable")
    def test_soffice_integration_writes_pdf_png_contact_sheet_and_report(self):
        with tempfile.TemporaryDirectory() as tmp:
            directory = Path(tmp)
            pptx = directory / "one-slide.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text = "Render smoke test"
            presentation.save(pptx)
            report = render_pptx(pptx, directory / "render", dpi=96)
            self.assertTrue(report.passed, report.errors)
            self.assertEqual(report.page_count, 1)
            self.assertTrue(Path(report.pdf).is_file())
            self.assertTrue(Path(report.pages[0]).is_file())
            self.assertTrue(Path(report.contact_sheet).is_file())
            self.assertTrue((directory / "render" / "render-report.json").is_file())


if __name__ == "__main__":
    unittest.main()
