from __future__ import annotations

import importlib.util
import json
import tempfile
import unittest
from unittest import mock
from collections import Counter, defaultdict
import hashlib
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches

from scripts.ghb_visual_quality import evaluate_final_gate
from scripts.render_ghb_pptx import make_contact_sheet


ROOT = Path(__file__).resolve().parents[1]
CASES_PATH = ROOT / "tests" / "fixtures" / "visual_quality_cases.json"
PREFERENCES_PATH = ROOT / "tests" / "fixtures" / "visual_preferences.json"
BUILDER_PATH = ROOT / "tests" / "fixtures" / "build_visual_benchmark.py"
FINAL_CONTRACT_PATH = ROOT / "tests" / "fixtures" / "final_deterministic_contract.json"
SCENARIOS_PATH = ROOT / "tests" / "fixtures" / "scenarios.json"


def load_builder():
    spec = importlib.util.spec_from_file_location("ghb_visual_benchmark", BUILDER_PATH)
    assert spec and spec.loader
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


class VisualBenchmarkContractTest(unittest.TestCase):
    def setUp(self):
        self.corpus = json.loads(CASES_PATH.read_text(encoding="utf-8"))
        self.preferences = json.loads(PREFERENCES_PATH.read_text(encoding="utf-8"))
        self.final_contract = json.loads(
            FINAL_CONTRACT_PATH.read_text(encoding="utf-8")
        )
        self.scenarios = json.loads(SCENARIOS_PATH.read_text(encoding="utf-8"))

    def test_corpus_is_stratified_and_partitioned(self):
        cases = self.corpus["cases"]
        self.assertGreaterEqual(len(cases), 30)
        self.assertEqual(len({case["case_id"] for case in cases}), len(cases))
        expected_purposes = {
            "architecture", "process", "comparison", "timeline", "metrics", "summary"
        }
        self.assertEqual({case["page_purpose"] for case in cases}, expected_purposes)
        source_ids = {
            (case["source"]["scenario_id"], case["source"]["body_slide_index"])
            for case in cases
        }
        self.assertEqual(len(source_ids), len(cases), "semantic sources must not leak across partitions")

        by_purpose = defaultdict(Counter)
        for case in cases:
            by_purpose[case["page_purpose"]][case["partition"]] += 1
            self.assertIn(case["density"], {"breathing", "anchor", "dense"})
            self.assertIn(case["content_length_band"], {"short", "medium", "long"})
            self.assertIn(case["item_count_band"], {"few", "standard", "many"})
            self.assertEqual(case["source"]["scenario_file"], "tests/fixtures/scenarios.json")
            self.assertIn("renderer", case["pre_change_evidence"]["provenance"])
            self.assertIn("dpi", case["pre_change_evidence"]["provenance"])
            self.assertIn("fonts", case["pre_change_evidence"]["provenance"])
            self.assertEqual(len(case["pre_change_evidence"]["authored_svg_sha256"]), 64)
            self.assertIsNone(case["pre_change_evidence"]["rendered_png"])
            self.assertEqual(
                case["pre_change_evidence"]["render_status"],
                "unavailable-without-render",
            )
        for purpose in expected_purposes:
            self.assertGreaterEqual(by_purpose[purpose]["calibration"], 2)
            self.assertGreaterEqual(by_purpose[purpose]["pilot-holdout"], 1)
            self.assertGreaterEqual(by_purpose[purpose]["final-holdout"], 2)
        self.assertEqual({case["density"] for case in cases}, {"breathing", "anchor", "dense"})
        self.assertEqual({case["content_length_band"] for case in cases}, {"short", "medium", "long"})
        self.assertEqual({case["item_count_band"] for case in cases}, {"few", "standard", "many"})

    def test_partition_digest_and_consumer_views_are_frozen(self):
        module = load_builder()
        module.validate_contracts(self.corpus, self.preferences)
        self.assertEqual(
            module.partition_digest(self.corpus["cases"]),
            self.corpus["partition_assignment_sha256"],
        )
        tuning = module.cases_for_consumer(self.corpus, "tuning")
        pilot = module.cases_for_consumer(self.corpus, "u11-pilot")
        final = module.cases_for_consumer(self.corpus, "final-evaluation")
        self.assertEqual({case["partition"] for case in tuning}, {"calibration"})
        self.assertEqual({case["partition"] for case in pilot}, {"pilot-holdout"})
        self.assertEqual({case["partition"] for case in final}, {"final-holdout"})
        with self.assertRaisesRegex(ValueError, "unknown benchmark consumer"):
            module.cases_for_consumer(self.corpus, "renderer")

    def test_builder_refuses_to_overwrite_and_needs_no_secret(self):
        module = load_builder()
        self.assertEqual(
            len(list(module.PRECHANGE_SVG_DIR.glob("*.svg"))),
            len(self.corpus["cases"]),
        )
        with tempfile.TemporaryDirectory() as tmp:
            output = Path(tmp) / "frozen"
            result = module.build(output, self.corpus, self.preferences)
            self.assertEqual(result["case_count"], len(self.corpus["cases"]))
            manifest = json.loads((output / "benchmark-manifest.json").read_text(encoding="utf-8"))
            self.assertEqual(manifest["network"], "disabled-not-required")
            self.assertEqual(manifest["model_adapter"], "not-used")
            self.assertEqual(manifest["credentials"], "not-used")
            self.assertEqual(len(manifest["generated_authored_svg"]), len(self.corpus["cases"]))
            self.assertEqual(len(list((output / "authored-svg").glob("*.svg"))), len(self.corpus["cases"]))
            with self.assertRaisesRegex(FileExistsError, "refusing to overwrite"):
                module.build(output, self.corpus, self.preferences)

    def test_source_identity_reuse_is_rejected(self):
        module = load_builder()
        changed = json.loads(json.dumps(self.corpus))
        changed["cases"][1]["source"] = dict(changed["cases"][0]["source"])
        with self.assertRaisesRegex(ValueError, "source identity"):
            module.validate_contracts(changed, self.preferences)

    def test_blind_review_protocol_and_frozen_false_positive_denominator(self):
        protocol = self.preferences["blind_review_protocol"]
        self.assertTrue(protocol["pair_order_randomized"])
        self.assertTrue(protocol["identity_anonymized"])
        self.assertGreaterEqual(protocol["minimum_independent_eligible_reviewers_per_page"], 3)
        self.assertEqual(protocol["tie_handling"], "exclude-from-preference-denominator")
        self.assertEqual(protocol["abstention_handling"], "exclude-from-preference-denominator")
        self.assertEqual(protocol["aggregation_order"], ["per-page", "across-pages"])
        self.assertEqual(protocol["structural_regression"], "page-level-veto")
        self.assertTrue(protocol["reject_duplicate_reviewer_page_judgments"])
        self.assertGreaterEqual(protocol["minimum_eligible_non_tie_judgments_per_page"], 2)

        gate = self.preferences["deterministic_pilot_false_positive_gate"]
        self.assertEqual(gate["blocking_false_positive_maximum"], 0)
        self.assertEqual(gate["advisory_false_positive_rate_maximum"], 0.10)
        self.assertEqual(gate["advisory_denominator_count"], len(gate["advisory_rule_case_pairs"]))
        self.assertEqual(len(gate["advisory_rule_case_pairs"]), len(set(gate["advisory_rule_case_pairs"])))
        module = load_builder()
        self.assertEqual(
            module.advisory_denominator_digest(gate["advisory_rule_case_pairs"]),
            gate["advisory_denominator_sha256"],
        )

    def test_final_evaluation_bundle_is_partition_isolated_blind_and_pending(self):
        module = load_builder()
        with tempfile.TemporaryDirectory() as tmp:
            output = Path(tmp) / "final"
            manifest = module.build_final_evaluation(output, self.corpus, self.preferences)

            final_ids = {
                case["case_id"]
                for case in self.corpus["cases"]
                if case["partition"] == "final-holdout"
            }
            self.assertEqual(manifest["decision"], "pending")
            self.assertEqual(manifest["eligible_case_count"], 12)
            self.assertEqual(set(manifest["eligible_case_ids"]), final_ids)
            self.assertEqual(len(list((output / "before").glob("*.svg"))), 12)
            self.assertEqual(len(list((output / "after").glob("*.svg"))), 12)
            self.assertEqual(len(list((output / "blind").glob("*/*.svg"))), 24)
            for blind_svg in (output / "blind").glob("*/*.svg"):
                self.assertNotIn("data-", blind_svg.read_text(encoding="utf-8"))
            self.assertEqual(
                {row["page_purpose"] for row in manifest["case_metrics"]},
                {"architecture", "process", "comparison", "timeline", "metrics", "summary"},
            )
            self.assertEqual(
                {row["family"] for row in manifest["case_metrics"]},
                {"iceberg", "layered_arch", "waterfall", "swimlane", "matrix", "timeline", "flywheel"},
            )

            public = json.loads((output / "blind-review-template.json").read_text(encoding="utf-8"))
            private = json.loads((output / "evaluator-record.json").read_text(encoding="utf-8"))
            self.assertNotIn("roles", json.dumps(public))
            self.assertNotIn("baseline", json.dumps(public))
            self.assertNotIn("optimized", json.dumps(public))
            self.assertEqual(public["judgments"], [])
            self.assertEqual(private["judgments"], [])
            self.assertEqual(private["decision"], "pending")
            self.assertEqual(private["reviewer_eligibility_roster"], [])
            self.assertEqual(
                public["eligibility_attestation_template"],
                {
                    "reviewer_id_hash": "<stable-opaque-reviewer-id>",
                    "independent": True,
                    "did_not_author": True,
                    "did_not_tune": True,
                },
            )
            self.assertTrue(all("roles" in pair for pair in private["pair_assignments"]))
            seed = private["randomization"]["seed"]
            self.assertGreaterEqual(len(seed), 64)
            self.assertNotIn(seed, json.dumps(public))
            self.assertNotIn("seed", json.dumps(public))
            self.assertNotIn("nonce", json.dumps(public))
            for assignment in private["pair_assignments"]:
                roles, left, right = module._blind_assignment(
                    assignment["page_case_id"], seed
                )
                self.assertEqual(roles, assignment["roles"])
                self.assertEqual((left, right), (
                    assignment["masked_left_id"], assignment["masked_right_id"]
                ))
            self.assertEqual(
                private["structural_evidence"]["contact_sheet"]["status"],
                "unavailable",
            )
            self.assertIn(
                "office-clipping",
                public["unreviewable_dimensions"],
            )
            self.assertEqual(manifest["network"], "not-used")
            self.assertEqual(manifest["model_adapter"], "not-used")
            self.assertEqual(manifest["credentials"], "not-used")
            with self.assertRaisesRegex(FileExistsError, "refusing to overwrite"):
                module.build_final_evaluation(output, self.corpus, self.preferences)

    def test_final_evaluation_does_not_read_adapter_or_secret_environment(self):
        module = load_builder()
        with tempfile.TemporaryDirectory() as tmp, mock.patch.dict(
            "os.environ",
            {"OPENAI_API_KEY": "must-not-be-read", "GHB_VISUAL_REVIEW_ADAPTER": "must-not-run"},
            clear=False,
        ), mock.patch("subprocess.run", side_effect=AssertionError("offline builder invoked a process")):
            manifest = module.build_final_evaluation(
                Path(tmp) / "final", self.corpus, self.preferences
            )
        self.assertEqual(manifest["decision"], "pending")

    def test_final_schema_does_not_invent_focal_intent_from_density(self):
        module = load_builder()
        schema = module._final_page_schema(
            {"density": "anchor", "page_purpose": "timeline"},
            {
                "layout_type": "timeline",
                "items": ["准备", "试点", "复盘", "规模化"],
            },
        )

        self.assertEqual(schema["density"], "balanced")
        self.assertEqual(schema["emphasis"], "distributed")
        self.assertNotIn("focal_target", schema)

    def test_deterministic_fixture_expectations_use_real_codes_or_limitations(self):
        module = load_builder()
        report = module.evaluate_deterministic_fixtures()
        rows = {row["fixture_id"]: row for row in report["fixtures"]}
        self.assertEqual(report["schema"], "ghb.visual-benchmark-deterministic.v1")
        self.assertEqual(rows["approved-balanced"]["unexpected_issue_codes"], [])
        self.assertEqual(rows["fake-diversity"]["observed_issue_codes"], ["visual-composition-repeated"])
        self.assertIn("visual-occupancy-below-min", rows["underscaled-content"]["observed_issue_codes"])
        self.assertIn("visual-explicit-bounds-violation", rows["long-title-overflow"]["observed_issue_codes"])
        self.assertIn("visual-primary-color-overuse", rows["primary-color-overuse"]["observed_issue_codes"])
        self.assertIn("visual-focal-zone-streak", rows["repeated-focal-zones"]["observed_issue_codes"])
        self.assertEqual(rows["empty-text"]["expected_issue_codes"], ["empty-text-box"])
        self.assertEqual(rows["empty-text"]["observed_issue_codes"], ["empty-text-box"])
        self.assertEqual(rows["empty-text"]["source"], "validate_ghb_pptx.validate_pptx")
        self.assertEqual(
            rows["font-limited-rendering"]["observed_issue_codes"],
            ["target-font-missing"],
        )
        self.assertEqual(
            rows["font-limited-rendering"]["source"],
            "render_ghb_pptx.font_evidence",
        )
        for row in report["fixtures"]:
            self.assertIn("metric_tolerances", row)
            self.assertFalse(row.get("fabricated_pass", False))
            self.assertEqual(row["unexpected_issue_codes"], [])
            self.assertEqual(row["missing_expected_issue_codes"], [])
            self.assertEqual(row["metric_regressions"], [])

    def test_deterministic_fixture_metrics_are_stable_within_declared_tolerances(self):
        module = load_builder()
        report = module.evaluate_deterministic_fixtures()
        for row in report["fixtures"]:
            for metric, tolerance in row["metric_tolerances"].items():
                self.assertAlmostEqual(
                    row["metrics"][metric],
                    row["expected_metrics"][metric],
                    delta=tolerance,
                )

    def test_final_gate_requires_eligibility_and_complete_structural_evidence(self):
        module = load_builder()
        bundle_tmp = tempfile.TemporaryDirectory()
        self.addCleanup(bundle_tmp.cleanup)
        output = Path(bundle_tmp.name) / "final"
        module.build_final_evaluation(output, self.corpus, self.preferences)
        review = json.loads((output / "evaluator-record.json").read_text(encoding="utf-8"))
        deterministic = json.loads(
            (output / "deterministic-fixtures.json").read_text(encoding="utf-8")
        )
        reviewers = [f"fixture-reviewer-{index}" for index in range(3)]
        review["reviewer_eligibility_roster"] = [
            {
                "reviewer_id_hash": reviewer,
                "independent": True,
                "did_not_author": True,
                "did_not_tune": True,
            }
            for reviewer in reviewers
        ]
        evidence_root = output
        pptx = evidence_root / "final.pptx"
        presentation = Presentation()
        final_cases = [
            case for case in self.corpus["cases"]
            if case["partition"] == "final-holdout"
        ]
        for case in final_cases:
            source = case["source"]
            slide_data = self.scenarios[source["scenario_id"]]["slides"][
                source["body_slide_index"] - 1
            ]
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            text_box = slide.shapes.add_textbox(
                Inches(1), Inches(1), Inches(10), Inches(5)
            )
            text_box.text = "\n".join(
                [slide_data["key_message"], *slide_data["items"]]
            )
        presentation.save(pptx)
        pages = []
        page_images = []
        for index in range(1, 13):
            page = evidence_root / f"slide-{index:02d}.png"
            image = Image.new("RGB", (800, 450), "white")
            ImageDraw.Draw(image).rectangle(
                (20 * index, 30, 20 * index + 120, 180), fill=(index * 10, 40, 80)
            )
            image.save(page)
            pages.append(str(page))
            page_images.append(image)
        contact_sheet = evidence_root / "contact-sheet.png"
        make_contact_sheet([Path(page) for page in pages], contact_sheet, columns=3)
        pdf = evidence_root / "render.pdf"
        page_images[0].save(
            pdf, "PDF", save_all=True, append_images=page_images[1:]
        )
        render_report = evidence_root / "render-report.json"
        render_report.write_text(
            json.dumps({
                "schema": "ghb.render-report.v1",
                "status": "passed",
                "passed": True,
                "pptx": str(pptx),
                "pptx_sha256": hashlib.sha256(pptx.read_bytes()).hexdigest(),
                "page_count": 12,
                "pages": pages,
                "pdf": str(pdf),
                "contact_sheet": str(contact_sheet),
                "outputs": [str(pdf), *pages, str(contact_sheet)],
                "font": {"status": "available", "warnings": []},
            }),
            encoding="utf-8",
        )
        review["structural_evidence"] = {}
        for key, artifact in {
            "pptx": pptx,
            "render": render_report,
            "target_fonts": render_report,
            "contact_sheet": contact_sheet,
        }.items():
            review["structural_evidence"][key] = {
                "status": "available",
                "artifact_path": artifact.name,
                "artifact_sha256": hashlib.sha256(artifact.read_bytes()).hexdigest(),
            }
        review["judgments"] = []
        for assignment in review["pair_assignments"]:
            optimized = next(
                masked for masked, role in assignment["roles"].items() if role == "optimized"
            )
            for reviewer in reviewers:
                review["judgments"].append({
                    "reviewer_id_hash": reviewer,
                    "page_case_id": assignment["page_case_id"],
                    "masked_left_id": assignment["masked_left_id"],
                    "masked_right_id": assignment["masked_right_id"],
                    "presented_order": [
                        assignment["masked_left_id"], assignment["masked_right_id"]
                    ],
                    "judgment": optimized,
                    "rubric": {dimension: 4 for dimension in self.preferences["blind_review_protocol"]["rubric_dimensions"]},
                    "structural_veto": False,
                    "recorded_at": "2026-07-16T00:00:00Z",
                })
        passed = evaluate_final_gate(
            self.preferences,
            review,
            deterministic,
            corpus=self.corpus,
            scenarios=self.scenarios,
            fixture_contract=self.final_contract,
            evidence_root=evidence_root,
        )
        self.assertEqual(passed["decision"], "passed")
        self.assertEqual(passed["preference_rate"], 1.0)

        missing = json.loads(json.dumps(review))
        missing["reviewer_eligibility_roster"] = []
        with self.assertRaisesRegex(ValueError, "missing-reviewer-eligibility"):
            evaluate_final_gate(
                self.preferences, missing, deterministic, corpus=self.corpus,
                scenarios=self.scenarios,
                fixture_contract=self.final_contract,
                evidence_root=evidence_root
            )
        false_attestation = json.loads(json.dumps(review))
        false_attestation["reviewer_eligibility_roster"][0]["did_not_tune"] = False
        with self.assertRaisesRegex(ValueError, "invalid-reviewer-eligibility"):
            evaluate_final_gate(
                self.preferences,
                false_attestation,
                deterministic,
                corpus=self.corpus,
                scenarios=self.scenarios,
                fixture_contract=self.final_contract,
                evidence_root=evidence_root,
            )
        duplicate = json.loads(json.dumps(review))
        duplicate["reviewer_eligibility_roster"].append(
            dict(duplicate["reviewer_eligibility_roster"][0])
        )
        with self.assertRaisesRegex(ValueError, "duplicate-reviewer-eligibility"):
            evaluate_final_gate(
                self.preferences, duplicate, deterministic, corpus=self.corpus,
                scenarios=self.scenarios,
                fixture_contract=self.final_contract,
                evidence_root=evidence_root
            )
        duplicate_page = json.loads(json.dumps(review))
        duplicate_page["judgments"].append(dict(duplicate_page["judgments"][0]))
        with self.assertRaisesRegex(ValueError, "duplicate-reviewer-page-judgment"):
            evaluate_final_gate(
                self.preferences,
                duplicate_page,
                deterministic,
                corpus=self.corpus,
                scenarios=self.scenarios,
                fixture_contract=self.final_contract,
                evidence_root=evidence_root,
            )
        incomplete_evidence = json.loads(json.dumps(review))
        incomplete_evidence["structural_evidence"]["contact_sheet"] = {
            "status": "unavailable", "limitations": ["no-office-render"]
        }
        pending = evaluate_final_gate(
            self.preferences,
            incomplete_evidence,
            deterministic,
            corpus=self.corpus,
            scenarios=self.scenarios,
            fixture_contract=self.final_contract,
            evidence_root=evidence_root,
        )
        self.assertEqual(pending["decision"], "pending")
        self.assertIn("structural-evidence-incomplete", pending["reasons"])

        forged = json.loads(json.dumps(review))
        forged["structural_evidence"]["pptx"]["artifact_sha256"] = "a" * 64
        with self.assertRaisesRegex(ValueError, "structural-evidence-digest-mismatch"):
            evaluate_final_gate(
                self.preferences, forged, deterministic, corpus=self.corpus,
                scenarios=self.scenarios,
                fixture_contract=self.final_contract,
                evidence_root=evidence_root
            )

        tampered_roles = json.loads(json.dumps(review))
        first = tampered_roles["pair_assignments"][0]
        first["roles"] = {
            masked: "baseline" if role == "optimized" else "optimized"
            for masked, role in first["roles"].items()
        }
        with self.assertRaisesRegex(ValueError, "final-role-assignment-mismatch"):
            evaluate_final_gate(
                self.preferences,
                tampered_roles,
                deterministic,
                corpus=self.corpus,
                scenarios=self.scenarios,
                fixture_contract=self.final_contract,
                evidence_root=evidence_root,
            )

        truncated = json.loads(json.dumps(deterministic))
        truncated["fixtures"] = truncated["fixtures"][:1]
        with self.assertRaisesRegex(ValueError, "final-deterministic-contract-mismatch"):
            evaluate_final_gate(
                self.preferences,
                review,
                truncated,
                corpus=self.corpus,
                scenarios=self.scenarios,
                fixture_contract=self.final_contract,
                evidence_root=evidence_root,
            )

        review_backup = json.loads(json.dumps(review))
        pptx_backup = pptx.read_bytes()
        render_backup = render_report.read_text(encoding="utf-8")
        unrelated = Presentation()
        for _ in range(12):
            unrelated.slides.add_slide(unrelated.slide_layouts[6])
        unrelated.save(pptx)
        review["structural_evidence"]["pptx"]["artifact_sha256"] = hashlib.sha256(
            pptx.read_bytes()
        ).hexdigest()
        payload = json.loads(render_backup)
        payload["pptx_sha256"] = review["structural_evidence"]["pptx"][
            "artifact_sha256"
        ]
        render_report.write_text(json.dumps(payload), encoding="utf-8")
        render_sha = hashlib.sha256(render_report.read_bytes()).hexdigest()
        review["structural_evidence"]["render"]["artifact_sha256"] = render_sha
        review["structural_evidence"]["target_fonts"]["artifact_sha256"] = render_sha
        with self.assertRaisesRegex(ValueError, "structural-pptx-content-mismatch"):
            evaluate_final_gate(
                self.preferences,
                review,
                deterministic,
                corpus=self.corpus,
                scenarios=self.scenarios,
                fixture_contract=self.final_contract,
                evidence_root=evidence_root,
            )
        pptx.write_bytes(pptx_backup)
        render_report.write_text(render_backup, encoding="utf-8")
        review = review_backup

        blind_path = evidence_root / review["pair_assignments"][0]["blind_artifacts"][
            review["pair_assignments"][0]["masked_left_id"]
        ]["artifact_path"]
        blind_path.write_text("<svg/>", encoding="utf-8")
        with self.assertRaisesRegex(ValueError, "final-blind-evidence-digest-mismatch"):
            evaluate_final_gate(
                self.preferences,
                review,
                deterministic,
                corpus=self.corpus,
                scenarios=self.scenarios,
                fixture_contract=self.final_contract,
                evidence_root=evidence_root,
            )


if __name__ == "__main__":
    unittest.main()
