#!/usr/bin/env python3
"""Validate a final GHB PPTX and emit console, JSON, and Markdown reports."""

from __future__ import annotations

import argparse
import hashlib
import json
import math
import os
import re
import sys
import tempfile
import zipfile
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any, Iterable
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
from scripts.merge_template_master import (  # noqa: E402
    NS,
    parse_rels,
    presentation_slide_parts,
    qn,
    relation_type,
    resolve_target,
    slide_layout_part,
)
from scripts.review_visual_quality import is_passive_review_text  # noqa: E402


CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
MOJIBAKE = ("\ufffd", "Ã", "Â", "â€", "ðŸ", "ï»¿")
PLACEHOLDER_PATTERNS = (
    re.compile(r"\{\{[^}]+\}\}"),
    re.compile(r"\b(?:XXX+|TBD|TODO)\b", re.IGNORECASE),
    re.compile(r"202X"),
)


def _file_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


@dataclass(frozen=True)
class Issue:
    severity: str
    code: str
    message: str
    slide: int | None = None


@dataclass
class SlideSummary:
    slide: int
    role: str
    layout_part: str
    master_part: str
    text_objects: int = 0
    text_chars: int = 0
    shape_objects: int = 0
    group_objects: int = 0
    picture_objects: int = 0
    table_objects: int = 0
    chart_objects: int = 0
    min_font_pt: float | None = None
    min_font_by_role: dict[str, float] = field(default_factory=dict)
    title: str = ""
    full_slide_pictures: int = 0
    out_of_bounds_objects: int = 0
    full_white_rectangles: int = 0
    notes_chars: int = 0
    empty_text_boxes: int = 0
    text_boxes_too_small: int = 0
    possible_text_overlaps: int = 0
    page_number_found: bool = False


@dataclass
class ValidationReport:
    pptx: str
    passed: bool
    file_size: int
    page_count: int
    body_count: int
    has_ending: bool
    slide_size: dict[str, Any]
    issues: list[Issue] = field(default_factory=list)
    slides: list[SlideSummary] = field(default_factory=list)
    package: dict[str, Any] = field(default_factory=dict)
    known_limitations: list[str] = field(default_factory=list)
    manual_review: list[str] = field(default_factory=list)
    quality: dict[str, Any] = field(default_factory=dict)

    @property
    def errors(self) -> list[Issue]:
        return [issue for issue in self.issues if issue.severity == "error"]

    @property
    def warnings(self) -> list[Issue]:
        return [issue for issue in self.issues if issue.severity == "warning"]


def _issue(issues: list[Issue], severity: str, code: str, message: str, slide: int | None = None) -> None:
    issues.append(Issue(severity, code, message, slide))


def _package(path: Path, issues: list[Issue]) -> tuple[dict[str, bytes], list[str]]:
    try:
        with zipfile.ZipFile(path) as archive:
            infos = [info for info in archive.infolist() if not info.is_dir()]
            names = [info.filename for info in infos]
            duplicates = sorted({name for name in names if names.count(name) > 1})
            for name in duplicates:
                _issue(issues, "error", "duplicate-part", f"duplicate ZIP part: {name}")
            bad = archive.testzip()
            if bad:
                _issue(issues, "error", "corrupt-zip-member", f"corrupt ZIP member: {bad}")
            return ({info.filename: archive.read(info.filename) for info in infos}, names)
    except zipfile.BadZipFile as exc:
        _issue(issues, "error", "invalid-zip", f"invalid PPTX ZIP: {exc}")
        return {}, []


def _owner_for_rels(name: str) -> str | None:
    if "/_rels/" in name:
        prefix, filename = name.split("/_rels/", 1)
        return f"{prefix}/{filename[:-5]}"
    if name.startswith("_rels/"):
        return name[len("_rels/") : -5]
    return None


def _check_relationships(parts: dict[str, bytes], issues: list[Issue]) -> dict[str, Any]:
    rel_count = 0
    external_count = 0
    for name, payload in parts.items():
        if not name.endswith(".rels"):
            continue
        owner = _owner_for_rels(name)
        if owner is None:
            continue
        try:
            root = ET.fromstring(payload)
        except ET.ParseError as exc:
            _issue(issues, "error", "invalid-rels", f"{name}: {exc}")
            continue
        ids = [rel.get("Id", "") for rel in root]
        for rid in sorted({rid for rid in ids if ids.count(rid) > 1}):
            _issue(issues, "error", "duplicate-rid", f"{name}: duplicate relationship ID {rid}")
        for rel in root:
            rel_count += 1
            if rel.get("TargetMode") == "External":
                external_count += 1
                continue
            target = resolve_target(owner, rel.get("Target", ""))
            if target not in parts:
                _issue(
                    issues,
                    "error",
                    "dangling-relationship",
                    f"{name}: {rel.get('Id')} targets missing part {target}",
                )
    return {"relationship_count": rel_count, "external_relationship_count": external_count}


def _check_content_types(parts: dict[str, bytes], issues: list[Issue]) -> dict[str, Any]:
    name = "[Content_Types].xml"
    if name not in parts:
        _issue(issues, "error", "missing-content-types", f"missing {name}")
        return {}
    raw = parts[name]
    if b'<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">' not in raw:
        _issue(
            issues,
            "error",
            "noncanonical-content-types-namespace",
            "[Content_Types].xml must use the OPC namespace as the default; "
            "LibreOffice rejects the equivalent ns0-prefixed wire form",
        )
    try:
        root = ET.fromstring(raw)
    except ET.ParseError as exc:
        _issue(issues, "error", "invalid-content-types", str(exc))
        return {}
    defaults = [node.get("Extension", "").lower() for node in root if node.tag.endswith("Default")]
    overrides = [node.get("PartName", "") for node in root if node.tag.endswith("Override")]
    for ext in sorted({ext for ext in defaults if defaults.count(ext) > 1}):
        _issue(issues, "error", "duplicate-content-default", f"duplicate content-type Default for .{ext}")
    for part in sorted({part for part in overrides if overrides.count(part) > 1}):
        _issue(issues, "error", "duplicate-content-override", f"duplicate content-type Override for {part}")
    override_set = {part.lstrip("/") for part in overrides}
    for part in parts:
        if re.fullmatch(r"ppt/(?:slides/slide|slideLayouts/slideLayout|slideMasters/slideMaster|theme/theme)\d+\.xml", part):
            if part not in override_set:
                _issue(issues, "error", "missing-content-override", f"missing content-type Override for {part}")
        if part.startswith("ppt/media/") and "." in part:
            ext = part.rsplit(".", 1)[-1].lower()
            if ext not in defaults:
                _issue(issues, "error", "missing-media-default", f"missing content-type Default for .{ext}")
    return {"default_count": len(defaults), "override_count": len(overrides)}


def _iter_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)


def _iter_shapes_with_group_context(shapes: Iterable[Any], *, nested: bool = False) -> Iterable[tuple[Any, bool]]:
    """Yield shapes and whether their coordinates are local to a parent group."""
    for shape in shapes:
        yield shape, nested
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes_with_group_context(shape.shapes, nested=True)


def _shape_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    try:
        return "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()
    except Exception:
        return ""


def _font_sizes(shape: Any) -> list[float]:
    if not getattr(shape, "has_text_frame", False):
        return []
    values: list[float] = []
    try:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size is not None:
                    values.append(round(run.font.size.pt, 2))
    except Exception:
        return values
    return values


def _shape_typography_role(shape: Any) -> str | None:
    """Map stable SVG/DrawingML shape names back to typography roles."""
    name = str(getattr(shape, "name", "") or "").strip().lower().replace("_", "-")
    if name == "main-title" or name.startswith("main-title-"):
        return "title"
    for role in ("body", "caption", "source", "footer"):
        if name == role or name.startswith(f"{role}-"):
            return role
    return None


def _effective_text_box(
    shape: Any,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
) -> tuple[int, int, int, int, str]:
    """Approximate the visible ink band for converter-made single lines."""
    try:
        frame = shape.text_frame
        is_single_line = len(frame.paragraphs) == 1 and "\n" not in text
        is_converter_frame = (
            frame.word_wrap is False
            and frame.auto_size == MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        )
    except Exception:
        return left, top, width, height, text
    sizes = _font_sizes(shape)
    if not is_single_line or not is_converter_frame or not sizes:
        return left, top, width, height, text

    # The SVG converter deliberately gives a single-line frame roughly 1.6em
    # of vertical room so glyphs do not clip across Office renderers. Treating
    # that transparent safety area as visible text creates false overlap
    # warnings for stacked value/unit and number/label typography.
    ink_height = int(round(max(sizes) * 12700))
    if ink_height <= 0 or ink_height >= height:
        return left, top, width, height, text
    ink_top = top + min(int(round(ink_height * 0.05)), height - ink_height)
    return left, ink_top, width, ink_height, text


def _is_white_fill(shape: Any) -> bool:
    try:
        color = shape.fill.fore_color.rgb
        return color is not None and str(color).upper() == "FFFFFF"
    except Exception:
        return False


def _notes_text(slide: Any) -> str:
    try:
        notes_slide = slide.notes_slide
    except Exception:
        return ""
    return "\n".join(filter(None, (_shape_text(shape) for shape in _iter_shapes(notes_slide.shapes))))


def _summarize_slide(
    slide: Any,
    *,
    index: int,
    role: str,
    layout_part: str,
    master_part: str,
    slide_width: int,
    slide_height: int,
    issues: list[Issue],
    typography: dict[str, Any] | None = None,
) -> SlideSummary:
    summary = SlideSummary(index, role, layout_part, master_part)
    positioned_text: list[tuple[int, int, str]] = []
    named_title_text: list[tuple[int, int, str]] = []
    text_boxes: list[tuple[int, int, int, int, str]] = []
    font_sizes: list[float] = []
    font_sizes_by_role: dict[str, list[float]] = {}
    slide_area = max(slide_width * slide_height, 1)
    for shape, nested_in_group in _iter_shapes_with_group_context(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            summary.group_objects += 1
        else:
            summary.shape_objects += 1
        text = _shape_text(shape)
        has_text_frame = bool(getattr(shape, "has_text_frame", False))
        if (
            has_text_frame
            and shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
            and not text
            and not bool(getattr(shape, "is_placeholder", False))
        ):
            summary.empty_text_boxes += 1
        if text:
            summary.text_objects += 1
            summary.text_chars += len(re.sub(r"\s+", "", text))
            positioned_text.append((int(getattr(shape, "top", 0) or 0), int(getattr(shape, "left", 0) or 0), text))
            explicit_font_sizes = _font_sizes(shape)
            font_sizes.extend(explicit_font_sizes)
            typography_role = _shape_typography_role(shape)
            if typography_role == "title":
                named_title_text.append(
                    (int(getattr(shape, "top", 0) or 0), int(getattr(shape, "left", 0) or 0), text)
                )
            if typography_role and explicit_font_sizes:
                font_sizes_by_role.setdefault(typography_role, []).extend(explicit_font_sizes)
            for marker in MOJIBAKE:
                if marker in text:
                    _issue(issues, "error", "mojibake", f"corrupt text marker {marker!r}: {text[:80]}", index)
            for pattern in PLACEHOLDER_PATTERNS:
                match = pattern.search(text)
                if match:
                    _issue(issues, "error", "placeholder", f"unreplaced placeholder {match.group(0)!r}", index)
        if getattr(shape, "has_table", False):
            summary.table_objects += 1
        if getattr(shape, "has_chart", False):
            summary.chart_objects += 1
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            summary.picture_objects += 1
        left = int(getattr(shape, "left", 0) or 0)
        top = int(getattr(shape, "top", 0) or 0)
        width = int(getattr(shape, "width", 0) or 0)
        height = int(getattr(shape, "height", 0) or 0)
        if text:
            text_boxes.append(_effective_text_box(shape, left, top, width, height, text))
            explicit_sizes = _font_sizes(shape)
            # python-pptx exposes grouped child geometry in the group's local
            # coordinate system. Comparing that raw height to points creates
            # false positives for otherwise valid native template groups.
            if not nested_in_group and explicit_sizes and height / 12700 < min(explicit_sizes) * 1.05:
                summary.text_boxes_too_small += 1
        if not nested_in_group and (
            left < -1 or top < -1 or left + width > slide_width + 1 or top + height > slide_height + 1
        ):
            summary.out_of_bounds_objects += 1
        area_ratio = max(width, 0) * max(height, 0) / slide_area
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and area_ratio >= 0.85:
            summary.full_slide_pictures += 1
        if area_ratio >= 0.90 and _is_white_fill(shape):
            summary.full_white_rectangles += 1
    if named_title_text:
        summary.title = " ".join(
            re.sub(r"\s+", " ", text).strip()
            for _top, _left, text in sorted(named_title_text)
        )
    elif positioned_text:
        summary.title = sorted(positioned_text)[0][2].splitlines()[0].strip()
    summary.page_number_found = any(
        re.search(r"\b\d{1,2}\s*/\s*\d{1,2}\b", text)
        for _left, _top, _width, _height, text in text_boxes
    )
    for left_index, left_box in enumerate(text_boxes):
        lx, ly, lw, lh, _lt = left_box
        left_area = max(lw, 0) * max(lh, 0)
        if not left_area:
            continue
        for right_box in text_boxes[left_index + 1 :]:
            rx, ry, rw, rh, _rt = right_box
            right_area = max(rw, 0) * max(rh, 0)
            if not right_area:
                continue
            intersection_w = max(0, min(lx + lw, rx + rw) - max(lx, rx))
            intersection_h = max(0, min(ly + lh, ry + rh) - max(ly, ry))
            if intersection_w * intersection_h / min(left_area, right_area) >= 0.60:
                summary.possible_text_overlaps += 1
    if font_sizes:
        summary.min_font_pt = min(font_sizes)
    summary.min_font_by_role = {
        role: min(values)
        for role, values in sorted(font_sizes_by_role.items())
        if values
    }
    summary.notes_chars = len(re.sub(r"\s+", "", _notes_text(slide)))

    if summary.shape_objects + summary.group_objects == 0:
        _issue(issues, "error", "blank-slide", "slide has no slide-level objects", index)
    if summary.out_of_bounds_objects and role == "body":
        _issue(issues, "error", "object-out-of-bounds", f"{summary.out_of_bounds_objects} object(s) exceed slide bounds", index)
    elif summary.out_of_bounds_objects:
        _issue(
            issues,
            "warning",
            "template-bleed",
            f"{summary.out_of_bounds_objects} template object(s) extend beyond slide bounds; verify intentional bleed visually",
            index,
        )
    if role == "body" and summary.full_slide_pictures and summary.text_objects == 0:
        _issue(issues, "error", "full-slide-image-body", "body slide is effectively a full-slide image with no editable text", index)
    if role == "body" and summary.full_white_rectangles:
        _issue(issues, "error", "full-white-rectangle", "body slide contains a near-full-slide white rectangle", index)
    if summary.min_font_pt is not None and summary.min_font_pt < 9:
        _issue(issues, "warning", "small-font", f"minimum explicit font size is {summary.min_font_pt:g} pt", index)
    if role == "body" and isinstance(typography, dict) and typography.get("enforcement") == "strict":
        floor_fields = {
            "title": "min_title_pt",
            "body": "min_body_pt",
            "caption": "min_caption_pt",
            "source": "min_source_pt",
            "footer": "min_footer_pt",
        }
        for typography_role, observed in summary.min_font_by_role.items():
            floor = typography.get(floor_fields[typography_role])
            if isinstance(floor, (int, float)) and not isinstance(floor, bool) and observed < float(floor):
                _issue(
                    issues,
                    "error",
                    f"typography-{typography_role}-below-min",
                    f"final PPTX {typography_role} text is {observed:g} pt; requires {float(floor):g} pt",
                    index,
                )
    if role == "body" and summary.empty_text_boxes:
        _issue(issues, "warning", "empty-text-box", f"{summary.empty_text_boxes} empty text box(es)", index)
    if role == "body" and summary.text_boxes_too_small:
        _issue(issues, "warning", "text-box-too-small", f"{summary.text_boxes_too_small} text box(es) are shorter than their explicit font size", index)
    if role == "body" and summary.possible_text_overlaps:
        _issue(issues, "warning", "possible-text-overlap", f"{summary.possible_text_overlaps} pair(s) of text boxes overlap by at least 60%", index)
    if role == "body" and not summary.page_number_found:
        _issue(issues, "warning", "missing-page-number", "no NN / total page-number text was detected", index)
    if role == "body" and summary.text_chars > 600:
        _issue(issues, "warning", "high-text-density", f"body slide has {summary.text_chars} non-whitespace text characters", index)
    if role == "body" and summary.shape_objects + summary.group_objects > 120:
        _issue(issues, "warning", "high-object-density", f"body slide has {summary.shape_objects + summary.group_objects} objects/groups", index)
    if summary.title and len(summary.title) > 80:
        _issue(issues, "warning", "long-title", f"title has {len(summary.title)} characters", index)
    return summary


def _layout_master(parts: dict[str, bytes], layout: str) -> str:
    rels = parse_rels(parts, layout)
    masters = [rel for rel in rels if relation_type(rel) == "slideMaster"]
    if len(masters) != 1:
        raise ValueError(f"{layout}: expected one slideMaster relationship, found {len(masters)}")
    return resolve_target(layout, masters[0].get("Target", ""))


def _check_mount_chain(
    parts: dict[str, bytes],
    ordered_slides: list[str],
    roles: list[str],
    issues: list[Issue],
) -> tuple[list[tuple[str, str]], dict[str, Any]]:
    chains: list[tuple[str, str]] = []
    used_layouts_by_master: dict[str, set[str]] = {}
    for index, (slide, role) in enumerate(zip(ordered_slides, roles), 1):
        try:
            layout = slide_layout_part(parts, slide)
            master = _layout_master(parts, layout)
        except (KeyError, ValueError, RuntimeError, ET.ParseError) as exc:
            _issue(issues, "error", "mount-chain", str(exc), index)
            chains.append(("", ""))
            continue
        chains.append((layout, master))
        used_layouts_by_master.setdefault(master, set()).add(layout)
    for master, used_layouts in used_layouts_by_master.items():
        try:
            rels = parse_rels(parts, master)
            registered_by_rel = {
                resolve_target(master, rel.get("Target", "")): rel.get("Id")
                for rel in rels if relation_type(rel) == "slideLayout"
            }
            theme_targets = [
                resolve_target(master, rel.get("Target", ""))
                for rel in rels if relation_type(rel) == "theme"
            ]
            master_xml = ET.fromstring(parts[master])
            registered_ids = {
                node.get(qn("r", "id")) for node in master_xml.findall(".//p:sldLayoutId", NS)
            }
        except (KeyError, ET.ParseError, RuntimeError) as exc:
            _issue(issues, "error", "master-structure", f"{master}: {exc}")
            continue
        missing = sorted(used_layouts - set(registered_by_rel))
        if missing:
            _issue(issues, "error", "unregistered-used-layout", f"{master} does not register used layout(s): {missing}")
        if registered_ids != set(registered_by_rel.values()):
            _issue(issues, "error", "master-layout-list-mismatch", f"{master} layout ID list does not match its relationships")
        if len(theme_targets) != 1 or theme_targets[0] not in parts:
            _issue(issues, "error", "master-theme", f"{master} must reference one existing theme")
    masters = {master for _layout, master in chains if master}
    if len(masters) != 1:
        _issue(issues, "error", "multiple-role-masters", f"cover/body/ending use {len(masters)} masters: {sorted(masters)}")
    return chains, {"masters_used": sorted(masters), "layouts_used": sorted({layout for layout, _master in chains if layout})}


def _load_layout_plan(path: Path | None, issues: list[Issue]) -> list[dict[str, Any]]:
    if path is None:
        return []
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        _issue(issues, "error", "invalid-layout-plan", str(exc))
        return []
    if not isinstance(payload, list):
        _issue(issues, "error", "invalid-layout-plan", "layout plan top level must be a list")
        return []
    return [entry for entry in payload if isinstance(entry, dict)]


def _load_warning_waivers(
    path: Path | None, issues: list[Issue]
) -> list[dict[str, Any]]:
    if path is None:
        return []
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        _issue(issues, "error", "invalid-warning-waivers", str(exc))
        return []
    rows = payload.get("waivers") if isinstance(payload, dict) else None
    if (
        not isinstance(payload, dict)
        or payload.get("schema") != "ghb.warning-waivers.v1"
        or not isinstance(rows, list)
    ):
        _issue(
            issues,
            "error",
            "invalid-warning-waivers",
            "expected ghb.warning-waivers.v1 with a waivers list",
        )
        return []
    accepted: list[dict[str, Any]] = []
    seen: set[tuple[str, str]] = set()
    for index, row in enumerate(rows, 1):
        if not isinstance(row, dict) or set(row) != {"code", "slide", "reason", "approved_by"}:
            _issue(
                issues,
                "error",
                "invalid-warning-waivers",
                f"waiver {index} must contain code, slide, reason, and approved_by",
            )
            continue
        code = row.get("code")
        slide = row.get("slide")
        reason = row.get("reason")
        approved_by = row.get("approved_by")
        valid_slide = slide is None or (
            isinstance(slide, int) and not isinstance(slide, bool) and slide > 0
        ) or (
            isinstance(slide, str) and bool(re.fullmatch(r"[A-Za-z0-9][A-Za-z0-9_-]{0,127}", slide))
        )
        if (
            not isinstance(code, str)
            or not re.fullmatch(r"[a-z][a-z0-9-]{2,80}", code)
            or not valid_slide
            or not isinstance(reason, str)
            or not 1 <= len(reason.strip()) <= 512
            or not isinstance(approved_by, str)
            or not 1 <= len(approved_by.strip()) <= 128
            or not is_passive_review_text(reason, maximum=512)
            or not is_passive_review_text(approved_by, maximum=128)
        ):
            _issue(
                issues,
                "error",
                "invalid-warning-waivers",
                f"waiver {index} contains invalid or unsafe values",
            )
            continue
        key = (code, str(slide))
        if key in seen:
            _issue(
                issues,
                "error",
                "invalid-warning-waivers",
                f"duplicate waiver for {code!r} slide {slide!r}",
            )
            continue
        seen.add(key)
        accepted.append({
            "code": code,
            "slide": slide,
            "reason": reason.strip(),
            "approved_by": approved_by.strip(),
        })
    return accepted


def _renderer_matches(target: str, actual: str | None) -> bool:
    if target == "auto":
        return bool(actual)
    normalized = (actual or "").strip().lower()
    if target == "libreoffice":
        return normalized in {"libreoffice", "soffice"}
    return normalized == target


def _waiver_matches(waiver: dict[str, Any], code: str, slide: int | str | None) -> bool:
    return waiver["code"] == code and (
        waiver["slide"] is None or str(waiver["slide"]) == str(slide)
    )


def validate_pptx(
    path: Path,
    *,
    expected_body_count: int | None = None,
    expect_ending: bool | None = None,
    source_svg_dir: Path | None = None,
    layout_plan_path: Path | None = None,
    cover_text: list[str] | None = None,
    render_dir: Path | None = None,
    expect_body_notes: bool = False,
    svg_report_paths: list[Path] | None = None,
    readback_markdown_path: Path | None = None,
    freshness: dict[str, Any] | None = None,
    review_report_path: Path | None = None,
    review_required: bool = False,
    quality_policy: str = "draft",
    warning_waivers_path: Path | None = None,
    target_renderer: str = "auto",
    visual_profile_path: Path | None = None,
    font_embed_report_path: Path | None = None,
) -> ValidationReport:
    issues: list[Issue] = []
    if quality_policy not in {"draft", "release"}:
        _issue(issues, "error", "invalid-quality-policy", f"unknown policy {quality_policy!r}")
    if target_renderer not in {"auto", "libreoffice", "powerpoint", "wps"}:
        _issue(issues, "error", "invalid-target-renderer", f"unknown renderer {target_renderer!r}")
    visual_profile: dict[str, Any] = {}
    if visual_profile_path is not None:
        try:
            loaded_profile = json.loads(visual_profile_path.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError) as exc:
            _issue(issues, "error", "invalid-visual-profile", str(exc))
        else:
            if isinstance(loaded_profile, dict):
                visual_profile = loaded_profile
            else:
                _issue(issues, "error", "invalid-visual-profile", "visual profile must be an object")
    typography = visual_profile.get("typography")
    if not isinstance(typography, dict):
        typography = {}
    if not path.is_file():
        return ValidationReport(
            str(path), False, 0, 0, 0, bool(expect_ending), {},
            issues=[Issue("error", "missing-file", f"PPTX not found: {path}")],
        )
    size = path.stat().st_size
    if size < 10_000:
        _issue(issues, "error", "file-too-small", f"PPTX size is suspiciously small: {size} bytes")
    parts, names = _package(path, issues)
    required = ("[Content_Types].xml", "ppt/presentation.xml", "ppt/_rels/presentation.xml.rels")
    for required_part in required:
        if required_part not in parts:
            _issue(issues, "error", "missing-required-part", f"missing required part: {required_part}")
    package_info = {"part_count": len(names), **_check_relationships(parts, issues), **_check_content_types(parts, issues)}
    presentation_xml = parts.get("ppt/presentation.xml", b"").decode(
        "utf-8", errors="ignore"
    )
    font_parts = sorted(
        name
        for name in names
        if name.startswith("ppt/fonts/") and name.endswith(".fntdata")
    )
    embedded_font_count = presentation_xml.count("<p:embeddedFont>")
    save_subset_fonts = (
        True if 'saveSubsetFonts="1"' in presentation_xml
        else False if 'saveSubsetFonts="0"' in presentation_xml
        else None
    )
    embedding_enabled = (
        'embedTrueTypeFonts="1"' in presentation_xml
        and save_subset_fonts is not None
    )
    font_embedding: dict[str, Any] = {
        "fonts_embedded": len(font_parts),
        "embedded_font_names": [],
        "fsType_ok": None,
        "embedding_enabled": embedding_enabled,
        "save_subset_fonts": save_subset_fonts,
        "font_parts": font_parts,
    }
    if font_parts and (not embedding_enabled or embedded_font_count != len(font_parts)):
        _issue(
            issues,
            "error",
            "invalid-embedded-font-contract",
            "embedded font parts, presentation flags, and embeddedFontLst entries disagree",
        )
    if font_embed_report_path is not None:
        try:
            font_report = json.loads(font_embed_report_path.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError) as exc:
            _issue(issues, "error", "invalid-font-embed-report", str(exc))
        else:
            valid_font_report = (
                isinstance(font_report, dict)
                and font_report.get("schema") == "ghb.font-embed-report.v1"
                and isinstance(font_report.get("fonts_embedded"), int)
                and isinstance(font_report.get("embedded_font_names"), list)
                and isinstance(font_report.get("fsType_ok"), bool)
            )
            if not valid_font_report:
                _issue(
                    issues,
                    "error",
                    "invalid-font-embed-report",
                    "font embed report does not match ghb.font-embed-report.v1",
                )
            else:
                font_embedding.update({
                    "embedded_font_names": font_report["embedded_font_names"],
                    "fsType_ok": font_report["fsType_ok"],
                })
                if font_report["fonts_embedded"] != len(font_parts):
                    _issue(
                        issues,
                        "error",
                        "font-embed-report-mismatch",
                        "font embed report count does not match PPTX font parts",
                    )
    package_info["font_embedding"] = font_embedding

    ordered_slides: list[str] = []
    if all(name in parts for name in required[1:]):
        try:
            ordered_slides = presentation_slide_parts(parts)
            pres = ET.fromstring(parts["ppt/presentation.xml"])
            slide_ids = [node.get("id") for node in pres.findall(".//p:sldId", NS)]
            master_ids = [node.get("id") for node in pres.findall(".//p:sldMasterId", NS)]
            if len(slide_ids) != len(set(slide_ids)):
                _issue(issues, "error", "duplicate-slide-id", "presentation has duplicate slide IDs")
            if len(master_ids) != len(set(master_ids)):
                _issue(issues, "error", "duplicate-master-id", "presentation has duplicate master IDs")
        except (ET.ParseError, RuntimeError, KeyError) as exc:
            _issue(issues, "error", "presentation-graph", str(exc))

    svg_count = None
    if source_svg_dir is not None:
        svg_count = len(list(source_svg_dir.glob("*.svg"))) if source_svg_dir.is_dir() else 0
        if svg_count == 0:
            _issue(issues, "error", "missing-source-svg", f"no source SVG files found: {source_svg_dir}")
    if expected_body_count is None:
        expected_body_count = svg_count
    if expect_ending is None:
        expect_ending = True
    inferred_body_count = max(len(ordered_slides) - 1 - int(expect_ending), 0)
    if expected_body_count is not None and inferred_body_count != expected_body_count:
        _issue(
            issues,
            "error",
            "page-count",
            f"expected 1 cover + {expected_body_count} body + {int(expect_ending)} ending, got {len(ordered_slides)} pages",
        )
    if svg_count is not None and inferred_body_count != svg_count:
        _issue(issues, "error", "svg-page-count", f"body count {inferred_body_count} does not match source SVG count {svg_count}")
    roles = ["cover"] + ["body"] * inferred_body_count + (["ending"] if expect_ending else [])
    if len(roles) != len(ordered_slides):
        roles = ["cover"] + ["body"] * max(len(ordered_slides) - 1, 0)

    chains: list[tuple[str, str]] = [("", "")] * len(ordered_slides)
    if ordered_slides:
        chains, mount_info = _check_mount_chain(parts, ordered_slides, roles, issues)
        package_info.update(mount_info)

    slides: list[SlideSummary] = []
    slide_full_texts: list[str] = []
    slide_size: dict[str, Any] = {}
    try:
        presentation = Presentation(path)
        width = int(presentation.slide_width)
        height = int(presentation.slide_height)
        ratio = width / height if height else 0
        slide_size = {"width_emu": width, "height_emu": height, "ratio": round(ratio, 6), "is_16_9": abs(ratio - 16 / 9) < 0.01}
        if not slide_size["is_16_9"]:
            _issue(issues, "error", "slide-size", f"slide ratio is {ratio:.4f}, expected 16:9")
        if len(presentation.slides) != len(ordered_slides):
            _issue(issues, "error", "python-pptx-page-count", f"python-pptx sees {len(presentation.slides)} slides, package graph sees {len(ordered_slides)}")
        for index, slide in enumerate(presentation.slides, 1):
            role = roles[index - 1] if index <= len(roles) else "body"
            layout, master = chains[index - 1] if index <= len(chains) else ("", "")
            slides.append(
                _summarize_slide(
                    slide,
                    index=index,
                    role=role,
                    layout_part=layout,
                    master_part=master,
                    slide_width=width,
                    slide_height=height,
                    issues=issues,
                    typography=typography,
                )
            )
            slide_full_texts.append(
                "\n".join(filter(None, (_shape_text(shape) for shape in _iter_shapes(slide.shapes))))
            )
    except Exception as exc:
        _issue(issues, "error", "python-pptx-open", f"python-pptx could not open deck: {exc}")

    all_text = "\n".join(summary.title for summary in slides) + "\n"
    for slide in getattr(locals().get("presentation"), "slides", []):
        all_text += "\n".join(_shape_text(shape) for shape in _iter_shapes(slide.shapes)) + "\n"
    if not all_text.strip():
        _issue(issues, "error", "empty-text-roundtrip", "no text could be read back from the final deck")
    for required_text in cover_text or []:
        if required_text and required_text not in all_text:
            _issue(issues, "error", "missing-cover-text", f"cover text not found: {required_text!r}", 1)

    plan = _load_layout_plan(layout_plan_path, issues)
    if plan and len(plan) != inferred_body_count:
        _issue(issues, "error", "layout-plan-count", f"layout plan has {len(plan)} entries for {inferred_body_count} body slides")
    for offset, entry in enumerate(plan, 2):
        if offset - 1 >= len(slides):
            break
        message = str(entry.get("key_message") or entry.get("message") or "").strip()
        slide_text = slide_full_texts[offset - 1] if offset - 1 < len(slide_full_texts) else ""
        normalized_slide_text = re.sub(r"\s+", "", slide_text)
        if message and re.sub(r"\s+", "", message) not in normalized_slide_text:
            _issue(issues, "error", "missing-planned-text", f"planned key message was not preserved: {message!r}", offset)
        for item in entry.get("items", []):
            if not isinstance(item, (str, int, float)):
                continue
            item_text = str(item).strip()
            if item_text and re.sub(r"\s+", "", item_text) not in normalized_slide_text:
                _issue(
                    issues,
                    "error",
                    "missing-planned-item",
                    f"planned body item was not preserved: {item_text!r}",
                    offset,
                )

    if expect_body_notes:
        for summary in slides:
            if summary.role == "body" and summary.notes_chars == 0:
                _issue(issues, "error", "missing-speaker-notes", "expected body speaker notes are missing", summary.slide)

    slide_xml = b"\n".join(payload for name, payload in parts.items() if re.fullmatch(r"ppt/slides/slide\d+\.xml", name))
    decoded_slide_xml = slide_xml.decode("utf-8", errors="replace")
    if re.search(r'typeface="(?:楷体|KaiTi)"', decoded_slide_xml, re.IGNORECASE):
        _issue(issues, "error", "cover-font", "KaiTi/楷体 remains in slide text runs")
    if "AB1F29" not in decoded_slide_xml.upper():
        _issue(issues, "warning", "brand-color", "GHB primary color #AB1F29 was not found in slide XML")
    for index, slide_part in enumerate(ordered_slides, 1):
        xml = parts.get(slide_part, b"").decode("utf-8", errors="replace")
        run_properties = re.findall(
            r"<a:rPr\b[^>]*(?:/>|>.*?</a:rPr>)",
            xml,
            re.DOTALL,
        )
        if any(
            "<a:noFill" in block or re.search(r'<a:alpha\s+val="0"', block)
            for block in run_properties
        ):
            _issue(issues, "warning", "possibly-invisible-text", "text run may have no fill or zero alpha", index)

    for left, right in zip(slides, slides[1:]):
        if left.title and left.title == right.title:
            _issue(issues, "warning", "adjacent-duplicate-title", f"adjacent slides repeat title {left.title!r}", right.slide)

    min_font_by_role: dict[str, float] = {}
    for summary in slides:
        if summary.role != "body":
            continue
        for typography_role, value in summary.min_font_by_role.items():
            min_font_by_role[typography_role] = min(
                value,
                min_font_by_role.get(typography_role, value),
            )
    package_info["min_font_by_role"] = min_font_by_role

    known_limitations: list[str] = []
    manual_review = [
        "Review typography, visual hierarchy, intentional overlaps, and aesthetic balance page by page.",
        "Confirm PowerPoint rendering on the target enterprise desktop environment.",
    ]
    review_outcome = "unavailable"
    render_status = "unavailable"
    render_provenance: dict[str, Any] = {}
    visual_findings: list[dict[str, Any]] = []
    if render_dir is None:
        known_limitations.append("No render directory was supplied; final visual appearance was not verified by this validation run.")
    else:
        package_info["rendered_pages"] = 0
        render_report_path = render_dir / "render-report.json"
        if render_report_path.is_file():
            try:
                render_payload = json.loads(render_report_path.read_text(encoding="utf-8"))
            except (OSError, json.JSONDecodeError) as exc:
                _issue(issues, "warning", "invalid-render-report", str(exc))
            else:
                render_status = str(render_payload.get("status") or ("passed" if render_payload.get("passed") else "error"))
                review_outcome = "skipped" if render_status == "passed" else "unavailable"
                if render_status != "passed":
                    # The current report is authoritative. Files left by an older
                    # successful render are not evidence for this failed attempt.
                    package_info["rendered_pages"] = 0
                    stale_png_count = (
                        sum(1 for _ in render_dir.glob("slide-*.png"))
                        if render_dir.is_dir()
                        else 0
                    )
                    if stale_png_count:
                        _issue(
                            issues,
                            "warning",
                            "stale-render-files",
                            f"ignored {stale_png_count} page PNG(s) left by an earlier render",
                        )
                else:
                    reported_pptx = render_payload.get("pptx")
                    try:
                        pptx_matches = bool(reported_pptx) and Path(
                            str(reported_pptx)
                        ).resolve() == path.resolve()
                    except (OSError, RuntimeError, ValueError):
                        pptx_matches = False
                    reported_digest = render_payload.get("pptx_sha256")
                    current_digest = _file_sha256(path)
                    if not pptx_matches or reported_digest != current_digest:
                        review_outcome = "stale"
                        message = (
                            "render report is not bound to the PPTX being validated; "
                            "ignored existing page images"
                        )
                        known_limitations.append(message)
                        _issue(issues, "warning", "render-pptx-mismatch", message)
                    else:
                        declared_pages = []
                        for value in render_payload.get("outputs", []):
                            if not isinstance(value, str):
                                continue
                            candidate = Path(value)
                            if not candidate.is_absolute():
                                candidate = render_dir / candidate
                            if (
                                candidate.name.startswith("slide-")
                                and candidate.suffix.lower() == ".png"
                                and candidate.is_file()
                            ):
                                declared_pages.append(candidate)
                        package_info["rendered_pages"] = len(declared_pages)
                        if len(declared_pages) != len(slides):
                            review_outcome = "limited"
                            _issue(
                                issues,
                                "warning",
                                "render-count",
                                f"render report declares {len(declared_pages)} page PNGs "
                                f"for {len(slides)} slides",
                            )
                render_provenance = {
                    "renderer": render_payload.get("renderer"),
                    "dpi": render_payload.get("dpi"),
                    "font": render_payload.get("font"),
                    "outputs": render_payload.get("outputs", []),
                    "status": render_status,
                }
                package_info["render"] = {
                    "renderer": render_payload.get("renderer"),
                    "page_count": render_payload.get("page_count"),
                    "passed": render_payload.get("passed"),
                    "status": render_status,
                }
                for warning in render_payload.get("warnings", []):
                    known_limitations.append(str(warning))
                    _issue(issues, "warning", "render-warning", str(warning))
                for error in render_payload.get("errors", []):
                    message = str(error)
                    known_limitations.append(message)
                    _issue(issues, "warning", f"render-{render_status}", message)
        else:
            message = "render directory has no authoritative render-report.json"
            known_limitations.append(message)
            _issue(issues, "warning", "missing-render-report", message)

    review_requirement_satisfied = not review_required
    review_freshness = "unavailable"
    if review_report_path is not None:
        try:
            review_payload = json.loads(review_report_path.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError):
            review_payload = None
        allowed_report = {
            "schema", "outcome", "deterministic_status", "completion_status",
            "freshness", "findings", "dimension_reviewability", "limitations",
            "provenance", "error", "request_digest", "reviewer_metadata",
        }
        active = re.compile(
            r"(?:[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]|<[^>]+>|\[[^\]]+\]\([^)]+\)|(?:https?|file)://)",
            re.IGNORECASE,
        )
        success_outcomes = {"passed", "needs-revision", "limited"}
        review_outcome_value = (
            review_payload.get("outcome") if isinstance(review_payload, dict) else None
        )
        valid = (
            isinstance(review_payload, dict)
            and set(review_payload).issubset(allowed_report)
            and {
                "schema", "outcome", "freshness", "findings",
                "dimension_reviewability", "limitations", "provenance",
            }.issubset(review_payload)
            and review_payload.get("schema") == "ghb.visual-review-report.v1"
            and review_payload.get("outcome")
            in {"passed", "needs-revision", "limited", "unavailable", "skipped", "error"}
            and isinstance(review_payload.get("findings"), list)
            and len(review_payload.get("findings", [])) <= 100
            and isinstance(review_payload.get("dimension_reviewability"), list)
            and len(review_payload.get("dimension_reviewability", [])) <= 6
            and isinstance(review_payload.get("limitations"), list)
            and len(review_payload.get("limitations", [])) <= 32
            and isinstance(review_payload.get("provenance"), dict)
        )
        if valid and review_outcome_value in success_outcomes:
            valid = (
                set(review_payload)
                == {
                    "schema", "outcome", "deterministic_status", "completion_status",
                    "freshness", "findings", "dimension_reviewability", "limitations",
                    "provenance", "request_digest", "reviewer_metadata",
                }
                and review_payload.get("freshness") == "fresh"
                and review_payload.get("deterministic_status") in {"passed", "failed"}
                and review_payload.get("deterministic_status")
                == (
                    "failed"
                    if any(issue.severity == "error" for issue in issues)
                    else "passed"
                )
                and review_payload.get("completion_status")
                == (
                    "failed"
                    if review_payload.get("deterministic_status") == "failed"
                    else "completed"
                )
                and isinstance(review_payload.get("request_digest"), str)
                and bool(re.fullmatch(r"[0-9a-f]{64}", review_payload["request_digest"]))
                and isinstance(review_payload.get("reviewer_metadata"), dict)
                and set(review_payload["reviewer_metadata"]) == {"adapter_version"}
                and isinstance(
                    review_payload["reviewer_metadata"].get("adapter_version"), str
                )
                and is_passive_review_text(
                    review_payload["reviewer_metadata"]["adapter_version"], maximum=256
                )
                and render_status == "passed"
                and package_info.get("rendered_pages") == len(slides)
                and isinstance(freshness, dict)
                and freshness.get("states", {}).get("adapter-review") == "fresh"
            )
        safe_findings: list[dict[str, Any]] = []
        safe_provenance: dict[str, Any] = {}
        if valid:
            serialized_metadata = json.dumps(
                {
                    "provenance": review_payload.get("provenance"),
                    "limitations": review_payload.get("limitations"),
                    "dimensions": review_payload.get("dimension_reviewability", []),
                },
                ensure_ascii=False,
            )
            valid = len(serialized_metadata) <= 65536 and not active.search(serialized_metadata)
        if valid:
            provenance = review_payload["provenance"]
            if provenance == {}:
                safe_provenance = {}
            elif provenance == {"adapter": "absent"}:
                safe_provenance = {"adapter": "absent"}
            else:
                allowed_provenance = {
                    "adapter_sha256", "launcher_sha256", "capability", "model_id",
                    "tool_contract", "credentials", "disclosure",
                    "direct_subprocess_isolation",
                }
                digest_pattern = re.compile(r"[0-9a-f]{64}")
                credentials = provenance.get("credentials")
                disclosure = provenance.get("disclosure")
                valid = (
                    set(provenance) == allowed_provenance
                    and isinstance(provenance.get("adapter_sha256"), str)
                    and bool(digest_pattern.fullmatch(provenance["adapter_sha256"]))
                    and (
                        provenance.get("launcher_sha256") is None
                        or (
                            isinstance(provenance.get("launcher_sha256"), str)
                            and bool(digest_pattern.fullmatch(provenance["launcher_sha256"]))
                        )
                    )
                    and provenance.get("capability") in {"local", "remote"}
                    and isinstance(provenance.get("model_id"), str)
                    and is_passive_review_text(provenance["model_id"], maximum=256)
                    and isinstance(provenance.get("tool_contract"), str)
                    and is_passive_review_text(provenance["tool_contract"], maximum=256)
                    and provenance.get("direct_subprocess_isolation") == "trusted-same-user"
                    and isinstance(credentials, list)
                    and len(credentials) <= 32
                    and all(
                        isinstance(item, dict)
                        and set(item) == {"name", "present"}
                        and isinstance(item.get("name"), str)
                        and bool(re.fullmatch(r"[A-Z][A-Z0-9_]{0,127}", item["name"]))
                        and isinstance(item.get("present"), bool)
                        for item in credentials
                    )
                )
                if valid and disclosure is not None:
                    valid = (
                        isinstance(disclosure, dict)
                        and set(disclosure)
                        == {
                            "provider", "destination", "retention", "slide_ids",
                            "authorization_digest",
                        }
                        and all(
                            isinstance(disclosure.get(key), str)
                            and is_passive_review_text(disclosure[key], maximum=512)
                            for key in ("provider", "destination", "retention")
                        )
                        and isinstance(disclosure.get("slide_ids"), list)
                        and all(
                            isinstance(item, str) for item in disclosure.get("slide_ids", [])
                        )
                        and isinstance(disclosure.get("authorization_digest"), str)
                        and bool(
                            digest_pattern.fullmatch(disclosure["authorization_digest"])
                        )
                    )
                if valid:
                    safe_provenance = provenance
            if valid and review_outcome_value in success_outcomes:
                valid = safe_provenance not in ({}, {"adapter": "absent"})
        if valid:
            valid = all(
                isinstance(item, str)
                and is_passive_review_text(item)
                for item in review_payload["limitations"]
            )
        if valid:
            allowed_dimensions = {
                "hierarchy", "spacing", "typography", "cjk", "geometry", "composition"
            }
            for dimension in review_payload["dimension_reviewability"]:
                if (
                    not isinstance(dimension, dict)
                    or set(dimension) != {"dimension", "status", "limitations"}
                    or dimension.get("dimension") not in allowed_dimensions
                    or dimension.get("status") not in {"reviewed", "limited", "unavailable"}
                    or not isinstance(dimension.get("limitations"), list)
                    or any(
                        not isinstance(item, str) or not is_passive_review_text(item)
                        for item in dimension.get("limitations", [])
                    )
                ):
                    valid = False
                    break
            if valid and review_outcome_value in success_outcomes:
                valid = {
                    item["dimension"] for item in review_payload["dimension_reviewability"]
                } == allowed_dimensions
        if valid:
            allowed_finding = {
                "code", "slide_id", "dimension", "reviewability", "severity",
                "location", "evidence", "action",
            }
            approved_review_slides = {
                f"slide-{index:02d}" for index in range(1, len(slides) + 1)
            }
            for finding in review_payload["findings"]:
                if (
                    not isinstance(finding, dict)
                    or set(finding) != allowed_finding
                    or finding.get("severity") != "advisory"
                    or not isinstance(finding.get("code"), str)
                    or not re.fullmatch(r"[a-z][a-z0-9-]{2,80}", finding["code"])
                    or not isinstance(finding.get("slide_id"), str)
                    or not re.fullmatch(r"[A-Za-z0-9][A-Za-z0-9_-]{0,127}", finding["slide_id"])
                    or finding["slide_id"] not in approved_review_slides
                    or finding.get("dimension") not in allowed_dimensions
                    or finding.get("reviewability") not in {"reviewed", "limited", "unavailable"}
                    or not isinstance(finding.get("evidence"), str)
                    or not is_passive_review_text(finding["evidence"])
                    or not isinstance(finding.get("action"), str)
                    or not is_passive_review_text(finding["action"])
                ):
                    valid = False
                    break
                location = finding.get("location")
                if (
                    not isinstance(location, dict)
                    or set(location) != {"x", "y", "width", "height"}
                    or any(
                        isinstance(location[key], bool)
                        or not isinstance(location[key], (int, float))
                        or not math.isfinite(float(location[key]))
                        or not 0 <= float(location[key]) <= 1
                        for key in location
                    )
                    or float(location["x"]) + float(location["width"]) > 1.000001
                    or float(location["y"]) + float(location["height"]) > 1.000001
                ):
                    valid = False
                    break
                safe_findings.append({
                    "code": finding["code"],
                    "severity": "warning",
                    "slide_id": finding["slide_id"],
                    "evidence": {
                        "dimension": finding.get("dimension"),
                        "reviewability": finding.get("reviewability"),
                        "location": finding.get("location"),
                        "observation": finding["evidence"],
                    },
                    "expected": {"source": "optional-visual-review"},
                    "suggested_action": finding["action"],
                    "source": "optional-visual-review",
                })
        if not valid:
            review_outcome = "error"
            review_freshness = "stale"
            known_limitations.append("optional review report failed projection validation")
            _issue(
                issues,
                "warning",
                "invalid-visual-review-report",
                "optional review report failed projection validation",
            )
        else:
            review_outcome = str(review_payload["outcome"])
            review_freshness = str(review_payload.get("freshness", "stale"))
            review_requirement_satisfied = (
                review_outcome == "passed"
                and review_freshness == "fresh"
            )
            visual_findings.extend(safe_findings)
            for limitation in review_payload.get("limitations", []):
                if isinstance(limitation, str):
                    known_limitations.append(limitation)
            render_provenance = {
                **render_provenance,
                "optional_review": safe_provenance,
            }
            package_info["optional_review"] = {
                "outcome": review_outcome,
                "freshness": review_freshness,
                "finding_count": len(safe_findings),
            }

    svg_quality: list[dict[str, Any]] = []
    for svg_report_path in svg_report_paths or []:
        try:
            svg_payload = json.loads(svg_report_path.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError) as exc:
            _issue(issues, "error", "invalid-svg-report", f"{svg_report_path}: {exc}")
            continue
        summary = {
            "stage": svg_payload.get("stage"),
            "passed": svg_payload.get("passed"),
            "file_count": len(svg_payload.get("files", [])),
            "error_count": svg_payload.get("error_count", 0),
            "warning_count": svg_payload.get("warning_count", 0),
            "path": str(svg_report_path),
        }
        svg_quality.append(summary)
        for file_payload in svg_payload.get("files", []):
            if not isinstance(file_payload, dict):
                continue
            for finding in file_payload.get("visual_findings", []):
                if not isinstance(finding, dict):
                    continue
                code = finding.get("code")
                severity = finding.get("severity")
                if not isinstance(code, str) or severity not in {"error", "warning"}:
                    continue
                visual_findings.append({
                    **finding,
                    "source": "svg-visual-quality",
                    "stage": svg_payload.get("stage"),
                    "file": file_payload.get("file"),
                    "slide_id": finding.get("slide_id") or file_payload.get("slide_id"),
                })
        if not svg_payload.get("passed"):
            _issue(issues, "error", "svg-quality-gate", f"SVG {summary['stage']} report failed with {summary['error_count']} errors")
    if svg_quality:
        package_info["svg_quality"] = svg_quality

    if readback_markdown_path is not None:
        if not readback_markdown_path.is_file():
            _issue(
                issues,
                "error",
                "missing-ppt-readback",
                f"ppt_to_md readback was not produced: {readback_markdown_path}",
            )
        else:
            try:
                readback = readback_markdown_path.read_text(encoding="utf-8")
            except OSError as exc:
                _issue(issues, "error", "invalid-ppt-readback", str(exc))
            else:
                readback_chars = len(re.sub(r"\s+", "", readback))
                readback_slides = len(re.findall(r"^## Slide \d+\s*$", readback, re.MULTILINE))
                package_info["ppt_readback"] = {
                    "path": str(readback_markdown_path),
                    "chars": readback_chars,
                    "slide_sections": readback_slides,
                }
                if readback_chars < 20:
                    _issue(issues, "error", "empty-ppt-readback", "ppt_to_md readback is empty")
                if readback_slides != len(slides):
                    _issue(
                        issues,
                        "error",
                        "ppt-readback-page-count",
                        f"ppt_to_md readback contains {readback_slides} slide sections for {len(slides)} pages",
                    )

    freshness_payload = freshness or {"status": "fresh", "states": {}, "issues": []}
    if freshness_payload.get("status") != "fresh" or freshness_payload.get("issues"):
        codes = [
            str(item.get("code", "unknown"))
            for item in freshness_payload.get("issues", [])
            if isinstance(item, dict)
        ]
        _issue(
            issues,
            "error",
            "stale-evidence",
            "freshness evidence is not valid: " + ", ".join(codes or ["stale"]),
        )

    waivers = _load_warning_waivers(warning_waivers_path, issues)
    render_evidence_bound = (
        render_status == "passed"
        and package_info.get("rendered_pages") == len(slides)
    )
    actual_renderer = (
        str(render_provenance.get("renderer") or "") or None
        if render_evidence_bound
        else None
    )
    if quality_policy == "release" and not _renderer_matches(target_renderer, actual_renderer):
        code = "target-renderer-unverified" if actual_renderer is None else "target-renderer-mismatch"
        _issue(
            issues,
            "error",
            code,
            f"release target renderer {target_renderer!r} does not match "
            f"validated renderer {actual_renderer or 'unavailable'!r}",
        )

    warning_rows: list[tuple[str, int | str | None]] = [
        (issue.code, issue.slide) for issue in issues if issue.severity == "warning"
    ] + [
        (str(item.get("code")), item.get("slide_id"))
        for item in visual_findings
        if item.get("severity") == "warning" and item.get("code")
    ]
    resolved_warning_rows = [
        (code, slide)
        for code, slide in warning_rows
        if any(_waiver_matches(waiver, code, slide) for waiver in waivers)
    ]
    unresolved_warning_rows = [
        (code, slide)
        for code, slide in warning_rows
        if not any(_waiver_matches(waiver, code, slide) for waiver in waivers)
    ]
    if quality_policy == "release" and unresolved_warning_rows:
        _issue(
            issues,
            "error",
            "release-unresolved-warnings",
            f"release policy has {len(unresolved_warning_rows)} unresolved warning(s); "
            "fix them or provide explicit warning waivers",
        )

    release_policy = {
        "policy": quality_policy,
        "target_renderer": target_renderer,
        "actual_renderer": actual_renderer,
        "render_evidence_bound": render_evidence_bound,
        "target_renderer_satisfied": _renderer_matches(target_renderer, actual_renderer),
        "waiver_file": str(warning_waivers_path) if warning_waivers_path else None,
        "waiver_count": len(waivers),
        "resolved_warning_count": len(resolved_warning_rows),
        "unresolved_warning_count": len(unresolved_warning_rows),
        "unresolved_warnings": [
            {"code": code, "slide": slide} for code, slide in unresolved_warning_rows
        ],
    }

    report = ValidationReport(
        pptx=str(path.resolve()),
        passed=not any(issue.severity == "error" for issue in issues),
        file_size=size,
        page_count=len(ordered_slides),
        body_count=inferred_body_count,
        has_ending=bool(expect_ending),
        slide_size=slide_size,
        issues=issues,
        slides=slides,
        package=package_info,
        known_limitations=known_limitations,
        manual_review=manual_review,
    )
    blocking = [asdict(issue) for issue in report.errors] + [
        item for item in visual_findings if item["severity"] == "error"
    ]
    advisory = [asdict(issue) for issue in report.warnings] + [
        item for item in visual_findings if item["severity"] == "warning"
    ]
    per_slide = []
    for slide in slides:
        slide_issues = [issue for issue in issues if issue.slide == slide.slide]
        per_slide.append(
            {
                "slide": slide.slide,
                "role": slide.role,
                "blocking_count": sum(issue.severity == "error" for issue in slide_issues),
                "advisory_count": sum(issue.severity == "warning" for issue in slide_issues),
                "issue_codes": [issue.code for issue in slide_issues],
                "evidence": {
                    "text_objects": slide.text_objects,
                    "shape_objects": slide.shape_objects,
                    "min_font_pt": slide.min_font_pt,
                    "min_font_by_role": slide.min_font_by_role,
                    "out_of_bounds_objects": slide.out_of_bounds_objects,
                },
                "actions": [],
            }
        )
    visual_by_slide: dict[str, list[dict[str, Any]]] = {}
    for finding in visual_findings:
        slide_id = str(finding.get("slide_id") or "unknown")
        visual_by_slide.setdefault(slide_id, []).append(finding)
    for slide_id, findings in sorted(visual_by_slide.items()):
        per_slide.append({
            "slide": None,
            "slide_id": slide_id,
            "role": "body-visual",
            "blocking_count": sum(item["severity"] == "error" for item in findings),
            "advisory_count": sum(item["severity"] == "warning" for item in findings),
            "issue_codes": [item["code"] for item in findings],
            "evidence": [item.get("evidence", {}) for item in findings],
            "actions": [
                str(item["suggested_action"])
                for item in findings
                if item.get("suggested_action")
            ],
        })
    report.quality = {
        "deterministic_outcome": {
            "status": "passed" if report.passed else "failed",
            "passed": report.passed,
            "blocking_count": len(blocking),
            "advisory_count": len(advisory),
            "visual_evidence": package_info.get("svg_quality", []),
        },
        "freshness": freshness_payload,
        "reviewability": {
            "review_outcome": review_outcome,
            "review_freshness": review_freshness,
            "required": review_required,
            "requirement_satisfied": review_requirement_satisfied,
            "render_status": render_status,
            "limitations": list(known_limitations),
            "provenance": render_provenance,
        },
        "release_policy": release_policy,
        "font_embedding": font_embedding,
        "blocking_findings": blocking,
        "advisory_findings": advisory,
        "per_slide_evidence": per_slide,
    }
    return report


def report_dict(report: ValidationReport) -> dict[str, Any]:
    return {
        **asdict(report),
        "error_count": len(report.errors),
        "warning_count": len(report.warnings),
    }


def _write_text_atomic(path: Path, text: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    fd, temp_name = tempfile.mkstemp(prefix=f".{path.name}.", suffix=".tmp", dir=path.parent)
    temp_path = Path(temp_name)
    try:
        with os.fdopen(fd, "w", encoding="utf-8") as handle:
            handle.write(text)
            handle.flush()
            os.fsync(handle.fileno())
        os.replace(temp_path, path)
    finally:
        if temp_path.exists():
            temp_path.unlink()


def markdown_report(report: ValidationReport) -> str:
    quality = report.quality
    freshness = quality.get("freshness", {})
    reviewability = quality.get("reviewability", {})
    release_policy = quality.get("release_policy", {})
    font_embedding = quality.get("font_embedding", {})
    lines = [
        "# GHB PPTX Quality Report",
        "",
        "## Deterministic outcome",
        "",
        f"- Result: **{'PASS' if report.passed else 'FAIL'}**",
        f"- File: `{report.pptx}`",
        f"- Pages: {report.page_count} (cover 1 / body {report.body_count} / ending {int(report.has_ending)})",
        f"- Errors: {len(report.errors)}",
        f"- Warnings: {len(report.warnings)}",
        "",
        "## Freshness",
        "",
        f"- Status: `{freshness.get('status', 'fresh')}`",
        f"- States: `{json.dumps(freshness.get('states', {}), ensure_ascii=False, sort_keys=True)}`",
        f"- Issues: `{json.dumps(freshness.get('issues', []), ensure_ascii=False, sort_keys=True)}`",
        "",
        "## Reviewability and limitations",
        "",
        f"- Review outcome: `{reviewability.get('review_outcome', 'unavailable')}`",
        f"- Review freshness: `{reviewability.get('review_freshness', 'unavailable')}`",
        f"- Review required: `{str(reviewability.get('required', False)).lower()}`",
        f"- Review requirement satisfied: `{str(reviewability.get('requirement_satisfied', True)).lower()}`",
        f"- Render status: `{reviewability.get('render_status', 'unavailable')}`",
        f"- Provenance: `{json.dumps(reviewability.get('provenance', {}), ensure_ascii=False, sort_keys=True)}`",
    ]
    lines.extend(f"- Limitation: {item}" for item in reviewability.get("limitations", []))
    lines.extend([
        "",
        "## Release policy",
        "",
        f"- Policy: `{release_policy.get('policy', 'draft')}`",
        f"- Target renderer: `{release_policy.get('target_renderer', 'auto')}`",
        f"- Actual renderer: `{release_policy.get('actual_renderer') or 'unavailable'}`",
        f"- Render evidence bound to final PPTX: `{str(release_policy.get('render_evidence_bound', False)).lower()}`",
        f"- Renderer requirement satisfied: `{str(release_policy.get('target_renderer_satisfied', True)).lower()}`",
        f"- Warning waivers: {release_policy.get('waiver_count', 0)}",
        f"- Unresolved warnings: {release_policy.get('unresolved_warning_count', 0)}",
        "",
        "## Font embedding",
        "",
        f"- Fonts embedded: `{font_embedding.get('fonts_embedded', 0)}`",
        f"- Names: `{json.dumps(font_embedding.get('embedded_font_names', []), ensure_ascii=False)}`",
        f"- Embedding flags enabled: `{str(font_embedding.get('embedding_enabled', False)).lower()}`",
        f"- License fsType accepted: `{font_embedding.get('fsType_ok')}`",
        "",
        "## Blocking findings",
        "",
    ])
    if quality.get("blocking_findings"):
        for issue in quality["blocking_findings"]:
            where = f" slide {issue.get('slide') or issue.get('slide_id')}" if issue.get("slide") or issue.get("slide_id") else ""
            message = issue.get("message") or issue.get("suggested_action") or "See measured evidence."
            lines.append(f"- `ERROR` `{issue['code']}`{where}: {message}")
    else:
        lines.append("No blocking deterministic findings.")
    lines.extend(["", "## Advisory findings", ""])
    if quality.get("advisory_findings"):
        for issue in quality["advisory_findings"]:
            where = f" slide {issue.get('slide') or issue.get('slide_id')}" if issue.get("slide") or issue.get("slide_id") else ""
            message = issue.get("message") or issue.get("suggested_action") or "See measured evidence."
            lines.append(f"- `WARNING` `{issue['code']}`{where}: {message}")
    else:
        lines.append("No advisory deterministic findings.")
    lines.extend(["", "## Per-slide evidence and actions", ""])
    for item in quality.get("per_slide_evidence", []):
        identity = item.get("slide_id") or item.get("slide")
        lines.append(
            f"- Slide {identity} ({item['role']}): blocking={item['blocking_count']}, "
            f"advisory={item['advisory_count']}, codes={','.join(item['issue_codes']) or 'none'}"
        )
    lines.extend([
        "",
        "## Per-slide object summary",
        "",
        "| Slide | Role | Text | Shapes | Groups | Pictures | Tables | Charts | Min font | OOB | Full image | Empty text | Overlaps | Page no. |",
        "|---:|---|---:|---:|---:|---:|---:|---:|---:|---:|---:|---:|---:|---|",
    ])
    svg_quality = report.package.get("svg_quality", [])
    if svg_quality:
        insertion = ["", "## SVG quality gates", ""]
        for item in svg_quality:
            insertion.append(
                f"- `{item['stage']}`: {'PASS' if item['passed'] else 'FAIL'}; "
                f"files={item['file_count']}, errors={item['error_count']}, warnings={item['warning_count']}"
            )
        # Keep the object table immediately after its heading; prepend this section.
        table_heading = lines.index("## Per-slide object summary")
        lines[table_heading:table_heading] = insertion
    for slide in report.slides:
        min_font = "" if slide.min_font_pt is None else f"{slide.min_font_pt:g}"
        lines.append(
            f"| {slide.slide} | {slide.role} | {slide.text_objects} | {slide.shape_objects} | "
            f"{slide.group_objects} | {slide.picture_objects} | {slide.table_objects} | "
            f"{slide.chart_objects} | {min_font} | {slide.out_of_bounds_objects} | "
            f"{slide.full_slide_pictures} | {slide.empty_text_boxes} | "
            f"{slide.possible_text_overlaps} | {'yes' if slide.page_number_found else 'no'} |"
        )
    lines.extend(["", "## Known limitations", ""])
    lines.extend(f"- {item}" for item in report.known_limitations)
    lines.extend(["", "## Manual review", ""])
    lines.extend(f"- {item}" for item in report.manual_review)
    return "\n".join(lines).rstrip() + "\n"


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--body-count", type=int)
    ending = parser.add_mutually_exclusive_group()
    ending.add_argument("--expect-ending", action="store_true")
    ending.add_argument("--no-ending", action="store_true")
    parser.add_argument("--source-svg-dir", type=Path)
    parser.add_argument("--visual-profile", type=Path)
    parser.add_argument("--font-embed-report", type=Path)
    parser.add_argument("--layout-plan", type=Path)
    parser.add_argument("--cover-text", action="append", default=[])
    parser.add_argument("--render-dir", type=Path)
    parser.add_argument("--expect-body-notes", action="store_true")
    parser.add_argument("--svg-report", type=Path, action="append", default=[])
    parser.add_argument("--readback-markdown", type=Path)
    parser.add_argument("--freshness-json", type=Path)
    parser.add_argument("--review-report", type=Path)
    parser.add_argument("--review-required", action="store_true")
    parser.add_argument("--quality-policy", choices=("draft", "release"), default="draft")
    parser.add_argument("--warning-waivers", type=Path)
    parser.add_argument(
        "--target-renderer",
        choices=("auto", "libreoffice", "powerpoint", "wps"),
        default="auto",
    )
    parser.add_argument("--json", action="store_true")
    parser.add_argument("--json-output", type=Path)
    parser.add_argument("--markdown-output", type=Path)
    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    expect_ending = False if args.no_ending else True if args.expect_ending else None
    freshness = None
    if args.freshness_json:
        try:
            freshness = json.loads(args.freshness_json.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError) as exc:
            freshness = {
                "status": "stale",
                "states": {},
                "issues": [{"code": "invalid-freshness-evidence", "message": str(exc)}],
            }
    report = validate_pptx(
        args.pptx,
        expected_body_count=args.body_count,
        expect_ending=expect_ending,
        source_svg_dir=args.source_svg_dir,
        layout_plan_path=args.layout_plan,
        cover_text=args.cover_text,
        render_dir=args.render_dir,
        expect_body_notes=args.expect_body_notes,
        svg_report_paths=args.svg_report,
        readback_markdown_path=args.readback_markdown,
        freshness=freshness,
        review_report_path=args.review_report,
        review_required=args.review_required,
        quality_policy=args.quality_policy,
        warning_waivers_path=args.warning_waivers,
        target_renderer=args.target_renderer,
        visual_profile_path=args.visual_profile,
        font_embed_report_path=args.font_embed_report,
    )
    payload = report_dict(report)
    reviewability = report.quality.get("reviewability", {})
    overall_passed = report.passed and (
        not reviewability.get("required", False)
        or reviewability.get("requirement_satisfied", False)
    )
    if args.json_output:
        _write_text_atomic(
            args.json_output,
            json.dumps(payload, ensure_ascii=False, indent=2) + "\n",
        )
    if args.markdown_output:
        _write_text_atomic(args.markdown_output, markdown_report(report))
    if args.json:
        print(json.dumps(payload, ensure_ascii=False, indent=2))
    else:
        print(f"[{'PASS' if overall_passed else 'FAIL'}] {args.pptx}")
        print(f"  pages={report.page_count} body={report.body_count} ending={int(report.has_ending)}")
        print(f"  errors={len(report.errors)} warnings={len(report.warnings)}")
        for issue in report.issues:
            where = f" slide={issue.slide}" if issue.slide else ""
            print(f"  {issue.severity.upper()} {issue.code}{where}: {issue.message}")
    return 0 if overall_passed else 1


if __name__ == "__main__":
    raise SystemExit(main())
