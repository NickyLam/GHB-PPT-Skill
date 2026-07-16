#!/usr/bin/env python3
"""Offline provisional SVG metrics and the auditable U11 pilot gate."""

from __future__ import annotations

import argparse
import hashlib
import hmac
import json
import math
import shutil
import statistics
import subprocess
import sys
import tempfile
import zipfile
from collections import defaultdict
from pathlib import Path
from typing import Any

from PIL import Image

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))
from scripts.ppt_master.visual_asset_checker import Box, measure_visible_geometry  # noqa: E402


REVIEW_SCHEMA = "ghb.visual-pilot-review.v1"
DETERMINISTIC_SCHEMA = "ghb.visual-pilot-deterministic.v1"
FINAL_REVIEW_SCHEMA = "ghb.visual-final-review.v1"
FINAL_DETERMINISTIC_SCHEMA = "ghb.visual-benchmark-deterministic.v1"


def _canonical_sha256(value: Any) -> str:
    payload = json.dumps(
        value, ensure_ascii=False, sort_keys=True, separators=(",", ":")
    ).encode("utf-8")
    return hashlib.sha256(payload).hexdigest()


def _union_area(boxes: list[Box]) -> float:
    xs = sorted({value for box in boxes for value in (box.x, box.x + box.width)})
    total = 0.0
    for left, right in zip(xs, xs[1:]):
        intervals = sorted(
            (box.y, box.y + box.height)
            for box in boxes
            if box.x < right and box.x + box.width > left and box.height > 0
        )
        if not intervals:
            continue
        start, end = intervals[0]
        height = 0.0
        for next_start, next_end in intervals[1:]:
            if next_start > end:
                height += end - start
                start, end = next_start, next_end
            else:
                end = max(end, next_end)
        total += (right - left) * (height + end - start)
    return total


def _metric(value: float | str | None, coverage: str) -> dict[str, Any]:
    return {
        "value": round(value, 6) if isinstance(value, float) else value,
        "coverage": coverage,
    }


def _geometry_fingerprint(observations: list[dict[str, Any]], body: list[float]) -> str:
    bx, by, width, height = body
    normalized = []
    for item in observations:
        x, y, box_w, box_h = item["box"]
        normalized.append(
            (
                item["role"],
                round((x - bx) / width, 2),
                round((y - by) / height, 2),
                round(box_w / width, 2),
                round(box_h / height, 2),
                bool(item["focal"]),
            )
        )
    normalized.sort()
    canonical = json.dumps(normalized, ensure_ascii=True, separators=(",", ":"))
    return "sha256:" + hashlib.sha256(canonical.encode("utf-8")).hexdigest()


def _focal_zone(observations: list[dict[str, Any]], body: list[float]) -> str | None:
    candidates = [item for item in observations if item["focal"]]
    candidates = [item for item in candidates if item["box"][2] * item["box"][3] > 0]
    if not candidates:
        return None
    focal = max(candidates, key=lambda item: item["box"][2] * item["box"][3])
    center = focal["box"][0] + focal["box"][2] / 2
    relative = (center - body[0]) / body[2]
    return "left" if relative < 1 / 3 else "right" if relative > 2 / 3 else "center"


def measure_svg(
    svg: str,
    *,
    primary_color: str = "#AB1F29",
    base_unit: float = 8.0,
) -> dict[str, Any]:
    """Return explainable metrics derived from visible geometry, never markers."""
    raw = measure_visible_geometry(svg)
    body = raw["body_canvas"]
    body_area = body[2] * body[3]
    observations = raw["observations"]
    boxes = [Box(*item["box"]) for item in observations]
    areas = [box.area for box in boxes if box.area > 0]
    focal_areas = [Box(*item["box"]).area for item in observations if item["focal"] and Box(*item["box"]).area > 0]
    non_focal = [Box(*item["box"]).area for item in observations if not item["focal"] and Box(*item["box"]).area > 0]
    if focal_areas and non_focal:
        focal_ratio = max(focal_areas) / max(statistics.median(non_focal), 1.0)
    elif len(areas) >= 2:
        focal_ratio = max(areas) / max(statistics.median(areas), 1.0)
    else:
        focal_ratio = 1.0 if areas else 0.0

    gaps = raw["gaps"]
    spacing_deviation = (
        statistics.pstdev(gaps) / statistics.mean(gaps)
        if len(gaps) >= 2 and statistics.mean(gaps)
        else 0.0
    )
    edges = [coordinate for box in boxes for coordinate in (box.x, box.y, box.x + box.width, box.y + box.height)]
    alignment_deviation = (
        statistics.mean(min(value % base_unit, base_unit - value % base_unit) for value in edges) / base_unit
        if edges and base_unit > 0
        else 0.0
    )
    primary = primary_color.strip().upper()
    primary_boxes = [
        Box(*item["box"])
        for item in observations
        if item["fill"] == primary and Box(*item["box"]).area > 0
    ]
    focal_primary_fill = any(item["focal"] and item["fill"] == primary for item in observations)
    title_sizes = [item["font_size"] for item in raw["text_sizes"] if item["role"] == "title"]
    body_sizes = [item["font_size"] for item in raw["text_sizes"] if item["role"] != "title"]
    if not title_sizes and len(raw["text_sizes"]) >= 2:
        largest = max(item["font_size"] for item in raw["text_sizes"])
        smaller = [item["font_size"] for item in raw["text_sizes"] if item["font_size"] < largest]
        if smaller:
            title_sizes = [largest]
            body_sizes = smaller
    title_body_ratio = (
        max(title_sizes) / statistics.median(body_sizes) if title_sizes and body_sizes else None
    )
    raw_boxes = [Box(*item["box"]) for item in raw["raw_observations"]]
    if raw_boxes:
        left = min(box.x for box in raw_boxes)
        top = min(box.y for box in raw_boxes)
        right = max(box.x + box.width for box in raw_boxes)
        bottom = max(box.y + box.height for box in raw_boxes)
        content_bounds: dict[str, float] | None = {
            "x": round(left, 6),
            "y": round(top, 6),
            "width": round(right - left, 6),
            "height": round(bottom - top, 6),
        }
    else:
        content_bounds = None
    coverage = raw["coverage"]
    status = coverage["status"]
    return {
        "schema": "ghb.visual-metrics.v1",
        "occupancy": _metric(raw["occupied_area"] / body_area if body_area else 0.0, status),
        "focal-dominance": _metric(focal_ratio, status),
        "focal-emphasis": {
            "area_ratio": round(focal_ratio, 6),
            "primary_fill": focal_primary_fill,
            "coverage": status,
        },
        "title-body-scale": _metric(title_body_ratio, "supported" if title_body_ratio is not None else "not-measurable"),
        "alignment-deviation": _metric(alignment_deviation, status),
        "spacing-consistency": _metric(spacing_deviation, status),
        "minimum-component-gap": _metric(min(gaps) if gaps else None, status if gaps else "not-measurable"),
        "emphasis-color-area": _metric(_union_area(primary_boxes) / body_area if body_area else 0.0, status),
        "composition-fingerprint": {
            "value": _geometry_fingerprint(observations, body),
            "coverage": status,
            "source": "geometry",
        },
        "focal-zone": _metric(_focal_zone(observations, body), status if observations else "not-measurable"),
        "content-bounds": {"value": content_bounds, "coverage": status if content_bounds else "not-measurable"},
        "coverage": coverage,
        "limitations": coverage["limitations"],
    }


def _finding(
    code: str,
    slide_id: str,
    evidence: dict[str, Any],
    expected: dict[str, Any],
    action: str,
    *,
    severity: str = "warning",
) -> dict[str, Any]:
    return {
        "code": code,
        "severity": severity,
        "slide_id": slide_id,
        "evidence": evidence,
        "expected": expected,
        "suggested_action": action,
    }


def _exceptions(page_schema: dict[str, Any]) -> set[str]:
    values = page_schema.get("policy_exceptions", [])
    return {item for item in values if isinstance(item, str)} if isinstance(values, list) else set()


def evaluate_page_quality(
    svg: str,
    *,
    slide_id: str,
    profile: dict[str, Any],
    page_schema: dict[str, Any],
) -> dict[str, Any]:
    """Apply advisory GHB policy without mutating the raw measurements."""
    brand = profile.get("brand") if isinstance(profile.get("brand"), dict) else {}
    spacing = profile.get("spacing") if isinstance(profile.get("spacing"), dict) else {}
    metrics = measure_svg(
        svg,
        primary_color=str(brand.get("primary", "#AB1F29")),
        base_unit=float(spacing.get("base_unit", 8)),
    )
    findings: list[dict[str, Any]] = []
    requested_exceptions = _exceptions(page_schema)
    occupancy_policy = profile.get("occupancy", {}).get("body", {}) if isinstance(profile.get("occupancy"), dict) else {}
    minimum = float(occupancy_policy.get("min", 0.42))
    maximum = float(occupancy_policy.get("max", 0.78))
    occupancy = metrics["occupancy"]
    if occupancy["coverage"] != "not-measurable" and occupancy["value"] < minimum:
        findings.append(_finding(
            "visual-occupancy-below-min", slide_id,
            {"occupancy": occupancy["value"], "coverage": occupancy["coverage"]},
            {"min": minimum, "max": maximum},
            "Increase the scale or spread of meaningful body components.",
        ))
    if occupancy["coverage"] != "not-measurable" and occupancy["value"] > maximum:
        findings.append(_finding(
            "visual-occupancy-above-max", slide_id,
            {"occupancy": occupancy["value"], "coverage": occupancy["coverage"]},
            {"min": minimum, "max": maximum},
            "Reduce or split body components to restore whitespace.",
        ))
    typography = profile.get("typography") if isinstance(profile.get("typography"), dict) else {}
    ratio = metrics["title-body-scale"]
    ratio_min = float(typography.get("min_title_body_ratio", 1.5))
    if ratio["value"] is not None and ratio["value"] < ratio_min:
        findings.append(_finding(
            "visual-title-body-scale-low", slide_id,
            {"title_body_ratio": ratio["value"]}, {"min": ratio_min},
            "Increase title scale relative to body text.",
        ))
    min_gap = float(spacing.get("min_component_gap", 16))
    gap_metric = metrics["minimum-component-gap"]
    if gap_metric["value"] is not None and gap_metric["value"] < min_gap:
        findings.append(_finding(
            "visual-component-gap-small", slide_id,
            {"minimum_gap": gap_metric["value"]}, {"min": min_gap},
            "Separate neighboring components while preserving equal sizing.",
        ))
    if metrics["spacing-consistency"]["value"] > 0.45:
        findings.append(_finding(
            "visual-spacing-inconsistent", slide_id,
            {"coefficient_of_variation": metrics["spacing-consistency"]["value"]},
            {"advisory_max": 0.45},
            "Use a consistent spacing scale between peer components.",
        ))
    if metrics["alignment-deviation"]["value"] > 0.35:
        findings.append(_finding(
            "visual-alignment-deviation", slide_id,
            {"grid_deviation": metrics["alignment-deviation"]["value"]},
            {"advisory_max": 0.35, "base_unit": spacing.get("base_unit", 8)},
            "Align component edges to the profile spacing grid.",
        ))
    if metrics["emphasis-color-area"]["value"] > 0.35:
        findings.append(_finding(
            "visual-primary-color-overuse", slide_id,
            {"primary_color_area": metrics["emphasis-color-area"]["value"]},
            {"advisory_max": 0.35},
            "Reserve the primary brand color for focal emphasis.",
        ))
    if (
        page_schema.get("emphasis") == "single-focal"
        and metrics["focal-dominance"]["value"] <= 1.1
        and not metrics["focal-emphasis"]["primary_fill"]
    ):
        findings.append(_finding(
            "visual-focal-dominance-low", slide_id,
            {"focal_ratio": metrics["focal-dominance"]["value"]},
            {"advisory_min": 1.1},
            "Increase focal contrast through scale, weight, or color without changing peer card sizes.",
        ))
    override = page_schema.get("bounds_override")
    observed_bounds = metrics["content-bounds"]
    if isinstance(override, dict) and observed_bounds["coverage"] == "supported" and isinstance(observed_bounds["value"], dict):
        expected_box = Box(*(float(override[key]) for key in ("x", "y", "width", "height")))
        value = observed_bounds["value"]
        observed_box = Box(*(float(value[key]) for key in ("x", "y", "width", "height")))
        tolerance = 0.5
        outside = (
            observed_box.x < expected_box.x - tolerance
            or observed_box.y < expected_box.y - tolerance
            or observed_box.x + observed_box.width > expected_box.x + expected_box.width + tolerance
            or observed_box.y + observed_box.height > expected_box.y + expected_box.height + tolerance
        )
        if outside:
            findings.append(_finding(
                "visual-explicit-bounds-violation", slide_id,
                {"observed": value, "coverage": "supported"},
                {"bounds_override": override, "tolerance": tolerance},
                "Move or resize content to satisfy the explicitly declared measurable bounds.",
                severity="error",
            ))
    if metrics["coverage"]["status"] == "partial":
        findings.append(_finding(
            "visual-coverage-partial", slide_id,
            metrics["coverage"], {"status": "supported"},
            "Add QA boxes or flatten supported geometry before relying on balance findings.",
        ))
    elif metrics["coverage"]["status"] == "not-measurable":
        findings.append(_finding(
            "visual-not-measurable", slide_id,
            metrics["coverage"], {"measured_elements": ">=1"},
            "Provide visible supported geometry or declared QA boxes.",
        ))
    suppressed = {
        item["code"]
        for item in findings
        if item["severity"] == "warning" and item["code"] in requested_exceptions
    }
    findings = [
        item
        for item in findings
        if not (item["severity"] == "warning" and item["code"] in requested_exceptions)
    ]
    return {
        "slide_id": slide_id,
        "page_schema": page_schema,
        "measurements": metrics,
        "coverage": metrics["coverage"],
        "findings": findings,
        "suppressed_issue_codes": sorted(suppressed),
    }


def analyze_deck_quality(pages: list[dict[str, Any]], *, profile: dict[str, Any]) -> dict[str, Any]:
    """Evaluate geometry/rhythm streaks; all results remain advisory."""
    fingerprints = [
        page["measurements"]["composition-fingerprint"]["value"]
        if page["coverage"]["status"] != "not-measurable" else None
        for page in pages
    ]
    focal_zones = [
        page["measurements"]["focal-zone"]["value"]
        if page["coverage"]["status"] != "not-measurable" else None
        for page in pages
    ]
    densities = [page.get("page_schema", {}).get("density") for page in pages]
    roles = [page.get("page_schema", {}).get("rhythm_role") for page in pages]
    variants = [page.get("page_schema", {}).get("layout_variant") for page in pages]
    findings: list[dict[str, Any]] = []

    def repeated_runs(values: list[Any], minimum: int) -> list[tuple[int, int, Any]]:
        runs: list[tuple[int, int, Any]] = []
        start = 0
        for index in range(1, len(values) + 1):
            if index < len(values) and values[index] == values[start] and values[start] is not None:
                continue
            if values[start] is not None and index - start >= minimum:
                runs.append((start, index - 1, values[start]))
            start = index
        return runs

    def add_run(code: str, run: tuple[int, int, Any], action: str) -> None:
        start, end, value = run
        affected = pages[start : end + 1]
        if any(code in _exceptions(page.get("page_schema", {})) for page in affected):
            return
        findings.append({
            "code": code,
            "severity": "warning",
            "slide_ids": [page["slide_id"] for page in affected],
            "evidence": {"value": value, "streak": end - start + 1},
            "expected": {"max_streak": end - start},
            "suggested_action": action,
        })

    for run in repeated_runs(fingerprints, 2):
        add_run("visual-composition-repeated", run, "Vary actual component geometry, not only data-layout markers.")
    for run in repeated_runs(focal_zones, 3):
        add_run("visual-focal-zone-streak", run, "Move the focal component to vary deck rhythm.")
    max_role = int(profile.get("deck_rhythm", {}).get("max_same_role_streak", 3)) if isinstance(profile.get("deck_rhythm"), dict) else 3
    for run in repeated_runs(roles, max_role + 1):
        add_run("visual-rhythm-role-streak", run, "Insert an anchor or transition page.")
    for run in repeated_runs(densities, 4):
        add_run("visual-density-rhythm-drift", run, "Alternate density where the narrative permits.")
    for run in repeated_runs(variants, 3):
        add_run("visual-variant-repetition", run, "Use a different semantic variant or component arrangement.")
    return {
        "schema": "ghb.visual-deck-metrics.v1",
        "measurements": {
            "composition_fingerprints": fingerprints,
            "focal_zones": focal_zones,
            "densities": densities,
            "rhythm_roles": roles,
            "layout_variants": variants,
        },
        "findings": findings,
    }


def _pending(reasons: list[str]) -> dict[str, Any]:
    return {
        "schema": "ghb.visual-pilot-gate.v1",
        "decision": "pending",
        "proceed": False,
        "preference_rate": None,
        "represented_purpose_results": {},
        "structural_veto_count": 0,
        "blocking_false_positive_count": None,
        "advisory_false_positive_rate": None,
        "reasons": reasons,
    }


def _aggregate_blind_preferences(
    preferences: dict[str, Any],
    review: dict[str, Any],
    *,
    review_schema: str,
    candidate_role: str,
    require_eligibility: bool = False,
    error_scope: str = "blind",
) -> dict[str, Any]:
    """Validate one frozen blind protocol and aggregate page-first preferences."""
    if preferences.get("schema") != "ghb.visual-preferences.v1":
        raise ValueError("invalid-visual-preferences-schema")
    if review.get("schema") != review_schema:
        raise ValueError(f"invalid-{error_scope}-review-schema")
    assignments = review.get("pair_assignments")
    judgments = review.get("judgments")
    eligible_case_ids = review.get("eligible_case_ids")
    if (
        not isinstance(assignments, list)
        or not isinstance(judgments, list)
        or not isinstance(eligible_case_ids, list)
    ):
        raise ValueError(f"invalid-{error_scope}-review-record")
    if not assignments or not judgments:
        return {"pending_reasons": ["insufficient-blind-review-evidence"]}

    baseline_role = "baseline"
    protocol = preferences.get("blind_review_protocol", {})
    assignment_by_case: dict[str, dict[str, Any]] = {}
    for assignment in assignments:
        case_id = assignment.get("page_case_id")
        roles = assignment.get("roles")
        if not isinstance(case_id, str) or case_id in assignment_by_case:
            raise ValueError("invalid-or-duplicate-pair-assignment")
        if not isinstance(roles, dict) or set(roles.values()) != {baseline_role, candidate_role}:
            raise ValueError("invalid-pair-role-assignment")
        masked_ids = {assignment.get("masked_left_id"), assignment.get("masked_right_id")}
        if None in masked_ids or len(masked_ids) != 2 or set(roles) != masked_ids:
            raise ValueError("pair-assignment-mask-mismatch")
        if not isinstance(assignment.get("page_purpose"), str) or not assignment["page_purpose"]:
            raise ValueError("invalid-pair-page-purpose")
        assignment_by_case[case_id] = assignment
    if (
        len(eligible_case_ids) < 3
        or len(eligible_case_ids) != len(set(eligible_case_ids))
        or set(eligible_case_ids) != set(assignment_by_case)
    ):
        raise ValueError(f"incomplete-{error_scope}-case-assignments")

    eligible_reviewers: set[str] | None = None
    if require_eligibility:
        roster = review.get("reviewer_eligibility_roster")
        if not isinstance(roster, list) or not roster:
            raise ValueError("missing-reviewer-eligibility")
        eligible_reviewers = set()
        for attestation in roster:
            reviewer = attestation.get("reviewer_id_hash") if isinstance(attestation, dict) else None
            if not isinstance(reviewer, str) or not reviewer:
                raise ValueError("invalid-reviewer-eligibility")
            if reviewer in eligible_reviewers:
                raise ValueError("duplicate-reviewer-eligibility")
            if any(attestation.get(field) is not True for field in ("independent", "did_not_author", "did_not_tune")):
                raise ValueError("invalid-reviewer-eligibility")
            eligible_reviewers.add(reviewer)

    grouped: dict[str, list[dict[str, Any]]] = defaultdict(list)
    seen: set[tuple[str, str]] = set()
    veto_count = 0
    required_audit = set(protocol.get("audit_fields", []))
    for judgment in judgments:
        missing = required_audit - set(judgment)
        if missing:
            raise ValueError(f"missing-review-audit-fields: {sorted(missing)}")
        case_id = judgment.get("page_case_id")
        reviewer = judgment.get("reviewer_id_hash")
        if case_id not in assignment_by_case or not isinstance(reviewer, str) or not reviewer:
            raise ValueError("invalid-review-identity")
        if eligible_reviewers is not None and reviewer not in eligible_reviewers:
            raise ValueError("missing-reviewer-eligibility")
        key = (case_id, reviewer)
        if key in seen:
            raise ValueError("duplicate-reviewer-page-judgment")
        seen.add(key)
        assignment = assignment_by_case[case_id]
        if judgment.get("masked_left_id") != assignment["masked_left_id"] or judgment.get("masked_right_id") != assignment["masked_right_id"]:
            raise ValueError("review-mask-mismatch")
        if judgment.get("presented_order") != [assignment["masked_left_id"], assignment["masked_right_id"]]:
            raise ValueError("review-presented-order-mismatch")
        choice = judgment.get("judgment")
        if choice not in {*assignment["roles"], "tie", "abstain"}:
            raise ValueError("invalid-review-judgment")
        rubric = judgment.get("rubric")
        if not isinstance(rubric, dict) or set(rubric) != set(protocol.get("rubric_dimensions", [])):
            raise ValueError("invalid-review-rubric")
        if any(
            value != "not-scored"
            and (
                isinstance(value, bool)
                or not isinstance(value, (int, float))
                or not math.isfinite(value)
                or not 1 <= value <= 5
            )
            for value in rubric.values()
        ):
            raise ValueError("invalid-review-rubric")
        if not isinstance(judgment.get("structural_veto"), bool):
            raise ValueError("invalid-review-structural-veto")
        if not isinstance(judgment.get("recorded_at"), str) or not judgment["recorded_at"]:
            raise ValueError("invalid-review-recorded-at")
        veto_count += bool(judgment.get("structural_veto"))
        grouped[case_id].append(judgment)

    minimum_reviewers = protocol.get("minimum_independent_eligible_reviewers_per_page", 3)
    minimum_non_tie = protocol.get("minimum_eligible_non_tie_judgments_per_page", 2)
    page_results: dict[str, str] = {}
    purpose_results: dict[str, list[str]] = defaultdict(list)
    for case_id, assignment in assignment_by_case.items():
        records = grouped.get(case_id, [])
        non_tie = [record for record in records if record["judgment"] not in {"tie", "abstain"}]
        if len(records) < minimum_reviewers or len(non_tie) < minimum_non_tie:
            return {"pending_reasons": ["insufficient-blind-review-evidence"]}
        votes = defaultdict(int)
        for record in non_tie:
            votes[assignment["roles"][record["judgment"]]] += 1
        if votes[candidate_role] == votes[baseline_role]:
            return {"pending_reasons": ["page-preference-tie"]}
        winner = candidate_role if votes[candidate_role] > votes[baseline_role] else baseline_role
        page_results[case_id] = winner
        purpose_results[assignment["page_purpose"]].append(winner)
    preference_rate = sum(result == candidate_role for result in page_results.values()) / len(page_results)
    represented = {
        purpose: "non-regressing" if baseline_role not in results else "regressed"
        for purpose, results in purpose_results.items()
    }
    return {
        "pending_reasons": [],
        "preference_rate": preference_rate,
        "represented_purpose_results": represented,
        "page_results": page_results,
        "structural_veto_count": veto_count,
    }


def evaluate_pilot_gate(
    preferences: dict[str, Any], review: dict[str, Any], deterministic: dict[str, Any]
) -> dict[str, Any]:
    """Evaluate frozen U11 evidence; absent human judgments can never pass."""

    if deterministic.get("schema") != DETERMINISTIC_SCHEMA:
        raise ValueError("invalid-pilot-deterministic-schema")
    aggregate = _aggregate_blind_preferences(
        preferences,
        review,
        review_schema=REVIEW_SCHEMA,
        candidate_role="pilot",
        error_scope="pilot",
    )
    if aggregate["pending_reasons"]:
        reasons = list(aggregate["pending_reasons"])
        if not review.get("judgments"):
            reasons.append("deterministic-audit-not-yet-eligible")
        return _pending(reasons)

    gate = preferences.get("deterministic_pilot_false_positive_gate", {})
    expected_pairs = set(gate.get("advisory_rule_case_pairs", []))
    results = deterministic.get("advisory_rule_case_results")
    if not isinstance(results, list):
        raise ValueError("invalid-advisory-audit")
    actual_pairs = [result.get("rule_case_pair") for result in results]
    if set(actual_pairs) != expected_pairs or len(actual_pairs) != len(expected_pairs):
        raise ValueError("incomplete-advisory-audit")
    allowed_dispositions = {"not-triggered", "true-positive", "false-positive"}
    if any(result.get("disposition") not in allowed_dispositions for result in results):
        raise ValueError("invalid-advisory-audit-disposition")
    blocking = deterministic.get("blocking_false_positives")
    if not isinstance(blocking, list):
        raise ValueError("invalid-blocking-false-positive-audit")
    advisory_false = sum(result["disposition"] == "false-positive" for result in results)
    advisory_rate = advisory_false / max(len(expected_pairs), 1)
    preference_rate = aggregate["preference_rate"]
    represented = aggregate["represented_purpose_results"]
    veto_count = aggregate["structural_veto_count"]
    reasons: list[str] = []
    if preference_rate < 0.70:
        reasons.append("pilot-preference-below-70-percent")
    if "regressed" in represented.values():
        reasons.append("represented-purpose-regression")
    if veto_count:
        reasons.append("structural-regression-veto")
    if len(blocking) > gate.get("blocking_false_positive_maximum", 0):
        reasons.append("blocking-false-positive-ceiling-exceeded")
    if advisory_rate > gate.get("advisory_false_positive_rate_maximum", 0.10):
        reasons.append("advisory-false-positive-ceiling-exceeded")
    return {
        "schema": "ghb.visual-pilot-gate.v1",
        "decision": "failed" if reasons else "passed",
        "proceed": not reasons,
        "preference_rate": round(preference_rate, 6),
        "represented_purpose_results": represented,
        "page_results": aggregate["page_results"],
        "structural_veto_count": veto_count,
        "blocking_false_positive_count": len(blocking),
        "advisory_false_positive_rate": round(advisory_rate, 6),
        "reasons": reasons,
    }


def _final_pending(reasons: list[str]) -> dict[str, Any]:
    return {
        "schema": "ghb.visual-final-gate.v1",
        "decision": "pending",
        "proceed": False,
        "preference_rate": None,
        "represented_purpose_results": {},
        "page_results": {},
        "structural_veto_count": 0,
        "deterministic_fixtures_clean": False,
        "reasons": reasons,
    }


def _file_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _structural_evidence_complete(
    evidence: Any,
    *,
    evidence_root: Path | None,
    expected_slide_texts: list[list[str]],
) -> bool:
    required = {"pptx", "render", "target_fonts", "contact_sheet"}
    if not isinstance(evidence, dict) or set(evidence) < required:
        return False
    if any(
        not isinstance(evidence[key], dict)
        or evidence[key].get("status") != "available"
        for key in required
    ):
        return False
    if evidence_root is None:
        raise ValueError("structural-evidence-root-required")
    root = evidence_root.resolve()
    resolved_paths: dict[str, Path] = {}
    for key in required:
        record = evidence[key]
        relative = record.get("artifact_path")
        expected = record.get("artifact_sha256")
        if (
            not isinstance(relative, str)
            or not relative
            or Path(relative).is_absolute()
            or not isinstance(expected, str)
            or len(expected) != 64
        ):
            raise ValueError("invalid-structural-evidence")
        candidate = root / relative
        if candidate.is_symlink():
            raise ValueError("invalid-structural-evidence")
        try:
            resolved = candidate.resolve(strict=True)
        except OSError as exc:
            raise ValueError("invalid-structural-evidence") from exc
        if not resolved.is_relative_to(root) or not resolved.is_file():
            raise ValueError("invalid-structural-evidence")
        resolved_paths[key] = resolved
        if _file_sha256(resolved) != expected:
            raise ValueError("structural-evidence-digest-mismatch")

    pptx = resolved_paths["pptx"]
    if not zipfile.is_zipfile(pptx):
        raise ValueError("invalid-structural-pptx")
    with zipfile.ZipFile(pptx) as archive:
        if not {"[Content_Types].xml", "ppt/presentation.xml"}.issubset(
            archive.namelist()
        ):
            raise ValueError("invalid-structural-pptx")
    try:
        from pptx import Presentation

        presentation = Presentation(pptx)
    except (OSError, ValueError, KeyError, zipfile.BadZipFile) as exc:
        raise ValueError("invalid-structural-pptx") from exc
    if len(presentation.slides) != len(expected_slide_texts):
        raise ValueError("invalid-structural-pptx")
    for slide, required_texts in zip(
        presentation.slides, expected_slide_texts, strict=True
    ):
        visible_chunks: list[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or not shape.has_text_frame:
                continue
            if (
                shape.width <= 0
                or shape.height <= 0
                or shape.left + shape.width <= 0
                or shape.top + shape.height <= 0
                or shape.left >= presentation.slide_width
                or shape.top >= presentation.slide_height
            ):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size is not None and run.font.size.pt < 10:
                        continue
                    visible_chunks.append(run.text)
        normalized = "".join("".join(visible_chunks).split())
        if any("".join(text.split()) not in normalized for text in required_texts):
            raise ValueError("structural-pptx-content-mismatch")
    try:
        render_payload = json.loads(resolved_paths["render"].read_text(encoding="utf-8"))
        font_payload = json.loads(
            resolved_paths["target_fonts"].read_text(encoding="utf-8")
        )
    except (OSError, json.JSONDecodeError) as exc:
        raise ValueError("invalid-structural-render-evidence") from exc
    if (
        not isinstance(render_payload, dict)
        or render_payload.get("schema") != "ghb.render-report.v1"
        or render_payload.get("status") != "passed"
        or render_payload.get("passed") is not True
        or render_payload.get("pptx_sha256") != _file_sha256(pptx)
    ):
        raise ValueError("invalid-structural-render-evidence")
    def reported_path(value: Any) -> Path:
        candidate = Path(str(value))
        if not candidate.is_absolute():
            candidate = root / candidate
        if candidate.is_symlink():
            raise ValueError("invalid-structural-render-evidence")
        resolved = candidate.resolve(strict=True)
        if not resolved.is_relative_to(root) or not resolved.is_file():
            raise ValueError("invalid-structural-render-evidence")
        return resolved

    try:
        reported_pptx = reported_path(render_payload.get("pptx"))
        reported_contact = reported_path(render_payload.get("contact_sheet"))
        reported_pdf = reported_path(render_payload.get("pdf"))
        reported_pages = [reported_path(value) for value in render_payload.get("pages", [])]
    except (OSError, RuntimeError, TypeError, ValueError) as exc:
        raise ValueError("invalid-structural-render-evidence") from exc
    if reported_pptx != pptx or reported_contact != resolved_paths["contact_sheet"]:
        raise ValueError("structural-render-binding-mismatch")
    if (
        render_payload.get("page_count") != len(expected_slide_texts)
        or len(reported_pages) != len(expected_slide_texts)
        or len(set(reported_pages)) != len(reported_pages)
    ):
        raise ValueError("invalid-structural-render-evidence")
    pdfinfo = shutil.which("pdfinfo")
    if pdfinfo is None:
        raise ValueError("pdf-page-validation-unavailable")
    completed = subprocess.run(
        [pdfinfo, str(reported_pdf)],
        capture_output=True,
        text=True,
        check=False,
        timeout=30,
    )
    page_line = next(
        (line for line in completed.stdout.splitlines() if line.startswith("Pages:")),
        "",
    )
    if completed.returncode or page_line.split()[-1:] != [str(len(expected_slide_texts))]:
        raise ValueError("invalid-structural-render-pdf")
    outputs = {str(value) for value in render_payload.get("outputs", [])}
    expected_outputs = {
        str(render_payload.get("pdf")),
        str(render_payload.get("contact_sheet")),
        *(str(value) for value in render_payload.get("pages", [])),
    }
    if not expected_outputs.issubset(outputs):
        raise ValueError("invalid-structural-render-evidence")
    page_size: tuple[int, int] | None = None
    page_digests: set[str] = set()
    for page in reported_pages:
        try:
            with Image.open(page) as image:
                if image.format != "PNG" or image.width < 640 or image.height < 360:
                    raise ValueError("invalid-structural-render-page")
                extrema = image.convert("RGB").getextrema()
                if all(low == high for low, high in extrema):
                    raise ValueError("invalid-structural-render-page")
                page_size = page_size or (image.width, image.height)
        except (OSError, ValueError) as exc:
            raise ValueError("invalid-structural-render-page") from exc
        page_digests.add(_file_sha256(page))
    if len(page_digests) != len(reported_pages):
        raise ValueError("duplicate-structural-render-pages")
    if resolved_paths["target_fonts"] != resolved_paths["render"]:
        raise ValueError("target-font-evidence-not-bound-to-render")
    font = font_payload.get("font", font_payload) if isinstance(font_payload, dict) else {}
    if not isinstance(font, dict) or font.get("status") != "available":
        raise ValueError("target-font-evidence-not-available")
    try:
        with Image.open(resolved_paths["contact_sheet"]) as image:
            if image.format != "PNG":
                raise ValueError("invalid-structural-contact-sheet")
            if page_size is None:
                raise ValueError("invalid-structural-contact-sheet")
            image.verify()
    except (OSError, ValueError) as exc:
        raise ValueError("invalid-structural-contact-sheet") from exc
    from scripts.render_ghb_pptx import make_contact_sheet

    with tempfile.TemporaryDirectory() as temporary:
        rebuilt = make_contact_sheet(
            reported_pages, Path(temporary) / "contact-sheet.png", columns=3
        )
        if _file_sha256(rebuilt) != _file_sha256(resolved_paths["contact_sheet"]):
            raise ValueError("structural-contact-sheet-binding-mismatch")
    return True


def _validate_final_deterministic_contract(
    contract_document: dict[str, Any], deterministic: dict[str, Any]
) -> bool:
    contract = contract_document.get("fixtures")
    fixtures = deterministic.get("fixtures")
    if (
        contract_document.get("schema")
        != "ghb.visual-final-deterministic-contract.v1"
        or not isinstance(contract, dict)
        or not contract
        or not isinstance(fixtures, list)
    ):
        raise ValueError("final-deterministic-contract-mismatch")
    if deterministic.get("contract_sha256") != _canonical_sha256(contract_document):
        raise ValueError("final-deterministic-contract-mismatch")
    rows: dict[str, dict[str, Any]] = {}
    for row in fixtures:
        fixture_id = row.get("fixture_id") if isinstance(row, dict) else None
        if not isinstance(fixture_id, str) or fixture_id in rows:
            raise ValueError("final-deterministic-contract-mismatch")
        rows[fixture_id] = row
    if set(rows) != set(contract):
        raise ValueError("final-deterministic-contract-mismatch")
    clean = True
    for fixture_id, expected in contract.items():
        row = rows[fixture_id]
        if (
            row.get("expected_issue_codes") != expected.get("expected_issue_codes")
            or row.get("expected_metrics") != expected.get("expected_metrics")
            or row.get("metric_tolerances") != expected.get("metric_tolerances")
        ):
            raise ValueError("final-deterministic-contract-mismatch")
        observed_codes = row.get("observed_issue_codes")
        metrics = row.get("metrics")
        if not isinstance(observed_codes, list) or not isinstance(metrics, dict):
            raise ValueError("final-deterministic-contract-mismatch")
        expected_codes = set(expected["expected_issue_codes"])
        observed_set = set(observed_codes)
        missing = sorted(expected_codes - observed_set)
        unexpected = sorted(observed_set - expected_codes)
        regressions = []
        for metric, expected_value in expected["expected_metrics"].items():
            value = metrics.get(metric)
            tolerance = expected["metric_tolerances"].get(metric)
            if (
                isinstance(value, bool)
                or not isinstance(value, (int, float))
                or isinstance(tolerance, bool)
                or not isinstance(tolerance, (int, float))
                or not math.isfinite(float(value))
                or not math.isfinite(float(tolerance))
                or abs(float(value) - float(expected_value)) > float(tolerance)
            ):
                regressions.append(metric)
        if (
            row.get("missing_expected_issue_codes") != missing
            or row.get("unexpected_issue_codes") != unexpected
            or row.get("metric_regressions") != sorted(regressions)
        ):
            raise ValueError("final-deterministic-contract-mismatch")
        clean = clean and not missing and not unexpected and not regressions
    return clean


def _validate_final_review_integrity(
    review: dict[str, Any], corpus: dict[str, Any], *, evidence_root: Path | None
) -> None:
    cases = corpus.get("cases") if isinstance(corpus, dict) else None
    if not isinstance(cases, list):
        raise ValueError("invalid-final-corpus")
    final_cases = {
        case.get("case_id"): case.get("page_purpose")
        for case in cases
        if isinstance(case, dict) and case.get("partition") == "final-holdout"
    }
    if (
        len(final_cases) < 12
        or set(review.get("eligible_case_ids", [])) != set(final_cases)
    ):
        raise ValueError("final-holdout-case-mismatch")
    randomization = review.get("randomization")
    if (
        not isinstance(randomization, dict)
        or randomization.get("algorithm") != "hmac-sha256"
        or not isinstance(randomization.get("seed"), str)
    ):
        raise ValueError("invalid-final-randomization")
    try:
        seed = bytes.fromhex(randomization["seed"])
    except ValueError as exc:
        raise ValueError("invalid-final-randomization") from exc
    if len(seed) != 32:
        raise ValueError("invalid-final-randomization")
    if evidence_root is None:
        raise ValueError("final-evidence-root-required")
    root = evidence_root.resolve()
    for assignment in review.get("pair_assignments", []):
        if not isinstance(assignment, dict):
            raise ValueError("final-role-assignment-mismatch")
        case_id = assignment.get("page_case_id")
        left = assignment.get("masked_left_id")
        right = assignment.get("masked_right_id")
        if (
            case_id not in final_cases
            or assignment.get("page_purpose") != final_cases[case_id]
            or left != f"{case_id}::A"
            or right != f"{case_id}::B"
        ):
            raise ValueError("final-role-assignment-mismatch")
        optimized_left = hmac.new(seed, str(case_id).encode(), hashlib.sha256).digest()[0] & 1 == 0
        expected_roles = {
            left: "optimized" if optimized_left else "baseline",
            right: "baseline" if optimized_left else "optimized",
        }
        if assignment.get("roles") != expected_roles:
            raise ValueError("final-role-assignment-mismatch")
        blind = assignment.get("blind_artifacts")
        if not isinstance(blind, dict) or set(blind) != {left, right}:
            raise ValueError("invalid-final-blind-evidence")
        for masked in (left, right):
            record = blind[masked]
            relative = record.get("artifact_path") if isinstance(record, dict) else None
            digest = record.get("sha256") if isinstance(record, dict) else None
            if (
                not isinstance(relative, str)
                or not relative
                or Path(relative).is_absolute()
                or not isinstance(digest, str)
                or len(digest) != 64
            ):
                raise ValueError("invalid-final-blind-evidence")
            candidate = root / relative
            if candidate.is_symlink():
                raise ValueError("invalid-final-blind-evidence")
            try:
                resolved = candidate.resolve(strict=True)
            except OSError as exc:
                raise ValueError("invalid-final-blind-evidence") from exc
            if not resolved.is_relative_to(root) or not resolved.is_file():
                raise ValueError("invalid-final-blind-evidence")
            if _file_sha256(resolved) != digest:
                raise ValueError("final-blind-evidence-digest-mismatch")


def _final_expected_slide_texts(
    corpus: dict[str, Any], scenarios: dict[str, Any]
) -> list[list[str]]:
    expected: list[list[str]] = []
    for case in corpus.get("cases", []):
        if not isinstance(case, dict) or case.get("partition") != "final-holdout":
            continue
        source = case.get("source")
        try:
            slide = scenarios[source["scenario_id"]]["slides"][
                source["body_slide_index"] - 1
            ]
            texts = [str(slide["key_message"]), *map(str, slide["items"])]
        except (KeyError, IndexError, TypeError) as exc:
            raise ValueError("invalid-final-scenario-corpus") from exc
        expected.append(texts)
    if len(expected) < 12:
        raise ValueError("invalid-final-scenario-corpus")
    return expected


def evaluate_final_gate(
    preferences: dict[str, Any],
    review: dict[str, Any],
    deterministic: dict[str, Any],
    *,
    corpus: dict[str, Any] | None = None,
    scenarios: dict[str, Any] | None = None,
    fixture_contract: dict[str, Any] | None = None,
    evidence_root: Path | None = None,
) -> dict[str, Any]:
    """Evaluate real final-holdout judgments without duplicating pilot aggregation."""
    if deterministic.get("schema") != FINAL_DETERMINISTIC_SCHEMA:
        raise ValueError("invalid-final-deterministic-schema")
    if corpus is None:
        raise ValueError("final-corpus-required")
    if scenarios is None:
        raise ValueError("final-scenarios-required")
    _validate_final_review_integrity(review, corpus, evidence_root=evidence_root)
    if fixture_contract is None:
        raise ValueError("final-deterministic-contract-required")
    deterministic_clean = _validate_final_deterministic_contract(
        fixture_contract, deterministic
    )
    aggregate = _aggregate_blind_preferences(
        preferences,
        review,
        review_schema=FINAL_REVIEW_SCHEMA,
        candidate_role="optimized",
        require_eligibility=True,
        error_scope="final",
    )
    evidence_complete = _structural_evidence_complete(
        review.get("structural_evidence"),
        evidence_root=evidence_root,
        expected_slide_texts=_final_expected_slide_texts(corpus, scenarios),
    )
    pending_reasons = list(aggregate["pending_reasons"])
    if not evidence_complete:
        pending_reasons.append("structural-evidence-incomplete")
    if pending_reasons:
        result = _final_pending(sorted(set(pending_reasons)))
        result["deterministic_fixtures_clean"] = deterministic_clean
        return result

    preference_rate = aggregate["preference_rate"]
    represented = aggregate["represented_purpose_results"]
    veto_count = aggregate["structural_veto_count"]
    reasons: list[str] = []
    if not deterministic_clean:
        reasons.append("deterministic-fixture-regression")
    if preference_rate < 0.70:
        reasons.append("final-preference-below-70-percent")
    if "regressed" in represented.values():
        reasons.append("purpose-level-regression")
    if veto_count:
        reasons.append("structural-regression-veto")
    return {
        "schema": "ghb.visual-final-gate.v1",
        "decision": "failed" if reasons else "passed",
        "proceed": not reasons,
        "preference_rate": round(preference_rate, 6),
        "represented_purpose_results": represented,
        "page_results": aggregate["page_results"],
        "structural_veto_count": veto_count,
        "deterministic_fixtures_clean": deterministic_clean,
        "reasons": reasons,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    subparsers = parser.add_subparsers(dest="command", required=True)
    measure = subparsers.add_parser("measure")
    measure.add_argument("svg", type=Path)
    gate = subparsers.add_parser("pilot-gate")
    gate.add_argument("--preferences", type=Path, required=True)
    gate.add_argument("--review", type=Path, required=True)
    gate.add_argument("--deterministic", type=Path, required=True)
    gate.add_argument("--output", type=Path)
    final_gate = subparsers.add_parser("final-gate")
    final_gate.add_argument("--preferences", type=Path, required=True)
    final_gate.add_argument("--review", type=Path, required=True)
    final_gate.add_argument("--deterministic", type=Path, required=True)
    final_gate.add_argument("--corpus", type=Path, required=True)
    final_gate.add_argument("--fixture-contract", type=Path, required=True)
    final_gate.add_argument("--scenarios", type=Path, required=True)
    final_gate.add_argument("--output", type=Path)
    args = parser.parse_args()
    if args.command == "measure":
        result = measure_svg(args.svg.read_text(encoding="utf-8"))
    elif args.command == "pilot-gate":
        result = evaluate_pilot_gate(
            json.loads(args.preferences.read_text(encoding="utf-8")),
            json.loads(args.review.read_text(encoding="utf-8")),
            json.loads(args.deterministic.read_text(encoding="utf-8")),
        )
    else:
        result = evaluate_final_gate(
            json.loads(args.preferences.read_text(encoding="utf-8")),
            json.loads(args.review.read_text(encoding="utf-8")),
            json.loads(args.deterministic.read_text(encoding="utf-8")),
            corpus=json.loads(args.corpus.read_text(encoding="utf-8")),
            scenarios=json.loads(args.scenarios.read_text(encoding="utf-8")),
            fixture_contract=json.loads(
                args.fixture_contract.read_text(encoding="utf-8")
            ),
            evidence_root=args.review.parent,
        )
    encoded = json.dumps(result, ensure_ascii=False, indent=2) + "\n"
    if getattr(args, "output", None):
        args.output.write_text(encoded, encoding="utf-8")
    else:
        print(encoded, end="")
    return 0 if result.get("decision") in {None, "passed"} else 2


if __name__ == "__main__":
    raise SystemExit(main())
